"""High-level SDK operations for ableppt.

These functions back both the Python SDK and the `ableppt` CLI.
"""

from __future__ import annotations

import contextlib
import io
import inspect
import json
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Inches, Pt

import ableppt.chart_builder as chart_builder
from ableppt.chart_builder import ChartParser
from ableppt.chart_builder.date_axis import (
    BIWEEKLY_TICKS,
    DAILY_TICKS,
    MONTHLY_TICKS,
    QUARTERLY_TICKS,
    WEEKLY_TICKS,
    YEARLY_TICKS,
    DateAxisConfig,
)
from ableppt.chart_builder.layout import ChartLayoutConfig, LegendConfig, ValueAxisConfig
from ableppt.chart_builder.styles import COLOR_SCHEMES, StyleConfig
from ableppt.connectors import ConnectorFactory
from ableppt.engine import PptEngine
from ableppt.models.job import DataSource, Job, Transform
from ableppt.parsers.ppt_parser import PPTParser
from ableppt.qa.deck_linter import DeckLinter
from ableppt.renderers.ppt_renderer import render_from_json_file
from ableppt.template.chart_presets import CHART_PRESET_FUNCTIONS
from ableppt.template.replacer import TemplateReplacer
from ableppt.transformers import DataFrameTransformer

LEGEND_POSITION_MAP = {
    "top": XL_LEGEND_POSITION.TOP,
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "left": XL_LEGEND_POSITION.LEFT,
    "right": XL_LEGEND_POSITION.RIGHT,
    "corner": XL_LEGEND_POSITION.CORNER,
}


def _read_json(path: Path) -> Any:
    with open(path, "r", encoding="utf-8") as handle:
        return json.load(handle)


def _new_widescreen_presentation() -> Presentation:
    """Return a 16:9 presentation for standalone chart-family renders."""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def parse_template(template_path: str | Path, flat: bool = False) -> dict:
    path = Path(template_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")
    replacer = TemplateReplacer(path)
    return replacer.extract_placeholders(by_page=not flat)


def _describe_chart(chart, chart_idx: int) -> dict:
    try:
        series_config, df, categories_col, layout_info = ChartParser(chart).parse()
    except Exception as exc:  # pragma: no cover - defensive summary path
        return {"chart_index": chart_idx, "error": str(exc)}

    series_summary = []
    for series in series_config:
        col_key = series.get("key", series.get("name", "unknown"))
        data_points = len(df[col_key]) if col_key in df.columns else 0
        series_summary.append(
            {
                "name": series.get("name", col_key),
                "type": series.get("type", "unknown"),
                "axis": series.get("axis", "primary"),
                "data_points": data_points,
            }
        )

    date_range = None
    if categories_col and categories_col in df.columns:
        try:
            first = df[categories_col].iloc[0]
            last = df[categories_col].iloc[-1]
            if hasattr(first, "strftime"):
                date_range = f"{first.strftime('%Y-%m-%d')} ~ {last.strftime('%Y-%m-%d')}"
            else:
                date_range = f"{first} ~ {last}"
        except Exception:  # pragma: no cover - best effort metadata
            pass

    title = ""
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text
    except Exception:  # pragma: no cover - best effort metadata
        pass

    return {
        "chart_index": chart_idx,
        "title": title,
        "series": series_summary,
        "categories_col": categories_col,
        "date_range": date_range,
        "data_rows": len(df),
    }


def describe_chart(input_pptx: str | Path, slide: int | None = None) -> dict:
    input_path = Path(input_pptx)
    if not input_path.exists():
        raise FileNotFoundError(f"文件不存在: {input_path}")

    prs = Presentation(str(input_path))
    result = {"slides": []}

    for slide_idx, slide_obj in enumerate(prs.slides):
        if slide is not None and (slide_idx + 1) != slide:
            continue

        slide_info = {"slide_index": slide_idx, "charts": []}
        chart_idx = 0
        for shape in slide_obj.shapes:
            if not hasattr(shape, "has_chart") or not shape.has_chart:
                continue
            summary = _describe_chart(shape.chart, chart_idx)
            summary["position"] = {
                "left_inches": round(shape.left / 914400, 2),
                "top_inches": round(shape.top / 914400, 2),
                "width_inches": round(shape.width / 914400, 2),
                "height_inches": round(shape.height / 914400, 2),
            }
            slide_info["charts"].append(summary)
            chart_idx += 1
        result["slides"].append(slide_info)

    return result


def list_presets(
    *,
    color_schemes: bool = False,
    date_axis: bool = False,
    chart_presets: bool = False,
) -> dict:
    result = {}
    if color_schemes:
        result["color_schemes"] = dict(COLOR_SCHEMES)
        return result
    if date_axis:
        result["date_axis_presets"] = {
            "DAILY_TICKS": _date_axis_to_dict(DAILY_TICKS),
            "WEEKLY_TICKS": _date_axis_to_dict(WEEKLY_TICKS),
            "BIWEEKLY_TICKS": _date_axis_to_dict(BIWEEKLY_TICKS),
            "MONTHLY_TICKS": _date_axis_to_dict(MONTHLY_TICKS),
            "QUARTERLY_TICKS": _date_axis_to_dict(QUARTERLY_TICKS),
            "YEARLY_TICKS": _date_axis_to_dict(YEARLY_TICKS),
        }
        return result
    if chart_presets:
        result["chart_presets"] = list(CHART_PRESET_FUNCTIONS.keys())
        result["semantic_chart_families"] = chart_builder.list_semantic_families()
        return result

    result["color_schemes"] = dict(COLOR_SCHEMES)
    result["date_axis_presets"] = {
        "DAILY_TICKS": _date_axis_to_dict(DAILY_TICKS),
        "WEEKLY_TICKS": _date_axis_to_dict(WEEKLY_TICKS),
        "BIWEEKLY_TICKS": _date_axis_to_dict(BIWEEKLY_TICKS),
        "MONTHLY_TICKS": _date_axis_to_dict(MONTHLY_TICKS),
        "QUARTERLY_TICKS": _date_axis_to_dict(QUARTERLY_TICKS),
        "YEARLY_TICKS": _date_axis_to_dict(YEARLY_TICKS),
    }
    result["chart_presets"] = list(CHART_PRESET_FUNCTIONS.keys())
    result["semantic_chart_families"] = chart_builder.list_semantic_families()
    return result


def _date_axis_to_dict(cfg: DateAxisConfig) -> dict:
    return {
        "base_unit": cfg.base_unit,
        "major_unit": cfg.major_unit,
        "major_unit_scale": cfg.major_unit_scale,
        "number_format": cfg.number_format,
    }


def validate_job(job_json: str | Path, dry_run: bool = False) -> dict:
    job_path = Path(job_json)
    if not job_path.exists():
        raise FileNotFoundError(f"文件不存在: {job_path}")

    raw = _read_json(job_path)
    job = Job(**raw)
    errors = []

    if job.mode == "composer":
        from ableppt.composer.layouts import LAYOUT_REGISTRY

        all_sources = set(job.datasources.keys())
        if job.transforms:
            all_sources.update(job.transforms.keys())

        for index, page in enumerate(job.pages):
            if page.layout not in LAYOUT_REGISTRY:
                available = ", ".join(LAYOUT_REGISTRY.keys())
                errors.append(f"pages[{index}]: 未知布局 '{page.layout}'。可用: {available}")

            if "datasource" in page.data and page.data["datasource"] not in all_sources:
                errors.append(
                    f"pages[{index}]: datasource '{page.data['datasource']}' 不存在于 datasources/transforms 中"
                )

            for key in ("left", "right"):
                if key in page.data and isinstance(page.data[key], dict):
                    sub = page.data[key]
                    ds_ref = sub.get("datasource") or sub.get("source")
                    if ds_ref and ds_ref not in all_sources:
                        errors.append(
                            f"pages[{index}].data.{key}: datasource/source '{ds_ref}' 不存在于 datasources/transforms 中"
                        )

    if errors:
        raise ValueError("; ".join(errors))

    summary: dict[str, Any] = {"mode": job.mode, "datasources": list(job.datasources.keys())}
    if job.mode == "template":
        summary["slides"] = len(job.slides)
        summary["template"] = job.template.path
    else:
        summary["pages"] = len(job.pages)
        summary["theme"] = job.theme or "able_finance"
        summary["layouts"] = [page.layout for page in job.pages]

    if job.transforms:
        summary["transforms"] = list(job.transforms.keys())

    if dry_run and job.datasources:
        ds_status = {}
        for name, ds in job.datasources.items():
            with contextlib.redirect_stdout(io.StringIO()):
                try:
                    df = ConnectorFactory.load_data(name, ds)
                    ds_status[name] = {"ok": True, "rows": len(df), "columns": list(df.columns)}
                except Exception as exc:
                    ds_status[name] = {"ok": False, "error": str(exc)}
        summary["datasource_check"] = ds_status

    return {"status": "ok", "summary": summary}


def render_job(job_json: str | Path, data_dir: str | Path | None = None) -> dict:
    job_path = Path(job_json)
    if not job_path.exists():
        raise FileNotFoundError(f"文件不存在: {job_path}")

    raw = _read_json(job_path)
    if data_dir:
        prefix = Path(data_dir)
        for ds in raw.get("datasources", {}).values():
            if ds.get("type") == "csv" and ds.get("path"):
                original = Path(ds["path"])
                if not original.is_absolute():
                    ds["path"] = str(prefix / original.name)

    job = Job(**raw)
    with contextlib.redirect_stdout(io.StringIO()):
        engine = PptEngine()
        output_path = engine.render(job)
    return {"status": "ok", "output": str(output_path)}


def parse_ppt(input_pptx: str | Path, output_json: str | Path | None = None) -> dict:
    input_path = Path(input_pptx)
    if not input_path.exists():
        raise FileNotFoundError(f"文件不存在: {input_path}")

    parser = PPTParser(input_path)
    with contextlib.redirect_stdout(io.StringIO()):
        if output_json:
            output_path = Path(output_json)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            parser.save_to_json(output_path)
            return {"status": "ok", "output": str(output_path)}
        return parser.parse()


def rebuild_ppt(input_json: str | Path, output_pptx: str | Path) -> dict:
    input_path = Path(input_json)
    output_path = Path(output_pptx)
    if not input_path.exists():
        raise FileNotFoundError(f"文件不存在: {input_path}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with contextlib.redirect_stdout(io.StringIO()):
        render_from_json_file(input_path, output_path)
    return {"status": "ok", "output": str(output_path)}


def validate_ppt(pptx_path: str | Path, write_notes: bool = False) -> dict:
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")

    linter = DeckLinter(str(path))
    report = linter.run()
    if write_notes:
        linter.write_to_notes(report)
    return report


def chart_engine_info() -> dict:
    compat_path = Path(inspect.getfile(chart_builder))
    combo_path = Path(inspect.getfile(chart_builder.create_combo_chart))
    waterfall_path = Path(inspect.getfile(chart_builder.create_waterfall_chart))
    prepare_path = Path(inspect.getfile(chart_builder.prepare_waterfall_dataframe))
    scatter_path = Path(inspect.getfile(chart_builder.create_scatter_chart))
    bubble_path = Path(inspect.getfile(chart_builder.create_bubble_chart))
    range_snapshot_path = Path(inspect.getfile(chart_builder.create_range_snapshot_chart))

    return {
        "status": "ok",
        "compatibility_layer": {
            "enabled": True,
            "import_path": "ableppt.chart_builder",
            "module": chart_builder.__name__,
            "module_path": str(compat_path),
            "implementation_root_exists": combo_path.exists() and waterfall_path.exists() and range_snapshot_path.exists(),
        },
        "operations": {
            "create_combo_chart": {
                "module": chart_builder.create_combo_chart.__module__,
                "file": str(combo_path),
            },
            "create_waterfall_chart": {
                "module": chart_builder.create_waterfall_chart.__module__,
                "file": str(waterfall_path),
            },
            "create_scatter_chart": {
                "module": chart_builder.create_scatter_chart.__module__,
                "file": str(scatter_path),
            },
            "create_bubble_chart": {
                "module": chart_builder.create_bubble_chart.__module__,
                "file": str(bubble_path),
            },
            "create_range_snapshot_chart": {
                "module": chart_builder.create_range_snapshot_chart.__module__,
                "file": str(range_snapshot_path),
            },
            "prepare_waterfall_dataframe": {
                "module": chart_builder.prepare_waterfall_dataframe.__module__,
                "file": str(prepare_path),
            },
        },
        "supported_surfaces": {
            "sdk": [
                "ableppt.create_combo_chart",
                "ableppt.create_waterfall_chart",
                "ableppt.create_scatter_chart",
                "ableppt.create_bubble_chart",
                "ableppt.create_range_snapshot_chart",
                "ableppt.create_semantic_chart",
                "ableppt.prepare_waterfall_dataframe",
                "ableppt.prepare_range_snapshot_dataframe",
            ],
            "cli": [
                "ableppt chart-engine-info",
                "ableppt render-family <config.json> <output.pptx>",
                "ableppt render-waterfall <config.json> <output.pptx>",
                "ableppt render-scatter <config.json> <output.pptx>",
                "ableppt render-bubble <config.json> <output.pptx>",
                "ableppt render-range-snapshot <config.json> <output.pptx>",
            ],
        },
        "semantic_families": chart_builder.list_semantic_families(),
        "notes": [
            "`ableppt.chart_builder` is a compatibility facade that resolves to the sibling `ablechart` package.",
            "Use `render-family` for demo01-derived semantic chart families, or the atomic render-* commands for base families.",
        ],
    }


def _add_heading_textbox(slide, title: str, subtitle: str | None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12.1), Inches(0.55))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_run = title_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(24)
    title_run.font.bold = True

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.82), Inches(12.1), Inches(0.3))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_run = subtitle_frame.paragraphs[0].runs[0]
        subtitle_run.font.size = Pt(11)


def render_waterfall(config_path: str | Path, output_pptx: str | Path) -> dict:
    config = Path(config_path)
    output = Path(output_pptx)
    if not config.exists():
        raise FileNotFoundError(f"配置文件不存在: {config}")

    raw = _read_json(config)
    csv_path = Path(raw["csv_path"])
    if not csv_path.is_absolute():
        csv_path = config.parent / csv_path
    if not csv_path.exists():
        raise FileNotFoundError(f"CSV 文件不存在: {csv_path}")

    df = pd.read_csv(csv_path)
    categories_col = raw["categories_col"]
    value_col = raw["value_col"]
    measure_col = raw.get("measure_col")
    total_categories = raw.get("total_categories")
    title = raw.get("title")
    subtitle = raw.get("subtitle")
    chart_title = raw.get("chart_title")
    chart_title = raw.get("chart_title")
    chart_title = raw.get("chart_title")

    prs = _new_widescreen_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if title:
        _add_heading_textbox(slide, title, subtitle)
        default_position = (0.6, 1.3)
        default_size = (12.1, 5.4)
    else:
        default_position = (0.6, 0.8)
        default_size = (12.1, 5.9)

    position_inches = raw.get("position_inches", default_position)
    size_inches = raw.get("size_inches", default_size)
    if len(position_inches) != 2 or len(size_inches) != 2:
        raise ValueError("position_inches 和 size_inches 必须是长度为 2 的数组")

    with contextlib.redirect_stdout(io.StringIO()):
        chart_builder.create_waterfall_chart(
            slide=slide,
            df=df,
            categories_col=categories_col,
            value_col=value_col,
            measure_col=measure_col,
            total_categories=total_categories,
            position=(Inches(position_inches[0]), Inches(position_inches[1])),
            size=(Inches(size_inches[0]), Inches(size_inches[1])),
            positive_color=raw.get("positive_color", "10B981"),
            negative_color=raw.get("negative_color", "EF4444"),
            total_color=raw.get("total_color", "1E2761"),
            show_legend=raw.get("show_legend", False),
        )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return {
        "status": "ok",
        "output": str(output),
        "summary": {
            "source_csv": str(csv_path),
            "categories_col": categories_col,
            "value_col": value_col,
            "measure_col": measure_col,
            "total_categories": total_categories or [],
            "prepared_rows": len(df),
        },
    }


def render_scatter(config_path: str | Path, output_pptx: str | Path) -> dict:
    return _render_xy_family(config_path, output_pptx, family="scatter")


def render_bubble(config_path: str | Path, output_pptx: str | Path) -> dict:
    return _render_xy_family(config_path, output_pptx, family="bubble")


def render_range_snapshot(config_path: str | Path, output_pptx: str | Path) -> dict:
    config = Path(config_path)
    output = Path(output_pptx)
    if not config.exists():
        raise FileNotFoundError(f"配置文件不存在: {config}")

    raw = _read_json(config)
    csv_path = Path(raw["csv_path"])
    if not csv_path.is_absolute():
        csv_path = config.parent / csv_path
    if not csv_path.exists():
        raise FileNotFoundError(f"CSV 文件不存在: {csv_path}")

    df = pd.read_csv(csv_path)
    categories_col = raw["categories_col"]
    min_col = raw["min_col"]
    max_col = raw["max_col"]
    average_col = raw["average_col"]
    current_col = raw["current_col"]
    orientation = raw.get("orientation", "vertical")
    title = raw.get("title")
    subtitle = raw.get("subtitle")

    prs = _new_widescreen_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if title:
        _add_heading_textbox(slide, title, subtitle)
        default_position = (0.6, 1.3)
        default_size = (12.1, 5.4)
    else:
        default_position = (0.6, 0.8)
        default_size = (12.1, 5.9)

    position_inches = raw.get("position_inches", default_position)
    size_inches = raw.get("size_inches", default_size)
    if len(position_inches) != 2 or len(size_inches) != 2:
        raise ValueError("position_inches 和 size_inches 必须是长度为 2 的数组")

    with contextlib.redirect_stdout(io.StringIO()):
        chart_builder.create_range_snapshot_chart(
            slide=slide,
            df=df,
            categories_col=categories_col,
            min_col=min_col,
            max_col=max_col,
            average_col=average_col,
            current_col=current_col,
            orientation=orientation,
            position=(Inches(position_inches[0]), Inches(position_inches[1])),
            size=(Inches(size_inches[0]), Inches(size_inches[1])),
            range_color=raw.get("range_color", "5F6772"),
            average_color=raw.get("average_color", "87A330"),
            current_color=raw.get("current_color", "1E88E5"),
            number_format=raw.get("number_format", "0.0x"),
            show_average_ticks=raw.get("show_average_ticks", True),
            show_current_markers=raw.get("show_current_markers", True),
            show_current_labels=raw.get("show_current_labels", True),
            axis_break=raw.get("axis_break"),
        )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return {
        "status": "ok",
        "family": "range_snapshot",
        "output": str(output),
        "summary": {
            "source_csv": str(csv_path),
            "categories_col": categories_col,
            "min_col": min_col,
            "max_col": max_col,
            "average_col": average_col,
            "current_col": current_col,
            "orientation": orientation,
            "rows": len(df),
        },
    }


def render_family(config_path: str | Path, output_pptx: str | Path) -> dict:
    config = Path(config_path)
    output = Path(output_pptx)
    if not config.exists():
        raise FileNotFoundError(f"配置文件不存在: {config}")

    raw = _read_json(config)
    family = raw["family"]
    if family not in chart_builder.SEMANTIC_FAMILY_REGISTRY:
        raise ValueError(
            f"未知 semantic family: {family}。可用: {', '.join(chart_builder.SEMANTIC_FAMILY_REGISTRY.keys())}"
        )

    if not chart_builder.SEMANTIC_FAMILY_REGISTRY[family].get("renderable", False):
        raise ValueError(f"{family} 当前不属于底层可编辑图表 family")

    csv_path = None
    df = None
    needs_top_level_df = family not in {
        chart_builder.TABLE_PLUS_CHART_COMPOSITE_FAMILY,
        chart_builder.FACTOR_ATTRIBUTION_PANEL_FAMILY,
        chart_builder.REGIME_TABLE_PANEL_FAMILY,
        chart_builder.MANAGER_TIMELINE_PROFILE_FAMILY,
        chart_builder.AWARD_TIMELINE_PANEL_FAMILY,
        chart_builder.SELECTION_TIMING_GRID_FAMILY,
        chart_builder.DUAL_CHART_PANEL_FAMILY,
        chart_builder.HOLDING_DETAIL_FAMILY,
    }
    if raw.get("csv_path"):
        csv_path = Path(raw["csv_path"])
        if not csv_path.is_absolute():
            csv_path = config.parent / csv_path
        if not csv_path.exists():
            raise FileNotFoundError(f"CSV 文件不存在: {csv_path}")
        df = pd.read_csv(csv_path)
    elif needs_top_level_df:
        rows = raw.get("rows")
        if not rows:
            raise ValueError("semantic family config 需要 csv_path 或 rows")
        df = pd.DataFrame(rows)

    title = raw.get("title")
    subtitle = raw.get("subtitle")
    chart_title = raw.get("chart_title")
    prs = _new_widescreen_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if title:
        _add_heading_textbox(slide, title, subtitle)
        default_position = (0.6, 1.3)
        default_size = (12.1, 5.4)
    else:
        default_position = (0.6, 0.8)
        default_size = (12.1, 5.9)

    position_inches = raw.get("position_inches", default_position)
    size_inches = raw.get("size_inches", default_size)
    if len(position_inches) != 2 or len(size_inches) != 2:
        raise ValueError("position_inches 和 size_inches 必须是长度为 2 的数组")

    common_position = (Inches(position_inches[0]), Inches(position_inches[1]))
    common_size = (Inches(size_inches[0]), Inches(size_inches[1]))

    kwargs: dict[str, Any] = {
        "slide": slide,
        "position": common_position,
        "size": common_size,
    }

    with contextlib.redirect_stdout(io.StringIO()):
        if family == chart_builder.PERFORMANCE_COMPARE_FAMILY:
            kwargs.update(
                df=df,
                categories_col=raw["categories_col"],
                series_entries=raw["series"],
                title=chart_title,
                number_format=raw.get("number_format", "0.0%"),
                date_number_format=raw.get("date_number_format", "yyyy-mm-dd"),
                color_scheme=raw.get("color_scheme", "demo01_compare"),
            )
            chart_builder.create_performance_compare_chart(**kwargs)
        elif family == chart_builder.DISTRIBUTION_PLUS_HISTORY_FAMILY:
            mode = raw.get("mode", "history")
            if mode == "snapshot":
                kwargs.update(
                    df=df,
                    category_col=raw["category_col"],
                    value_col=raw["value_col"],
                    snapshot_label=raw.get("snapshot_label", "当前分布"),
                    title=chart_title,
                    number_format=raw.get("number_format", "0%"),
                    color_scheme=raw.get("color_scheme", "demo01_distribution"),
                )
                chart_builder.create_distribution_snapshot_chart(**kwargs)
            else:
                kwargs.update(
                    df=df,
                    categories_col=raw["categories_col"],
                    series_columns=raw["series_columns"],
                    chart_type=raw.get("chart_type", "area"),
                    title=chart_title,
                    number_format=raw.get("number_format", "0%"),
                    date_number_format=raw.get("date_number_format", "yyyy-mm-dd"),
                    color_scheme=raw.get("color_scheme", "demo01_distribution"),
                )
                chart_builder.create_distribution_history_chart(**kwargs)
        elif family == chart_builder.STYLE_BOX_FAMILY:
            kwargs.update(
                df=df,
                x_col=raw["x_col"],
                y_col=raw["y_col"],
                series_name=raw.get("series_name"),
                color=raw.get("color", "5679CC"),
            )
            chart_builder.create_style_box_chart(**kwargs)
        elif family == chart_builder.STYLE_ALLOCATION_FAMILY:
            kwargs.update(
                df=df,
                mode=raw.get("mode", "history"),
                title=chart_title,
                categories_col=raw.get("categories_col"),
                series_columns=raw.get("series_columns"),
                category_col=raw.get("category_col"),
                value_col=raw.get("value_col"),
            )
            chart_builder.create_style_allocation_chart(**kwargs)
        elif family == chart_builder.FACTOR_EXPOSURE_FAMILY:
            kwargs.update(
                df=df,
                mode=raw.get("mode", "history"),
                categories_col=raw["categories_col"],
                series_entries=raw["series"],
                title=chart_title,
            )
            chart_builder.create_factor_exposure_chart(**kwargs)
        elif family in {chart_builder.SCORE_OVERLAY_FAMILY, chart_builder.CONCENTRATION_FAMILY}:
            kwargs.update(
                df=df,
                categories_col=raw["categories_col"],
                raw_series=raw["raw_series"],
                score_series=raw["score_series"],
                title=chart_title,
                raw_number_format=raw.get("raw_number_format", "0.0%"),
                score_number_format=raw.get("score_number_format", "0"),
                date_number_format=raw.get("date_number_format", "yyyy-mm-dd"),
                color_scheme=raw.get("color_scheme", "demo01_score"),
            )
            if family == chart_builder.SCORE_OVERLAY_FAMILY:
                chart_builder.create_score_overlay_chart(**kwargs)
            else:
                chart_builder.create_concentration_chart(**kwargs)
        elif family == chart_builder.EVENT_TIMELINE_FAMILY:
            kwargs.update(
                df=df,
                categories_col=raw["categories_col"],
                series_entries=raw["series"],
                events=raw["events"],
                title=chart_title,
                number_format=raw.get("number_format", "0.0%"),
                date_number_format=raw.get("date_number_format", "yyyy-mm-dd"),
                color_scheme=raw.get("color_scheme", "demo01_compare"),
                show_event_labels=raw.get("show_event_labels", False),
            )
            chart_builder.create_event_timeline_chart(**kwargs)
        elif family == chart_builder.ATTRIBUTION_DECOMPOSITION_FAMILY:
            kwargs.update(
                df=df,
                categories_col=raw["categories_col"],
                value_col=raw["value_col"],
                measure_col=raw.get("measure_col"),
                total_categories=raw.get("total_categories"),
                positive_color=raw.get("positive_color", "10B981"),
                negative_color=raw.get("negative_color", "EF4444"),
                total_color=raw.get("total_color", "1E2761"),
            )
            chart_builder.create_attribution_decomposition_chart(**kwargs)
        elif family in {chart_builder.RANKED_TILE_MATRIX_FAMILY, chart_builder.HEATMAP_MATRIX_FAMILY}:
            kwargs.update(
                df=df,
                row_col=raw["row_col"],
                column_col=raw["column_col"],
                value_col=raw["value_col"],
                value_format=raw.get("value_format", ".0f"),
                low_color=raw.get("low_color", "F4F6FB"),
                high_color=raw.get("high_color", "4F66A8"),
                text_dark=raw.get("text_dark", "2F3542"),
                text_light=raw.get("text_light", "FFFFFF"),
                row_header_fill=raw.get("row_header_fill", "FFFFFF"),
                col_header_fill=raw.get("col_header_fill", "FFFFFF"),
                grid_line_color=raw.get("grid_line_color", "DCE3EF"),
            )
            if family == chart_builder.RANKED_TILE_MATRIX_FAMILY:
                chart_builder.create_ranked_tile_matrix_chart(**kwargs)
            else:
                chart_builder.create_heatmap_matrix_chart(**kwargs)
        elif family == chart_builder.TABLE_PLUS_CHART_COMPOSITE_FAMILY:
            kwargs.update(
                chart_family=raw["chart_family"],
                chart_kwargs=raw["chart_spec"],
                headers=raw.get("headers"),
                rows=raw.get("table_rows"),
                chart_ratio=raw.get("chart_ratio", 0.62),
                gap_inches=raw.get("gap_inches", 0.30),
                left_title=raw.get("left_title"),
                right_title=raw.get("right_title"),
            )
            chart_builder.create_table_plus_chart_composite(**kwargs)
        elif family == chart_builder.FACTOR_ATTRIBUTION_PANEL_FAMILY:
            kwargs.update(
                chart_family=raw["chart_family"],
                chart_kwargs=raw["chart_spec"],
                sidebar_title=raw.get("sidebar_title"),
                summary_items=raw.get("summary_items"),
                bullets=raw.get("bullets"),
                chart_ratio=raw.get("chart_ratio", 0.66),
                gap_inches=raw.get("gap_inches", 0.32),
            )
            chart_builder.create_factor_attribution_panel(**kwargs)
        elif family == chart_builder.REGIME_TABLE_PANEL_FAMILY:
            kwargs.update(
                chart_kwargs=raw["chart_spec"],
                table_headers=raw.get("table_headers"),
                table_rows=raw.get("table_rows"),
                top_title=raw.get("top_title"),
                bottom_title=raw.get("bottom_title"),
                chart_height_ratio=raw.get("chart_height_ratio", 0.42),
                gap_inches=raw.get("gap_inches", 0.24),
            )
            chart_builder.create_regime_table_panel(**kwargs)
        elif family == chart_builder.MANAGER_TIMELINE_PROFILE_FAMILY:
            kwargs.update(
                chart_kwargs=raw["chart_spec"],
                manager_name=raw["manager_name"],
                tenure_start=raw["tenure_start"],
                tenure_days=raw["tenure_days"],
                tenure_return=raw["tenure_return"],
                summary_items=raw.get("summary_items"),
            )
            chart_builder.create_manager_timeline_profile(**kwargs)
        elif family == chart_builder.AWARD_TIMELINE_PANEL_FAMILY:
            kwargs.update(
                headers=raw["headers"],
                rows=raw.get("rows"),
                empty_title=raw.get("empty_title", "暂无数据"),
                empty_subtitle=raw.get("empty_subtitle", "暂无数据/结果"),
            )
            chart_builder.create_award_timeline_panel(**kwargs)
        elif family == chart_builder.SELECTION_TIMING_GRID_FAMILY:
            kwargs.update(
                section_titles=raw["section_titles"],
                level_labels=raw["level_labels"],
                row_labels=raw["row_labels"],
                sections=raw["sections"],
            )
            chart_builder.create_selection_timing_grid(**kwargs)
        elif family == chart_builder.HOLDING_DETAIL_FAMILY:
            kwargs.update(
                headers=raw["headers"],
                rows=raw["rows"],
                subtitle=raw.get("subtitle"),
                summary_text=raw.get("summary_text"),
            )
            chart_builder.create_holding_detail_panel(**kwargs)
        elif family == chart_builder.DUAL_CHART_PANEL_FAMILY:
            kwargs.update(
                left_chart_family=raw["left_chart_family"],
                left_chart_kwargs=raw["left_chart_spec"],
                right_chart_family=raw["right_chart_family"],
                right_chart_kwargs=raw["right_chart_spec"],
                left_title=raw.get("left_title"),
                right_title=raw.get("right_title"),
                gap_inches=raw.get("gap_inches", 0.28),
            )
            chart_builder.create_dual_chart_panel(**kwargs)
        else:
            raise ValueError(f"暂未支持的 semantic family: {family}")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return {
        "status": "ok",
        "family": family,
        "output": str(output),
        "summary": {
            "source_csv": str(csv_path) if csv_path else None,
            "rows": len(df) if df is not None else None,
        },
    }


def _render_xy_family(config_path: str | Path, output_pptx: str | Path, *, family: str) -> dict:
    config = Path(config_path)
    output = Path(output_pptx)
    if not config.exists():
        raise FileNotFoundError(f"配置文件不存在: {config}")

    raw = _read_json(config)
    csv_path = Path(raw["csv_path"])
    if not csv_path.is_absolute():
        csv_path = config.parent / csv_path
    if not csv_path.exists():
        raise FileNotFoundError(f"CSV 文件不存在: {csv_path}")

    df = pd.read_csv(csv_path)
    x_col = raw["x_col"]
    y_col = raw["y_col"]
    size_col = raw.get("size_col")
    title = raw.get("title")
    subtitle = raw.get("subtitle")

    prs = _new_widescreen_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if title:
        _add_heading_textbox(slide, title, subtitle)
        default_position = (0.6, 1.3)
        default_size = (12.1, 5.4)
    else:
        default_position = (0.6, 0.8)
        default_size = (12.1, 5.9)

    position_inches = raw.get("position_inches", default_position)
    size_inches = raw.get("size_inches", default_size)
    if len(position_inches) != 2 or len(size_inches) != 2:
        raise ValueError("position_inches 和 size_inches 必须是长度为 2 的数组")

    with contextlib.redirect_stdout(io.StringIO()):
        if family == "scatter":
            chart_builder.create_scatter_chart(
                slide=slide,
                df=df,
                x_col=x_col,
                y_col=y_col,
                series_name=raw.get("series_name"),
                position=(Inches(position_inches[0]), Inches(position_inches[1])),
                size=(Inches(size_inches[0]), Inches(size_inches[1])),
                color=raw.get("color", "1E2761"),
                marker_size=raw.get("marker_size", 9),
            )
        else:
            if not size_col:
                raise ValueError("bubble config 缺少 size_col")
            chart_builder.create_bubble_chart(
                slide=slide,
                df=df,
                x_col=x_col,
                y_col=y_col,
                size_col=size_col,
                series_name=raw.get("series_name"),
                position=(Inches(position_inches[0]), Inches(position_inches[1])),
                size=(Inches(size_inches[0]), Inches(size_inches[1])),
                color=raw.get("color", "1E2761"),
                marker_size=raw.get("marker_size", 9),
            )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    summary = {
        "source_csv": str(csv_path),
        "x_col": x_col,
        "y_col": y_col,
        "rows": len(df),
    }
    if size_col:
        summary["size_col"] = size_col

    return {"status": "ok", "family": family, "output": str(output), "summary": summary}


def fetch_data(
    config_path: str | Path,
    output_dir: str | Path,
    transforms_path: str | Path | None = None,
) -> dict:
    config = Path(config_path)
    if not config.exists():
        raise FileNotFoundError(f"配置文件不存在: {config}")

    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)

    raw_ds = _read_json(config)
    datasources = {name: DataSource(**spec) for name, spec in raw_ds.items()}

    raw_dfs = {}
    for name, ds in datasources.items():
        with contextlib.redirect_stdout(io.StringIO()):
            raw_dfs[name] = ConnectorFactory.load_data(name, ds)

    dfs = raw_dfs
    if transforms_path:
        tf_path = Path(transforms_path)
        if not tf_path.exists():
            raise FileNotFoundError(f"transforms 文件不存在: {tf_path}")
        raw_tf = _read_json(tf_path)
        transforms = {name: Transform(**spec) for name, spec in raw_tf.items()}
        dfs = DataFrameTransformer.apply_transforms(raw_dfs, transforms)

    result = {}
    for name, df in dfs.items():
        csv_path = output / f"{name}.csv"
        df.to_csv(csv_path, index=False, encoding="utf-8")
        result[name] = {"path": str(csv_path), "rows": len(df), "columns": list(df.columns)}

    return {"status": "ok", "datasets": result}


def build_chart_config(chart_def: dict, config_dir: Path) -> dict:
    csv_path = config_dir / chart_def["csv_path"]
    if not csv_path.exists():
        raise FileNotFoundError(f"CSV 文件不存在: {csv_path}")

    df = pd.read_csv(str(csv_path))
    categories_col = chart_def["categories_col"]
    if categories_col in df.columns:
        try:
            df[categories_col] = pd.to_datetime(df[categories_col])
        except Exception:
            pass

    style_def = chart_def.get("style", {})
    style_config = StyleConfig(
        color_scheme=style_def.get("color_scheme", "aim00"),
        line_width_pt=style_def.get("line_width_pt", 2.0),
        marker_style=style_def.get("marker_style", "none"),
    )

    layout_def = chart_def.get("layout", {})
    legend_def = layout_def.get("legend", {})
    legend_pos = LEGEND_POSITION_MAP.get(legend_def.get("position", "top"), XL_LEGEND_POSITION.TOP)
    legend_config = LegendConfig(
        position=legend_pos,
        font_size_pt=legend_def.get("font_size_pt", 9),
        font_name=legend_def.get("font_name", "黑体"),
    )
    value_axis = layout_def.get("value_axis", {})
    secondary_axis = layout_def.get("secondary_value_axis", {})
    layout_config = ChartLayoutConfig(
        title=layout_def.get("title", ""),
        legend_config=legend_config,
        value_axis_config=ValueAxisConfig(
            number_format=value_axis.get("number_format", "0%"),
            font_size_pt=value_axis.get("font_size_pt", 9),
            font_name=value_axis.get("font_name", "黑体"),
            has_major_gridlines=value_axis.get("has_major_gridlines", False),
        ),
        secondary_value_axis_config=ValueAxisConfig(
            number_format=secondary_axis.get("number_format", "#,##0"),
            font_size_pt=secondary_axis.get("font_size_pt", 9),
            font_name=secondary_axis.get("font_name", "黑体"),
            has_major_gridlines=secondary_axis.get("has_major_gridlines", False),
        ),
    )

    date_axis = layout_def.get("date_axis", {})
    if date_axis:
        layout_config.date_axis_config = DateAxisConfig(
            base_unit=date_axis.get("base_unit", "days"),
            major_unit=date_axis.get("major_unit", len(df) // 7),
            number_format=date_axis.get("number_format", "yyyy/mm"),
        )

    return {
        "df": df,
        "categories_col": categories_col,
        "series_config": chart_def["series_config"],
        "style_config": style_config,
        "layout_config": layout_config,
    }


def generate_ppt(template_path: str | Path, config_path: str | Path, output_path: str | Path) -> dict:
    template = Path(template_path)
    config = Path(config_path)
    output = Path(output_path)
    if not template.exists():
        raise FileNotFoundError(f"模板文件不存在: {template}")
    if not config.exists():
        raise FileNotFoundError(f"配置文件不存在: {config}")

    config_dir = config.parent
    raw = _read_json(config)
    chart_configs = {
        name: build_chart_config(chart_def, config_dir)
        for name, chart_def in raw.get("chart_configs", {}).items()
    }

    replacer = TemplateReplacer(template)
    with contextlib.redirect_stdout(io.StringIO()):
        replacer.replace(data=raw.get("text_data", {}), chart_configs=chart_configs)
        output.parent.mkdir(parents=True, exist_ok=True)
        replacer.save(output)
    return {"status": "ok", "output": str(output)}


def parse_waterfall(input_pptx: str | Path, slide_idx: int = 0, shape_idx: int = 0) -> dict:
    with contextlib.redirect_stdout(io.StringIO()):
        result = chart_builder.parse_waterfall_from_pptx(str(input_pptx), slide_idx=slide_idx, shape_idx=shape_idx)
    return {
        "categories_col": result.categories_col,
        "value_col": result.value_col,
        "measure_col": result.measure_col,
        "rows": result.df.to_dict(orient="records"),
    }


def parse_scatter(input_pptx: str | Path, slide_idx: int = 0, shape_idx: int = 0) -> dict:
    with contextlib.redirect_stdout(io.StringIO()):
        result = chart_builder.parse_scatter_from_pptx(str(input_pptx), slide_idx=slide_idx, shape_idx=shape_idx)
    return {
        "chart_family": result.chart_family,
        "x_col": result.x_col,
        "y_col": result.y_col,
        "size_col": result.size_col,
        "rows": result.df.to_dict(orient="records"),
    }


def parse_bubble(input_pptx: str | Path, slide_idx: int = 0, shape_idx: int = 0) -> dict:
    with contextlib.redirect_stdout(io.StringIO()):
        result = chart_builder.parse_bubble_from_pptx(str(input_pptx), slide_idx=slide_idx, shape_idx=shape_idx)
    return {
        "chart_family": result.chart_family,
        "x_col": result.x_col,
        "y_col": result.y_col,
        "size_col": result.size_col,
        "rows": result.df.to_dict(orient="records"),
    }


def parse_range_snapshot(input_pptx: str | Path, slide_idx: int = 0, shape_idx: int = 0) -> dict:
    with contextlib.redirect_stdout(io.StringIO()):
        result = chart_builder.parse_range_snapshot_from_pptx(
            str(input_pptx),
            slide_idx=slide_idx,
            shape_idx=shape_idx,
        )
    return {
        "chart_family": "range_snapshot",
        "categories_col": result.categories_col,
        "min_col": result.min_col,
        "max_col": result.max_col,
        "average_col": result.average_col,
        "current_col": result.current_col,
        "orientation": result.orientation,
        "rows": result.df.to_dict(orient="records"),
    }


def demo_waterfall(output_path: str | Path) -> dict:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    df = pd.DataFrame(
        {
            "项目": ["期初收益", "权益贡献", "债券贡献", "汇率拖累", "费用扣除", "期末收益"],
            "值": [8.5, 2.1, 1.3, -1.8, 0.0, 10.1],
            "measure": ["total", "relative", "relative", "relative", "relative", "total"],
        }
    )
    prs = _new_widescreen_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    with contextlib.redirect_stdout(io.StringIO()):
        chart_builder.create_waterfall_chart(
            slide=slide,
            df=df,
            categories_col="项目",
            value_col="值",
            measure_col="measure",
            position=(Inches(0.8), Inches(1.2)),
            size=(Inches(11.8), Inches(4.8)),
        )
        prs.save(output)
    return {"status": "ok", "output": str(output)}
