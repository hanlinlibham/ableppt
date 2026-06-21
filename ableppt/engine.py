"""
PPT 渲染引擎
"""

from pathlib import Path
from typing import Dict
from datetime import datetime
import hashlib
import json
import pandas as pd
from pptx import Presentation

from ableppt.models.job import Job
from ableppt.config import settings
from ableppt.connectors import ConnectorFactory
from ableppt.transformers import DataFrameTransformer
from ableppt.renderers import TextRenderer, TableRenderer
from ableppt.utils import get_slide_layouts
from ableppt.chart_builder import create_combo_chart


class PptEngine:
    """PPT 渲染引擎"""

    def __init__(self):
        self.text_renderer = TextRenderer()
        self.table_renderer = TableRenderer()

    def render(self, job: Job) -> Path:
        """
        渲染 PPT

        Args:
            job: 任务配置

        Returns:
            生成的 PPT 文件路径
        """
        # 1. 加载数据（共用）
        raw_dfs = {}
        for name, datasource in job.datasources.items():
            print(f"正在加载数据源: {name}")
            raw_dfs[name] = ConnectorFactory.load_data(name, datasource)

        # 2. 数据转换（共用）
        dfs = DataFrameTransformer.apply_transforms(raw_dfs, job.transforms or {})

        # 3. 按模式分发
        if job.mode == "composer":
            prs = self._render_composer(job, dfs)
        else:
            prs = self._render_template(job, dfs)

        # 4. 添加元数据
        if job.output.add_metadata:
            self._add_metadata(prs, job)

        # 5. 保存文件
        output_path = Path(job.output.path)
        if not output_path.is_absolute():
            output_path = settings.output_dir / output_path

        output_path.parent.mkdir(parents=True, exist_ok=True)

        if output_path.exists() and not job.output.overwrite:
            raise FileExistsError(f"输出文件已存在: {output_path}")

        prs.save(str(output_path))
        print(f"PPT 已保存: {output_path}")

        return output_path

    # ------------------------------------------------------------------
    # Template 模式（原有逻辑，零修改搬入）
    # ------------------------------------------------------------------

    def _render_template(self, job: Job, dfs: Dict[str, pd.DataFrame]) -> Presentation:
        """Template 模式渲染：基于模板 + slides 配置"""
        # 加载模板
        template_path = Path(job.template.path)
        if not template_path.is_absolute():
            template_path = settings.templates_dir / template_path

        if not template_path.exists():
            raise FileNotFoundError(f"模板文件不存在: {template_path}")

        prs = Presentation(template_path)

        # 准备渲染上下文
        context = {"params": job.params or {}}

        # 获取幻灯片版式
        slide_layouts = get_slide_layouts(prs)

        # 渲染每一页
        for i, slide_spec in enumerate(job.slides):
            print(f"正在渲染幻灯片: {slide_spec.id}")

            # 如果需要指定版式，则创建新幻灯片
            if slide_spec.layout:
                if slide_spec.layout not in slide_layouts:
                    raise ValueError(f"未找到版式: {slide_spec.layout}")
                slide = prs.slides.add_slide(slide_layouts[slide_spec.layout])
            elif i < len(prs.slides):
                # 使用现有幻灯片
                slide = prs.slides[i]
            else:
                # 使用默认版式
                slide = prs.slides.add_slide(prs.slide_layouts[0])

            # 渲染文本
            if slide_spec.texts:
                self.text_renderer.render(slide, slide_spec.texts, context)

            # 渲染表格
            if slide_spec.tables:
                for table_spec in slide_spec.tables:
                    if table_spec.source not in dfs:
                        print(f"警告: 数据源 '{table_spec.source}' 不存在")
                        continue
                    self.table_renderer.render(slide, table_spec, dfs[table_spec.source])

            # 渲染图表 - 使用新的 chart_builder
            if slide_spec.charts:
                for chart_spec in slide_spec.charts:
                    if chart_spec.source not in dfs:
                        print(f"警告: 数据源 '{chart_spec.source}' 不存在")
                        continue

                    df = dfs[chart_spec.source]

                    try:
                        create_combo_chart(
                            slide=slide,
                            data=df,
                            shape_name=chart_spec.shape_name if hasattr(chart_spec, 'shape_name') else None,
                        )
                    except Exception as e:
                        print(f"警告: 图表渲染失败: {e}")

        return prs

    # ------------------------------------------------------------------
    # Composer 模式（新增）
    # ------------------------------------------------------------------

    def _render_composer(self, job: Job, dfs: Dict[str, pd.DataFrame]) -> Presentation:
        """Composer 模式渲染：基于 PageComposer + pages 配置"""
        from ableppt.composer import PageComposer
        from ableppt.composer.themes import resolve_theme

        theme = resolve_theme(job.theme or "able_finance", aspect_ratio=job.aspect_ratio)
        composer = PageComposer(theme=theme)
        default_lc = job.default_layout_config

        page_specs = self._apply_deck_workflow(job)
        for layout_name, page_data in page_specs:
            data = self._resolve_data_refs(page_data, dfs, default_layout_config=default_lc)
            composer.add_page(layout_name, data)

        return composer.prs

    @staticmethod
    def _apply_deck_workflow(job: Job):
        """GTM deck 工作流编排：
        1. deck 级默认值（brand/market/source/section/tag）注入每个 gtm_* 页面
        2. 页码自增（未显式给 page_num 的页面按顺序编号）
        3. 章节延续（页面缺 section 时沿用上一页）
        4. gtm_toc 页 items 缺省时按后续页面的 section/title 自动生成
        """
        deck = dict(job.deck or {})
        inheritable = {k: deck[k] for k in ("brand", "market", "source", "tag", "section_color")
                       if k in deck}
        page_counter = int(deck.get("start_page", 1))
        last_section = deck.get("section")

        staged = []
        for page_spec in job.pages:
            data = dict(page_spec.data)
            is_gtm = page_spec.layout.startswith("gtm_")
            if is_gtm:
                for key, value in inheritable.items():
                    data.setdefault(key, value)
                if "section" in data:
                    last_section = data["section"]
                elif last_section and page_spec.layout in ("gtm_panels", "gtm_chart_text", "gtm_quilt"):
                    data["section"] = last_section
                data.setdefault("page_num", page_counter)
            staged.append((page_spec.layout, data))
            page_counter += 1

        # 目录自动生成：扫描其后的 gtm 内容页
        for idx, (layout_name, data) in enumerate(staged):
            if layout_name == "gtm_toc" and not data.get("items"):
                items, order = {}, []
                for later_layout, later in staged[idx + 1:]:
                    if not later_layout.startswith("gtm_") or later_layout in ("gtm_cover", "gtm_toc"):
                        continue
                    section = later.get("section", "")
                    if section not in items:
                        items[section] = []
                        order.append(section)
                    items[section].append({"title": later.get("title", ""),
                                           "page": later.get("page_num", "")})
                data["items"] = [{"section": s, "entries": items[s]} for s in order]

        return staged

    def _resolve_data_refs(self, data: dict, dfs: Dict[str, pd.DataFrame],
                           default_layout_config=None) -> dict:
        """解析 data 中的 datasource/source 引用，替换为实际 DataFrame

        语义规则:
        - datasource (新 canonical key): 始终视为数据源引用
        - source: 值匹配已知数据源时视为引用并 pop；否则保留为脚注文本
        - footnote: 始终保留，用于页脚渲染

        执行顺序（修复嵌套 layout_config 类型错误）:
        1. 顶层 datasource 解析
        2. 嵌套 datasource 解析（不转 config 对象）
        3. 注入 defaults（全部仍为 dict，安全合并）
        4. 统一转对象（顶层 + 嵌套）
        """
        resolved = self._resolve_nested_sources(dict(data), dfs)

        # Step 3: 注入 defaults（全部仍为 dict，安全合并）
        if default_layout_config:
            self._inject_default_layout_config(resolved, default_layout_config)

        # Step 4: 统一转对象（顶层 + 嵌套）
        self._resolve_config_objects(resolved)
        for key in ("left", "right", "top", "bottom"):
            if key in resolved and isinstance(resolved[key], dict):
                self._resolve_config_objects(resolved[key])

        return resolved

    def _resolve_nested_sources(self, node, dfs: Dict[str, pd.DataFrame]):
        """递归解析任意嵌套 dict/list 里的 datasource/source 引用。"""

        if isinstance(node, list):
            return [self._resolve_nested_sources(item, dfs) for item in node]

        if not isinstance(node, dict):
            return node

        resolved = dict(node)
        ds_key = None
        if "datasource" in resolved:
            ds_key = "datasource"
        elif "source" in resolved and resolved["source"] in dfs:
            ds_key = "source"

        if ds_key:
            resolved["df"] = dfs[resolved.pop(ds_key)]

        for key, value in list(resolved.items()):
            if isinstance(value, (dict, list)):
                resolved[key] = self._resolve_nested_sources(value, dfs)

        return resolved

    def _inject_default_layout_config(self, data: dict, defaults: dict) -> None:
        """将 Job 级 default_layout_config 合并到页面数据中（页面级优先覆盖）"""
        def _merge_lc(target):
            if "layout_config" not in target and "df" not in target:
                return  # 非图表数据，跳过
            if "layout_config" in target:
                merged = dict(defaults)
                merged.update(target["layout_config"])  # 页面级覆盖
                target["layout_config"] = merged
            elif "df" in target:
                target["layout_config"] = dict(defaults)

        _merge_lc(data)
        for key in ("left", "right", "top", "bottom"):
            if key in data and isinstance(data[key], dict):
                _merge_lc(data[key])

    def _resolve_config_objects(self, data: dict) -> None:
        """将 dict 形式的 style_config / layout_config 转为对象"""
        from ableppt.chart_builder.styles import StyleConfig
        from ableppt.chart_builder.layout import (
            ChartLayoutConfig, LegendConfig, ValueAxisConfig, CategoryAxisConfig,
        )

        if "style_config" in data and isinstance(data["style_config"], dict):
            data["style_config"] = StyleConfig(**data["style_config"])

        if "layout_config" in data and isinstance(data["layout_config"], dict):
            data["layout_config"] = self._build_layout_config(data["layout_config"])

    def _build_layout_config(self, cfg: dict):
        """将嵌套 dict 转为 ChartLayoutConfig，处理子对象实例化"""
        from ableppt.chart_builder.layout import (
            ChartLayoutConfig, LegendConfig, ValueAxisConfig, CategoryAxisConfig,
        )

        kwargs = dict(cfg)

        if "legend_config" in kwargs and isinstance(kwargs["legend_config"], dict):
            kwargs["legend_config"] = LegendConfig(**kwargs["legend_config"])

        if "category_axis_config" in kwargs and isinstance(kwargs["category_axis_config"], dict):
            kwargs["category_axis_config"] = CategoryAxisConfig(**kwargs["category_axis_config"])

        if "value_axis_config" in kwargs and isinstance(kwargs["value_axis_config"], dict):
            kwargs["value_axis_config"] = ValueAxisConfig(**kwargs["value_axis_config"])

        if "secondary_value_axis_config" in kwargs and isinstance(kwargs["secondary_value_axis_config"], dict):
            kwargs["secondary_value_axis_config"] = ValueAxisConfig(**kwargs["secondary_value_axis_config"])

        # date_axis_config 用预设名解析
        if "date_axis_config" in kwargs and isinstance(kwargs["date_axis_config"], str):
            from ableppt.chart_builder import (
                DAILY_TICKS, WEEKLY_TICKS, BIWEEKLY_TICKS,
                MONTHLY_TICKS, QUARTERLY_TICKS, YEARLY_TICKS,
            )
            presets = {
                "daily": DAILY_TICKS, "weekly": WEEKLY_TICKS,
                "biweekly": BIWEEKLY_TICKS, "monthly": MONTHLY_TICKS,
                "quarterly": QUARTERLY_TICKS, "yearly": YEARLY_TICKS,
            }
            preset_name = kwargs["date_axis_config"].lower()
            if preset_name in presets:
                kwargs["date_axis_config"] = presets[preset_name]
            else:
                raise ValueError(f"未知日期轴预设: {preset_name}。可用: {list(presets.keys())}")

        return ChartLayoutConfig(**kwargs)

    # ------------------------------------------------------------------
    # 共用工具
    # ------------------------------------------------------------------

    def _add_metadata(self, prs: Presentation, job: Job) -> None:
        """添加生成信息到 PPT 元数据"""
        core_props = prs.core_properties
        core_props.modified = datetime.now()

        metadata_lines = [
            "生成信息:",
            f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        ]

        if job.mode == "template" and job.template:
            metadata_lines.append(f"模板: {job.template.path}")
        elif job.mode == "composer":
            metadata_lines.append(f"模式: composer, 主题: {job.theme or 'able_finance'}")

        if job.params:
            metadata_lines.append(f"参数: {json.dumps(job.params, ensure_ascii=False)}")

        config_str = json.dumps(job.model_dump(), ensure_ascii=False, sort_keys=True, default=str)
        config_hash = hashlib.sha256(config_str.encode()).hexdigest()[:16]
        metadata_lines.append(f"配置哈希: {config_hash}")

        core_props.comments = "\n".join(metadata_lines)
