"""
PPT 视觉风格提取器

从参考 PPT 文件中提取完整的视觉风格（系列颜色、线宽、字体、背景色等），
产出可直接消费的 StyleConfig + ChartLayoutConfig + composer theme。
"""

from collections import Counter
from dataclasses import dataclass, field
from typing import List, Optional

from pptx import Presentation
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Pt

from ..chart_builder.styles import (
    StyleConfig, COLOR_SCHEMES, LINE_WIDTH_PT, DEFAULT_LINE_WIDTH,
)
from ablechart import register_scheme
from ..chart_builder.layout import (
    ChartLayoutConfig, LegendConfig, ValueAxisConfig,
)
from ..chart_builder.parser import ChartParser
from .xml_color_parser import (
    extract_series_colors,
    extract_series_line_widths,
    extract_series_marker_styles,
)


# 图例位置映射（python-pptx enum → 字符串）
_LEGEND_POS_MAP = {
    XL_LEGEND_POSITION.BOTTOM: "bottom",
    XL_LEGEND_POSITION.TOP: "top",
    XL_LEGEND_POSITION.LEFT: "left",
    XL_LEGEND_POSITION.RIGHT: "right",
    XL_LEGEND_POSITION.CORNER: "corner",
}

_LEGEND_POS_REVERSE = {v: k for k, v in _LEGEND_POS_MAP.items()}


def _color_brightness(hex_rgb: str) -> float:
    """计算颜色亮度 (0=黑, 1=白)"""
    try:
        r = int(hex_rgb[0:2], 16)
        g = int(hex_rgb[2:4], 16)
        b = int(hex_rgb[4:6], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) / 255
    except (ValueError, IndexError):
        return 0.5


@dataclass
class StyleProfile:
    """PPT 视觉 DNA — 从参考文件提取的完整视觉风格"""

    # 图表系列样式
    series_colors: List[str] = field(default_factory=lambda: ["C0C0C0", "C00000"])
    line_width_pt: float = 1.0
    marker_style: str = "none"

    # 布局样式
    legend_position: str = "bottom"
    legend_font_size_pt: float = 9
    axis_font_size_pt: float = 9
    axis_font_name: str = "黑体"
    primary_axis_format: Optional[str] = None
    secondary_axis_format: Optional[str] = None
    has_gridlines: bool = False

    # 文本样式
    header_font: str = "微软雅黑"
    body_font: str = "微软雅黑"
    title_size: int = 28
    body_size: int = 14

    # 颜色板
    primary_color: str = "1E2761"
    accent_color: str = "E8B931"
    bg_dark: str = "1A1A2E"
    bg_light: str = "F5F7FA"
    text_dark: str = "1E2761"
    text_light: str = "FFFFFF"

    def register_color_scheme(self, name: str):
        """将提取的系列颜色注册为 COLOR_SCHEMES 中的一个方案

        Args:
            name: 方案名称（如 "extracted" 或 "aim00"）
        """
        register_scheme(name, list(self.series_colors))

    def to_style_config(self, name: Optional[str] = None) -> StyleConfig:
        """转换为 StyleConfig（可直接传入 create_combo_chart）

        Args:
            name: 颜色方案名称。如为 None，自动注册为 "_extracted"

        Returns:
            StyleConfig 实例
        """
        scheme_name = name or "_extracted"
        self.register_color_scheme(scheme_name)

        # 找到最接近的标准线宽
        closest_pt = min(LINE_WIDTH_PT.keys(),
                         key=lambda pt: abs(pt - self.line_width_pt))

        return StyleConfig(
            color_scheme=scheme_name,
            line_width_pt=closest_pt,
            marker_style=self.marker_style,
        )

    def to_chart_layout_config(self, title: Optional[str] = None) -> ChartLayoutConfig:
        """转换为 ChartLayoutConfig

        Args:
            title: 图表标题（可选）

        Returns:
            ChartLayoutConfig 实例
        """
        legend_pos = _LEGEND_POS_REVERSE.get(self.legend_position,
                                              XL_LEGEND_POSITION.BOTTOM)
        legend_cfg = LegendConfig(
            position=legend_pos,
            font_size_pt=self.legend_font_size_pt,
            font_name=self.axis_font_name,
        )

        primary_axis_cfg = ValueAxisConfig(
            number_format=self.primary_axis_format,
            font_size_pt=self.axis_font_size_pt,
            font_name=self.axis_font_name,
            has_major_gridlines=self.has_gridlines,
        )

        secondary_axis_cfg = None
        if self.secondary_axis_format:
            secondary_axis_cfg = ValueAxisConfig(
                number_format=self.secondary_axis_format,
                font_size_pt=self.axis_font_size_pt,
                font_name=self.axis_font_name,
                has_major_gridlines=False,
            )

        return ChartLayoutConfig(
            title=title,
            legend_config=legend_cfg,
            value_axis_config=primary_axis_cfg,
            secondary_value_axis_config=secondary_axis_cfg,
        )

    def to_composer_theme(self, name: str = "extracted") -> dict:
        """转换为 composer theme 字典

        Args:
            name: 主题显示名称

        Returns:
            与 composer/themes.py THEMES 格式一致的字典
        """
        return {
            "name": name,
            "primary": self.primary_color,
            "secondary": self.series_colors[0] if self.series_colors else "CADCFC",
            "accent": self.accent_color,
            "bg_dark": self.bg_dark,
            "bg_light": self.bg_light,
            "card_bg": "FFFFFF",
            "text_dark": self.text_dark,
            "text_light": self.text_light,
            "text_muted": "8896AB",
            "border": "E2E8F0",
            "positive": "10B981",
            "negative": "EF4444",
            "header_font": self.header_font,
            "body_font": self.body_font,
            "margin": 0.6,
            "title_size": self.title_size,
            "subtitle_size": max(14, self.title_size - 10),
            "body_size": self.body_size,
            "caption_size": 10,
            "kpi_size": 44,
        }


class StyleExtractor:
    """从 PPT 文件提取完整视觉风格"""

    def __init__(self, pptx_path: str):
        """
        Args:
            pptx_path: PPT 文件路径
        """
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)

    def extract(self, slide_index: Optional[int] = None) -> StyleProfile:
        """提取完整视觉风格

        Args:
            slide_index: 指定幻灯片索引提取图表样式。
                         None 时自动扫描所有幻灯片找第一个图表。

        Returns:
            StyleProfile 实例
        """
        profile = StyleProfile()

        self._extract_chart_styles(profile, slide_index)
        self._extract_text_styles(profile)
        self._extract_background_colors(profile)

        return profile

    def _extract_chart_styles(self, profile: StyleProfile,
                               slide_index: Optional[int]):
        """从图表提取系列颜色/线宽/标记点 + 布局信息"""
        chart = self._find_first_chart(slide_index)
        if chart is None:
            return

        chart_element = chart._element

        # 系列颜色
        colors = extract_series_colors(chart_element)
        if colors:
            profile.series_colors = colors

        # 线宽（取第一个有效值）
        widths = extract_series_line_widths(chart_element)
        valid_widths = [w for w in widths if w is not None]
        if valid_widths:
            profile.line_width_pt = valid_widths[0]

        # 标记点（取第一个有效值）
        markers = extract_series_marker_styles(chart_element)
        valid_markers = [m for m in markers if m is not None]
        if valid_markers:
            profile.marker_style = valid_markers[0]

        # 布局信息 — 复用 ChartParser._parse_layout_info()
        try:
            parser = ChartParser(chart)
            layout_info = parser._parse_layout_info()

            if "legend" in layout_info:
                leg = layout_info["legend"]
                pos = leg.get("position")
                if pos is not None:
                    profile.legend_position = _LEGEND_POS_MAP.get(pos, "bottom")
                if leg.get("font_size_pt"):
                    profile.legend_font_size_pt = leg["font_size_pt"]

            if "value_axis" in layout_info:
                va = layout_info["value_axis"]
                fmt = va.get("number_format")
                if fmt and fmt != "General":
                    profile.primary_axis_format = fmt
                if va.get("font_size_pt"):
                    profile.axis_font_size_pt = va["font_size_pt"]
                if va.get("has_major_gridlines") is not None:
                    profile.has_gridlines = va["has_major_gridlines"]
        except Exception:
            pass

    def _extract_text_styles(self, profile: StyleProfile):
        """扫描全 PPT 的 text run，统计最常见的字体和字号"""
        font_counter = Counter()
        size_counter = Counter()

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            font_counter[run.font.name] += 1
                        if run.font.size:
                            size_counter[int(run.font.size.pt)] += 1

        # 最常见字体
        if font_counter:
            most_common_font = font_counter.most_common(1)[0][0]
            profile.header_font = most_common_font
            profile.body_font = most_common_font

        # 字号：最大 → title，次大 → body
        if size_counter:
            sorted_sizes = sorted(size_counter.keys(), reverse=True)
            profile.title_size = sorted_sizes[0]
            if len(sorted_sizes) >= 2:
                profile.body_size = sorted_sizes[1]
            else:
                profile.body_size = max(12, sorted_sizes[0] - 14)

    def _extract_background_colors(self, profile: StyleProfile):
        """从幻灯片背景提取颜色，按亮度分类"""
        bg_colors = []

        for slide in self.prs.slides:
            try:
                bg = slide.background
                if bg.fill and bg.fill.type is not None:
                    color = bg.fill.fore_color
                    if color and color.rgb:
                        hex_rgb = str(color.rgb)
                        bg_colors.append(hex_rgb)
            except Exception:
                continue

        if not bg_colors:
            return

        # 按亮度分类
        for hex_rgb in bg_colors:
            brightness = _color_brightness(hex_rgb)
            if brightness < 0.4:
                profile.bg_dark = hex_rgb
                profile.primary_color = hex_rgb
                profile.text_light = "FFFFFF"
            elif brightness > 0.7:
                profile.bg_light = hex_rgb
            else:
                profile.accent_color = hex_rgb

    def _find_first_chart(self, slide_index: Optional[int]):
        """找到第一个图表对象"""
        if slide_index is not None:
            slides = [self.prs.slides[slide_index]]
        else:
            slides = self.prs.slides

        for slide in slides:
            for shape in slide.shapes:
                if hasattr(shape, 'has_chart') and shape.has_chart:
                    return shape.chart

        return None
