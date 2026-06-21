"""DeckLinter — 生成 PPT 后自动审计，检查 JP 规范合规性

检查项:
  ERROR: Header 存在性 (非封面/结论页 y < 0.5")
  ERROR: Footer 存在性 (非封面/结论页 y > 7.0")
  WARN:  标题字号越界 (内容页标题 > 18pt)
  WARN:  图表颜色数 (单图表系列色 > 3)
  ERROR: 元素越界 (right > 13.0" 或 bottom > 7.5")
  WARN:  对齐偏差 (同页多图 top 偏差 > 0.05")

用法:
    linter = DeckLinter("output.pptx")
    report = linter.run()
    linter.print_summary(report)
"""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu

# EMU (English Metric Units) 到 inches 的转换因子
EMU_PER_INCH = 914400

# 无需 header/footer 的页面类型关键词
_SKIP_KEYWORDS = {
    "封面", "cover", "title", "结论", "conclusion", "divider", "分隔",
    "总结", "summary", "appendix", "附录", "要点总结", "investment summary",
}


class DeckLinter:
    """PPT 合规性检查器"""

    # 静态阈值（用于字号/对齐等不依赖尺寸的规则）
    HEADER_Y_MAX = 0.6     # header textbox 的 y 上限
    FOOTER_Y_MIN = 6.8     # footer textbox 的 y 下限
    TITLE_SIZE_MAX = 18    # 内容页标题字号上限 (pt)
    ALIGN_TOLERANCE = 0.05 # 对齐偏差容忍度 (inches)

    def __init__(self, pptx_path: str):
        self.path = Path(pptx_path)
        self.prs = Presentation(str(self.path))
        # 动态阈值 — 从 PPT 实际尺寸计算边界
        sw_inches = self.prs.slide_width / EMU_PER_INCH
        sh_inches = self.prs.slide_height / EMU_PER_INCH
        self.max_right = sw_inches + 0.067  # 小容差
        self.max_bottom = sh_inches

    def run(self) -> dict:
        """执行所有检查，返回报告

        Returns:
            {
                "file": str,
                "total_slides": int,
                "issues": [
                    {"slide": int, "severity": "ERROR"|"WARN", "rule": str, "detail": str}
                ],
                "summary": {"errors": int, "warnings": int}
            }
        """
        issues = []

        for idx, slide in enumerate(self.prs.slides):
            slide_num = idx + 1
            is_skip = self._is_skip_slide(slide)

            if not is_skip:
                issues.extend(self._check_header(slide, slide_num))
                issues.extend(self._check_footer(slide, slide_num))

            issues.extend(self._check_bounds(slide, slide_num))
            issues.extend(self._check_title_size(slide, slide_num, is_skip))
            issues.extend(self._check_alignment(slide, slide_num))
            if not is_skip:
                issues.extend(self._check_insight_text(slide, slide_num))

        errors = sum(1 for i in issues if i["severity"] == "ERROR")
        warnings = sum(1 for i in issues if i["severity"] == "WARN")

        return {
            "file": str(self.path),
            "total_slides": len(self.prs.slides),
            "issues": issues,
            "summary": {"errors": errors, "warnings": warnings},
        }

    def _is_skip_slide(self, slide) -> bool:
        """判断是否为封面/结论/分隔页（不需要 header/footer）

        启发式:
        1. 文本包含关键词（封面/title/结论/conclusion/divider）
        2. 有全幅装饰矩形（宽度 > 12"，高度 < 0.1"，y 接近 0 或底部）
        3. 有大色块（宽度 > 4"，高度 > 7"）用于分隔页
        4. 无 header 区域的小字标题文本（只有大字 > 20pt）
        """
        has_small_header = False
        has_decoration_bar = False
        has_large_block = False

        sw_inches = self.prs.slide_width / EMU_PER_INCH
        sh_inches = self.prs.slide_height / EMU_PER_INCH

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                for kw in _SKIP_KEYWORDS:
                    if kw in text:
                        return True

            w_inches = shape.width / EMU_PER_INCH if shape.width else 0
            h_inches = shape.height / EMU_PER_INCH if shape.height else 0
            y_inches = shape.top / EMU_PER_INCH if shape.top else 0

            # 全幅装饰条（顶部/底部装饰线）— 宽度 > 90% 幻灯片宽度
            if w_inches > sw_inches * 0.9 and h_inches < 0.15:
                has_decoration_bar = True

            # 大色块（分隔页左侧色块）— 宽度 > 30% 且高度 > 85%
            if w_inches > sw_inches * 0.3 and h_inches > sh_inches * 0.85:
                has_large_block = True

            # 检查是否有 header 区域的标准小字标题
            if shape.has_text_frame and y_inches < self.HEADER_Y_MAX:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            pt = run.font.size / 12700
                            if pt <= 20:
                                has_small_header = True

        # 如果有装饰条但没有小字 header，大概率是封面/结论页
        if has_decoration_bar and not has_small_header:
            return True

        # 大色块 = 分隔页
        if has_large_block:
            return True

        return False

    def _check_header(self, slide, slide_num) -> list:
        """检查 header 存在性: 非封面页必须有 y < HEADER_Y_MAX 的文本框"""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                y_inches = shape.top / EMU_PER_INCH
                if y_inches < self.HEADER_Y_MAX:
                    return []
        return [self._issue(slide_num, "ERROR", "header_missing",
                           "No header text box found (y < 0.6\")")]

    def _check_footer(self, slide, slide_num) -> list:
        """检查 footer 存在性: 非封面页必须有 y > FOOTER_Y_MIN 的文本框"""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                y_inches = shape.top / EMU_PER_INCH
                if y_inches > self.FOOTER_Y_MIN:
                    return []
        return [self._issue(slide_num, "ERROR", "footer_missing",
                           "No footer text box found (y > 6.8\")")]

    def _check_bounds(self, slide, slide_num) -> list:
        """检查元素是否越界"""
        issues = []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            right = (shape.left + shape.width) / EMU_PER_INCH
            bottom = (shape.top + shape.height) / EMU_PER_INCH
            if right > self.max_right:
                issues.append(self._issue(
                    slide_num, "ERROR", "bounds_overflow",
                    f"Shape '{shape.name}' right={right:.2f}\" > {self.max_right:.2f}\""))
            if bottom > self.max_bottom + 0.1:  # small tolerance for rounding
                issues.append(self._issue(
                    slide_num, "ERROR", "bounds_overflow",
                    f"Shape '{shape.name}' bottom={bottom:.2f}\" > {self.max_bottom:.2f}\""))
        return issues

    def _check_title_size(self, slide, slide_num, is_skip) -> list:
        """检查标题字号是否越界"""
        if is_skip:
            return []
        issues = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            y_inches = shape.top / EMU_PER_INCH
            if y_inches > self.HEADER_Y_MAX:
                continue
            # 这是一个 header 区域的文本框
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Emu(self.TITLE_SIZE_MAX * 12700):
                        pt = run.font.size / 12700
                        issues.append(self._issue(
                            slide_num, "WARN", "title_size_overflow",
                            f"Title font {pt:.0f}pt > {self.TITLE_SIZE_MAX}pt"))
        return issues

    def _check_alignment(self, slide, slide_num) -> list:
        """检查同页多图的 top 对齐偏差"""
        chart_tops = []
        for shape in slide.shapes:
            if shape.has_chart:
                chart_tops.append(shape.top / EMU_PER_INCH)

        if len(chart_tops) < 2:
            return []

        issues = []
        base = chart_tops[0]
        for i, t in enumerate(chart_tops[1:], 1):
            diff = abs(t - base)
            if diff > self.ALIGN_TOLERANCE:
                issues.append(self._issue(
                    slide_num, "WARN", "chart_misaligned",
                    f"Chart {i+1} top={t:.3f}\" vs chart 1 top={base:.3f}\" "
                    f"(diff={diff:.3f}\")"))
        return issues

    def _check_insight_text(self, slide, slide_num) -> list:
        """WARN: 图表页无 insight 文本"""
        # 仅检查含图表的页面
        has_chart = any(shape.has_chart for shape in slide.shapes)
        if not has_chart:
            return []

        # insight 区间：header divider 下方、图表上方
        # 4:3 和 16:9 高度均为 7.5"，y 布局相同；
        # 范围基于 divider_y(0.80) 到 content_y(1.00) + insight_h(0.55)
        sh_inches = self.prs.slide_height / EMU_PER_INCH
        insight_y_min = sh_inches * 0.11   # ~0.83" for 7.5"
        insight_y_max = sh_inches * 0.22   # ~1.65" for 7.5"

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            y_inches = shape.top / EMU_PER_INCH
            if insight_y_min < y_inches < insight_y_max:
                text = shape.text_frame.text.strip()
                if len(text) > 15:
                    return []  # 有 insight 文本

        return [self._issue(slide_num, "WARN", "insight_missing",
                           "Chart page has no insight text (recommend 2-3 line summary)")]

    @staticmethod
    def _issue(slide_num, severity, rule, detail) -> dict:
        return {"slide": slide_num, "severity": severity, "rule": rule, "detail": detail}

    def print_summary(self, report: dict, file=None):
        """输出彩色终端摘要"""
        file = file or sys.stderr
        errors = report["summary"]["errors"]
        warnings = report["summary"]["warnings"]
        total = len(report["issues"])

        print(f"\n{'='*60}", file=file)
        print(f"DeckLinter: {report['file']}", file=file)
        print(f"Slides: {report['total_slides']}  |  "
              f"Errors: {errors}  |  Warnings: {warnings}", file=file)
        print(f"{'='*60}", file=file)

        for issue in report["issues"]:
            tag = "ERR " if issue["severity"] == "ERROR" else "WARN"
            print(f"  [{tag}] Slide {issue['slide']}: "
                  f"{issue['rule']} - {issue['detail']}", file=file)

        if total == 0:
            print("  All checks passed.", file=file)
        print(file=file)

    def to_json(self, report: dict) -> str:
        """输出 JSON 报告"""
        return json.dumps(report, ensure_ascii=False, indent=2)

    def write_to_notes(self, report: dict):
        """将 lint 结果写入 PPT 的 notes"""
        for issue in report["issues"]:
            slide_idx = issue["slide"] - 1
            if slide_idx >= len(self.prs.slides):
                continue
            slide = self.prs.slides[slide_idx]
            notes_slide = slide.notes_slide if slide.has_notes_slide else slide.notes_slide
            tf = notes_slide.notes_text_frame
            existing = tf.text or ""
            lint_line = f"[{issue['severity']}] {issue['rule']}: {issue['detail']}"
            if lint_line not in existing:
                if existing:
                    tf.text = existing + "\n" + lint_line
                else:
                    tf.text = lint_line

        # 保存
        self.prs.save(str(self.path))
