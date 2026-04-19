import re
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import RGBColor

def read_excel_data(excel_file):
    """从Excel文件中读取数据，并格式化数值为两位小数"""
    workbook = load_workbook(excel_file, data_only=True)
    data = {}
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        sheet_data = {}
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.value is not None:
                    value = cell.value
                    if isinstance(value, (int, float)):
                        sheet_data[cell.coordinate] = f"{value:.2f}"
                    else:
                        sheet_data[cell.coordinate] = str(value)
        data[sheet] = sheet_data
    return data

def copy_font(target_font, source_font):
    """复制字体样式，包括更多属性和颜色类型"""
    # 复制字体的基础属性
    attributes = ['name', 'size', 'bold', 'italic', 'underline', 'strike', 'subscript', 'superscript']
    for attr in attributes:
        if hasattr(source_font, attr):
            setattr(target_font, attr, getattr(source_font, attr))

    # 复制颜色属性
    if source_font.color.type == MSO_COLOR_TYPE.RGB:
        target_font.color.rgb = source_font.color.rgb
    elif source_font.color.type == MSO_COLOR_TYPE.SCHEME:
        target_font.color.theme_color = source_font.color.theme_color
        target_font.color.brightness = source_font.color.brightness
    elif source_font.color.type == MSO_COLOR_TYPE.PRESET:
        target_font.color.theme_color = source_font.color.theme_color
    elif source_font.color.type == MSO_COLOR_TYPE.SYSTEM:
        target_font.color.rgb = source_font.color.rgb


def fonts_are_equal(font1, font2):
    """比较两个字体的样式是否相同，扩展更多的属性和颜色类型"""
    # 比较字体的基础属性
    attributes = ['name', 'size', 'bold', 'italic', 'underline', 'strike', 'subscript', 'superscript']
    for attr in attributes:
        if getattr(font1, attr, None) != getattr(font2, attr, None):
            return False

    # 比较颜色类型和颜色值
    if font1.color.type != font2.color.type:
        return False
    if font1.color.type == MSO_COLOR_TYPE.RGB:
        return font1.color.rgb == font2.color.rgb
    elif font1.color.type == MSO_COLOR_TYPE.SCHEME:
        return (font1.color.theme_color == font2.color.theme_color and
                font1.color.brightness == font2.color.brightness)
    elif font1.color.type == MSO_COLOR_TYPE.PRESET:
        return font1.color.theme_color == font2.color.theme_color  # 比较预设颜色
    elif font1.color.type == MSO_COLOR_TYPE.SYSTEM:
        return font1.color.rgb == font2.color.rgb  # 比较系统颜色

    return True

def replace_text_with_style(text_frame, placeholder_pattern, excel_data):
    """替换占位符并保留格式，支持百分比样式"""
    for paragraph in text_frame.paragraphs:
        char_list = []
        for run in paragraph.runs:
            for char in run.text:
                char_list.append({'char': char, 'font': run.font})

        full_text = ''.join([char_info['char'] for char_info in char_list])
        matches = list(placeholder_pattern.finditer(full_text))

        if not matches:
            continue

        new_char_list = []
        idx = 0

        for match in matches:
            start, end = match.span()
            placeholder = match.group(1)

            while idx < start:
                new_char_list.append(char_list[idx])
                idx += 1

            try:
                sheet_name, cell_ref = placeholder.split('!')
                replacement_text = excel_data.get(sheet_name.strip(), {}).get(cell_ref.strip(), '')
            except ValueError:
                print(f"占位符格式错误: {placeholder}")
                replacement_text = ''

            placeholder_font = char_list[idx - 1]['font'] if idx > 0 else paragraph.runs[0].font

            percent_style = None
            if end < len(char_list) and char_list[end]['char'] == '%':
                percent_style = char_list[end]['font']
                end += 1  # 跳过 '%' 字符

            for char in str(replacement_text):
                new_char_list.append({'char': char, 'font': percent_style or placeholder_font})

            if percent_style:
                new_char_list.append({'char': '%', 'font': percent_style})

            idx = end

        while idx < len(char_list):
            new_char_list.append(char_list[idx])
            idx += 1

        for run in paragraph.runs:
            run.text = ''

        if not new_char_list:
            continue

        current_run = paragraph.add_run()
        current_run.text = new_char_list[0]['char']
        copy_font(current_run.font, new_char_list[0]['font'])

        for char_info in new_char_list[1:]:
            if not fonts_are_equal(char_info['font'], current_run.font):
                current_run = paragraph.add_run()
                current_run.text = char_info['char']
                copy_font(current_run.font, char_info['font'])
            else:
                current_run.text += char_info['char']

def process_shape(shape, excel_data, placeholder_pattern):
    """处理形状（包括文本框、表格和组合形状）"""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            process_shape(subshape, excel_data, placeholder_pattern)
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_text_with_style(cell.text_frame, placeholder_pattern, excel_data)
    elif hasattr(shape, 'text_frame'):
        replace_text_with_style(shape.text_frame, placeholder_pattern, excel_data)

def process_slide(slide, excel_data, placeholder_pattern):
    """处理幻灯片"""
    for shape in slide.shapes:
        process_shape(shape, excel_data, placeholder_pattern)
    
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                replace_text_with_style(shape.text_frame, placeholder_pattern, excel_data)

def replace_placeholders_in_ppt(ppt_file, excel_file, output_file):
    """主函数：替换PPT中的占位符"""
    excel_data = read_excel_data(excel_file)
    placeholder_pattern = re.compile(r"\{([^\}]+)\}")

    presentation = Presentation(ppt_file)
    for slide in presentation.slides:
        process_slide(slide, excel_data, placeholder_pattern)

    presentation.save(output_file)
    print(f"已生成替换后的 PPT：{output_file}")


