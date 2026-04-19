"""
PPT 完整解析器 - 将 PPTX 文件解析为 JSON 格式

支持的元素类型：
- 文本框
- 表格
- 图表
- 组合形状
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from typing import Dict, List, Any
import json

from ..chart_builder import ChartParser


class PPTParser:
    """PPT 完整解析器"""
    
    def __init__(self, pptx_path: Path):
        """
        初始化解析器
        
        Args:
            pptx_path: PPTX 文件路径
        """
        self.pptx_path = pptx_path
        self.prs = Presentation(str(pptx_path))
    
    def parse(self) -> Dict:
        """
        解析整个 PPT 文件
        
        Returns:
            包含所有幻灯片和元素的字典
        """
        result = {
            "metadata": {
                "source_file": str(self.pptx_path.name),
                "slide_count": len(self.prs.slides),
                "slide_width": self.prs.slide_width,  # EMU 单位
                "slide_height": self.prs.slide_height,  # EMU 单位
            },
            "slides": []
        }
        
        for slide_idx, slide in enumerate(self.prs.slides):
            slide_data = self._parse_slide(slide, slide_idx)
            result["slides"].append(slide_data)
        
        return result
    
    def _parse_slide(self, slide, slide_idx: int) -> Dict:
        """解析单个幻灯片"""
        slide_data = {
            "slide_index": slide_idx,
            "elements": []
        }
        
        print(f"\n解析幻灯片 {slide_idx + 1}...")
        
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                element = self._parse_shape(shape, shape_idx)
                if element:
                    slide_data["elements"].append(element)
            except Exception as e:
                print(f"  ⚠️ 形状 {shape_idx + 1} 解析失败: {e}")
                continue
        
        print(f"  ✅ 解析了 {len(slide_data['elements'])} 个元素")
        
        return slide_data
    
    def _parse_shape(self, shape, shape_idx: int) -> Dict:
        """解析单个形状"""
        
        # 基础信息
        element = {
            "shape_index": shape_idx,
            "shape_type": type(shape).__name__,
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
        }
        
        # 图表
        if hasattr(shape, 'has_chart') and shape.has_chart:
            element["type"] = "chart"
            element["data"] = self._parse_chart(shape.chart)
        
        # 表格
        elif hasattr(shape, 'has_table') and shape.has_table:
            element["type"] = "table"
            element["data"] = self._parse_table(shape.table)
        
        # 文本框
        elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:  # 只保存非空文本
                element["type"] = "textbox"
                element["data"] = self._parse_textbox(shape.text_frame)
            else:
                return None  # 跳过空文本框
        
        # 图片（跳过）
        elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element["type"] = "picture"
            element["data"] = {"note": "图片已跳过"}
            return None  # 不保存图片
        
        # 组合形状
        elif hasattr(shape, 'shapes'):  # GroupShape
            element["type"] = "group"
            element["data"] = self._parse_group(shape)
        
        # 其他形状
        else:
            element["type"] = "other"
            element["data"] = {}
            return None  # 跳过其他类型
        
        return element
    
    def _parse_chart(self, chart) -> Dict:
        """解析图表"""
        try:
            parser = ChartParser(chart)
            series_config, df, categories_col, layout_info = parser.parse()
            
            # 将 DataFrame 转换为 JSON 可序列化的格式
            df_dict = df.to_dict(orient='records')
            
            # 转换日期为字符串
            for record in df_dict:
                for key, value in record.items():
                    if hasattr(value, 'strftime'):  # datetime 对象
                        record[key] = value.strftime('%Y-%m-%d')
            
            return {
                "title": chart.chart_title.text_frame.text if chart.has_title else "",
                "chart_type": str(chart.chart_type),
                "series_config": series_config,
                "data": df_dict,
                "categories_col": categories_col,
                "layout_info": layout_info,
            }
        except Exception as e:
            return {"error": str(e)}
    
    def _parse_table(self, table) -> Dict:
        """解析表格"""
        rows_data = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                row_data.append(cell_text)
            rows_data.append(row_data)
        
        return {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "data": rows_data,
        }
    
    def _parse_textbox(self, text_frame) -> Dict:
        """解析文本框"""
        paragraphs = []
        
        for para in text_frame.paragraphs:
            para_data = {
                "text": para.text,
                "level": para.level,
                "alignment": str(para.alignment) if para.alignment else None,
                "runs": []
            }
            
            for run in para.runs:
                run_data = {
                    "text": run.text,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "underline": run.font.underline,
                    "font_size": run.font.size.pt if run.font.size else None,
                    "font_name": run.font.name,  # ⭐ 保留原始字体名称
                    "color": self._parse_color(run.font.color),
                }
                para_data["runs"].append(run_data)
            
            paragraphs.append(para_data)
        
        return {
            "paragraphs": paragraphs,
            "full_text": text_frame.text,
        }
    
    def _parse_color(self, color_format):
        """解析颜色"""
        try:
            if color_format and color_format.type:
                # 尝试获取 RGB 值
                if hasattr(color_format, 'rgb'):
                    rgb = color_format.rgb
                    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        except:
            pass
        return None
    
    def _parse_group(self, group_shape) -> Dict:
        """解析组合形状"""
        elements = []
        
        for shape in group_shape.shapes:
            element = self._parse_shape(shape, 0)
            if element:
                elements.append(element)
        
        return {
            "element_count": len(elements),
            "elements": elements,
        }
    
    def save_to_json(self, output_path: Path):
        """
        解析并保存为 JSON 文件
        
        Args:
            output_path: JSON 输出路径
        """
        data = self.parse()
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        
        print(f"✅ JSON 已保存: {output_path}")
        return data


if __name__ == "__main__":
    # 测试
    from pathlib import Path
    
    aim01_path = Path(__file__).parent.parent.parent / "aim" / "aim01.pptx"
    parser = PPTParser(aim01_path)
    
    output_path = Path(__file__).parent.parent.parent / "output" / "aim01_parsed.json"
    parser.save_to_json(output_path)

