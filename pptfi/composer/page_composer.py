"""PageComposer — Flow E 核心：按页组合生成 PPT

用法:
    composer = PageComposer(theme="midnight")
    composer.add_page("title_dark", {"title": "报告标题", "date": "2026-02"})
    composer.add_page("kpi_cards", {"title": "核心指标", "cards": [...]})
    composer.add_page("chart_full", {"title": "图表", "df": df, ...})
    composer.save("output.pptx")
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu

from .layouts import LAYOUT_REGISTRY
from .themes import THEMES, DEFAULT_THEME, _LAYOUT_DEFAULTS
from .helpers import add_page_footer

# 不自动注入 header/footer 的布局
_SKIP_FOOTER_LAYOUTS = frozenset({
    "title_dark", "title_light",
    "conclusion_dark",
    "section_divider",
    "gtm_panels", "gtm_cover", "gtm_toc", "gtm_quilt", "gtm_chart_text",  # GTM 自带骨架
    "research_shell_4_3",
    "dashboard_shell_16_9",
    "factsheet_shell_4_3",
    "summary_shell_16_9",
    "section_cover_4_3",
    "chapter_divider_16_9",
    "profile_factsheet_4_3",
})


class PageComposer:
    """页面级组合 PPT 生成器"""

    def __init__(self, theme=None, slide_width=None, slide_height=None):
        """
        Args:
            theme: 主题名称(str) 或主题字典(dict)，包含 slide_w/slide_h tokens
            slide_width: 自定义幻灯片宽度（优先级最高）
            slide_height: 自定义幻灯片高度（优先级最高）
        """
        if isinstance(theme, str):
            self.theme = THEMES.get(theme, DEFAULT_THEME)
        elif isinstance(theme, dict):
            merged = {**_LAYOUT_DEFAULTS, **DEFAULT_THEME, **theme}
            self.theme = merged
        else:
            self.theme = DEFAULT_THEME

        # 从 theme 读取尺寸，slide_width/slide_height 参数可覆盖
        _sw = slide_width or Inches(self.theme.get("slide_w", 13.333))
        _sh = slide_height or Inches(self.theme.get("slide_h", 7.5))
        self.prs = Presentation()
        self.prs.slide_width = _sw
        self.prs.slide_height = _sh

        self._page_count = 0

    @property
    def available_layouts(self):
        """返回所有已注册的布局名称"""
        return list(LAYOUT_REGISTRY.keys())

    def add_page(self, layout_name: str, data: dict):
        """添加一页

        Args:
            layout_name: 布局名称，必须在 LAYOUT_REGISTRY 中
            data: 该布局所需的数据字典
                  可选字段 footnote (str): 数据来源/脚注，自动传给 footer
                  可选字段 source (str): 兼容回退，footnote 优先

        Returns:
            self (支持链式调用)
        """
        if layout_name not in LAYOUT_REGISTRY:
            available = ", ".join(LAYOUT_REGISTRY.keys())
            raise ValueError(f"未知布局 '{layout_name}'。可用布局: {available}")

        layout_fn = LAYOUT_REGISTRY[layout_name]

        # 添加空白幻灯片
        blank_layout = self._get_blank_layout()
        slide = self.prs.slides.add_slide(blank_layout)

        # 执行布局函数
        layout_fn(slide, data, self.theme)
        self._page_count += 1

        # 自动注入 footer（非封面/结论/分隔页）
        if layout_name not in _SKIP_FOOTER_LAYOUTS:
            footer_text = data.get("footnote", "") or data.get("source", "")
            add_page_footer(
                slide, self.theme,
                source=footer_text,
                page_num=self._page_count,
            )

        print(f"  [+] Page {self._page_count}: {layout_name}", file=sys.stderr)
        return self

    def add_custom_page(self, render_fn, data=None):
        """添加自定义页面 — render_fn(slide, theme) 自行绘制

        Args:
            render_fn: callable(slide, theme) 或 callable(slide, data, theme)
            data: 传给 render_fn 的数据
        """
        blank_layout = self._get_blank_layout()
        slide = self.prs.slides.add_slide(blank_layout)

        if data is not None:
            render_fn(slide, data, self.theme)
        else:
            render_fn(slide, self.theme)

        self._page_count += 1
        print(f"  [+] Page {self._page_count}: custom", file=sys.stderr)
        return self

    def save(self, path, lint=True):
        """保存 PPT

        Args:
            path: 输出文件路径
            lint: 是否自动运行 DeckLinter（默认 True）
        """
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        print(f"\n[OK] PPT saved: {path} ({self._page_count} pages)", file=sys.stderr)

        # 自动运行 lint
        if lint:
            try:
                from pptfi.qa.deck_linter import DeckLinter
                linter = DeckLinter(str(path))
                report = linter.run()
                linter.print_summary(report)
            except ImportError:
                pass  # qa 模块尚未安装时静默跳过
            except Exception as e:
                print(f"  [WARN] Lint skipped: {e}", file=sys.stderr)

        return str(path)

    def _get_blank_layout(self):
        """获取空白幻灯片布局"""
        # 优先找名为 'Blank' 的布局
        for layout in self.prs.slide_layouts:
            if layout.name == "Blank":
                return layout
        # 否则用最后一个（通常是空白）
        return self.prs.slide_layouts[-1]

    @staticmethod
    def register_layout(name: str, fn):
        """注册自定义布局函数

        Args:
            name: 布局名称
            fn: callable(slide, data, theme)
        """
        LAYOUT_REGISTRY[name] = fn
