"""图表类布局 — 复用 pptfi ChartBuilder"""

from pptx.util import Inches
from ..helpers import set_slide_bg, add_text, add_rect, add_page_header, setup_content_page


def _add_chart_to_slide(slide, data, position, size, theme=None):
    """调用 pptfi 的 create_combo_chart，自动关联主题配色"""
    from pptfi.chart_builder.api import create_combo_chart
    from pptfi.chart_builder.styles import StyleConfig

    # 如果 data 没有显式 style_config，从 theme 自动构建
    style_config = data.get("style_config")
    if style_config is None and theme is not None:
        scheme = theme.get("chart_default_scheme", "default")
        style_config = StyleConfig(
            color_scheme=scheme,
            line_width_pt=1.5,
            marker_style="none",
        )

    create_combo_chart(
        slide=slide,
        df=data["df"],
        categories_col=data["categories_col"],
        series_config=data["series_config"],
        position=position,
        size=size,
        style_config=style_config,
        layout_config=data.get("layout_config"),
    )


def layout_chart_full(slide, data, theme):
    """全幅图表页 — 标题 + 占满的图表

    data keys:
        title (str): 页面标题
        subtitle (str, optional): 副标题
        df (DataFrame): 图表数据
        categories_col (str): 分类列名
        series_config (list): 系列配置
        style_config (StyleConfig, optional): 样式
        layout_config (ChartLayoutConfig, optional): 布局
        insight (str, optional): 结论先行文本（2-3行，副标题与图表之间）
        footnote (str, optional): 数据来源/脚注（auto-footer 渲染）
    """
    m, sw, content_y, content_w, footer_y = setup_content_page(slide, data, theme)

    subtitle_h = 0
    if data.get("subtitle"):
        chart_sub_size = theme.get("chart_subtitle_size", 11)
        add_text(slide, data["subtitle"],
                 x=m, y=content_y - 0.05, w=content_w, h=0.3,
                 font_size=chart_sub_size, color=theme["text_muted"],
                 font_name=theme["body_font"])
        subtitle_h = 0.35

    # insight 插槽
    insight_h = 0
    if data.get("insight"):
        add_text(slide, data["insight"],
                 x=m, y=content_y + subtitle_h, w=content_w, h=0.50,
                 font_size=theme.get("body_size", 12),
                 color=theme["text_muted"], font_name=theme["body_font"])
        insight_h = 0.55

    # 图表
    chart_top = content_y + subtitle_h + insight_h
    chart_h = footer_y - chart_top - 0.25
    _add_chart_to_slide(
        slide, data,
        position=(Inches(m), Inches(chart_top)),
        size=(Inches(content_w), Inches(chart_h)),
        theme=theme,
    )


def layout_chart_text(slide, data, theme):
    """左图表 + 右文字 — 图文并排

    data keys:
        title (str): 页面标题
        df, categories_col, series_config, style_config, layout_config: 图表参数
        text_title (str, optional): 右侧小标题
        text_body (str): 右侧正文
        text_bullets (list[str], optional): 右侧要点列表
        insight (str, optional): 结论先行文本（header 之后、图文之前）
        footnote (str, optional): 数据来源/脚注（auto-footer 渲染）
    """
    m, sw, content_y, content_w, footer_y = setup_content_page(slide, data, theme)

    # insight 插槽
    insight_h = 0
    if data.get("insight"):
        add_text(slide, data["insight"],
                 x=m, y=content_y, w=content_w, h=0.50,
                 font_size=theme.get("body_size", 12),
                 color=theme["text_muted"], font_name=theme["body_font"])
        insight_h = 0.55

    # 左侧图表 (占 58%)
    chart_w = content_w * 0.58
    chart_top = content_y + insight_h + 0.05
    chart_h = footer_y - chart_top - 0.25
    _add_chart_to_slide(
        slide, data,
        position=(Inches(m), Inches(chart_top)),
        size=(Inches(chart_w), Inches(chart_h)),
        theme=theme,
    )

    # 右侧文字
    text_x = m + chart_w + 0.4
    text_w = sw - text_x - m

    if data.get("text_title"):
        add_text(slide, data["text_title"],
                 x=text_x, y=content_y + 0.1, w=text_w, h=0.5,
                 font_size=theme["subtitle_size"], color=theme["primary"],
                 bold=True, font_name=theme["header_font"])

    body_y = content_y + 0.7 if data.get("text_title") else content_y + 0.1
    add_text(slide, data.get("text_body", ""),
             x=text_x, y=body_y, w=text_w, h=2.5,
             font_size=theme["body_size"], color=theme["text_dark"],
             font_name=theme["body_font"])

    if data.get("text_bullets"):
        bullet_y = body_y + 2.8
        for j, bullet in enumerate(data["text_bullets"]):
            add_text(slide, f"  {bullet}",
                     x=text_x, y=bullet_y + j * 0.45, w=text_w, h=0.4,
                     font_size=theme["body_size"] - 1, color=theme["text_dark"],
                     font_name=theme["body_font"])


def layout_two_charts(slide, data, theme):
    """双图表页 — 左右各一个图表

    data keys:
        title (str): 页面标题
        left (dict): 左图表 (df, categories_col, series_config, ...)
        right (dict): 右图表 (同上)
        left_title (str, optional): 左图小标题
        right_title (str, optional): 右图小标题
        insight (str, optional): 结论先行文本（header 之后、图表之前）
        footnote (str, optional): 数据来源/脚注（auto-footer 渲染）
    """
    m, sw, content_y, content_w, footer_y = setup_content_page(slide, data, theme)

    # insight 插槽
    insight_h = 0
    if data.get("insight"):
        add_text(slide, data["insight"],
                 x=m, y=content_y, w=content_w, h=0.50,
                 font_size=theme.get("body_size", 12),
                 color=theme["text_muted"], font_name=theme["body_font"])
        insight_h = 0.55

    half_w = (content_w - 0.4) / 2
    chart_top = content_y + insight_h + 0.05
    chart_h = footer_y - chart_top - 0.25

    # 小标题
    sub_base_y = content_y + insight_h
    if data.get("left_title"):
        chart_sub_size = theme.get("chart_subtitle_size", 11)
        add_text(slide, data["left_title"],
                 x=m, y=sub_base_y - 0.05, w=half_w, h=0.25,
                 font_size=chart_sub_size, color=theme["text_muted"],
                 font_name=theme["body_font"])
        chart_top = sub_base_y + 0.25
        chart_h = footer_y - chart_top - 0.25

    if data.get("right_title"):
        chart_sub_size = theme.get("chart_subtitle_size", 11)
        add_text(slide, data["right_title"],
                 x=m + half_w + 0.4, y=sub_base_y - 0.05, w=half_w, h=0.25,
                 font_size=chart_sub_size, color=theme["text_muted"],
                 font_name=theme["body_font"])

    # 左图表
    _add_chart_to_slide(
        slide, data["left"],
        position=(Inches(m), Inches(chart_top)),
        size=(Inches(half_w), Inches(chart_h)),
        theme=theme,
    )

    # 右图表
    _add_chart_to_slide(
        slide, data["right"],
        position=(Inches(m + half_w + 0.4), Inches(chart_top)),
        size=(Inches(half_w), Inches(chart_h)),
        theme=theme,
    )


def layout_two_charts_vertical(slide, data, theme):
    """上下双图 — 共享日期轴，最大化水平空间

    data keys:
        title (str): 页面标题
        top (dict): 上图表 (df, categories_col, series_config, ...)
        bottom (dict): 下图表 (同上)
        top_title (str, optional): 上图小标题
        bottom_title (str, optional): 下图小标题
        insight (str, optional): 结论先行文本（header 与上图之间）
        footnote (str, optional): 数据来源/脚注（auto-footer 渲染）
    """
    m, sw, content_y, content_w, footer_y = setup_content_page(slide, data, theme)

    # insight 插槽
    insight_h = 0
    if data.get("insight"):
        add_text(slide, data["insight"],
                 x=m, y=content_y, w=content_w, h=0.50,
                 font_size=theme.get("body_size", 12),
                 color=theme["text_muted"], font_name=theme["body_font"])
        insight_h = 0.55

    chart_area_top = content_y + insight_h
    chart_area_h = footer_y - chart_area_top - 0.25
    gap = 0.15
    each_h = (chart_area_h - gap) / 2

    # 上图小标题
    sub_h = 0
    chart_sub_size = theme.get("chart_subtitle_size", 11)
    if data.get("top_title"):
        add_text(slide, data["top_title"],
                 x=m, y=chart_area_top, w=content_w, h=0.25,
                 font_size=chart_sub_size, color=theme["text_muted"],
                 font_name=theme["body_font"])
        sub_h = 0.28
        each_h = (chart_area_h - gap - sub_h * 2) / 2

    top_chart_y = chart_area_top + sub_h
    _add_chart_to_slide(
        slide, data["top"],
        position=(Inches(m), Inches(top_chart_y)),
        size=(Inches(content_w), Inches(each_h)),
        theme=theme,
    )

    # 下图小标题
    bottom_label_y = top_chart_y + each_h + gap
    if data.get("bottom_title"):
        add_text(slide, data["bottom_title"],
                 x=m, y=bottom_label_y, w=content_w, h=0.25,
                 font_size=chart_sub_size, color=theme["text_muted"],
                 font_name=theme["body_font"])
        bottom_chart_y = bottom_label_y + sub_h
    else:
        bottom_chart_y = bottom_label_y

    _add_chart_to_slide(
        slide, data["bottom"],
        position=(Inches(m), Inches(bottom_chart_y)),
        size=(Inches(content_w), Inches(each_h)),
        theme=theme,
    )


def layout_chart_table(slide, data, theme):
    """左图表 + 右表格 — 趋势可视化 + 截面数据快照

    data keys:
        title (str): 页面标题
        df (DataFrame): 图表数据
        categories_col (str): 分类列名
        series_config (list): 系列配置
        style_config (StyleConfig, optional): 图表样式
        layout_config (ChartLayoutConfig, optional): 图表布局
        headers (list[str]): 表格表头
        rows (list[list]): 表格数据行
        table_df (DataFrame, optional): 若提供，自动转为 headers/rows
        insight (str, optional): 结论先行文本
        chart_ratio (float, optional): 图表宽度占比（默认 0.62）
        footnote (str, optional): 数据来源/脚注（auto-footer 渲染）
    """
    from ..helpers import add_table

    m, sw, content_y, content_w, footer_y = setup_content_page(slide, data, theme)

    # insight 插槽
    insight_h = 0
    if data.get("insight"):
        add_text(slide, data["insight"],
                 x=m, y=content_y, w=content_w, h=0.50,
                 font_size=theme.get("body_size", 12),
                 color=theme["text_muted"], font_name=theme["body_font"])
        insight_h = 0.55

    area_top = content_y + insight_h
    area_h = footer_y - area_top - 0.25

    # 图表/表格分割
    chart_ratio = data.get("chart_ratio", 0.62)
    gap = 0.4
    chart_w = content_w * chart_ratio
    table_w = content_w - chart_w - gap

    # 左侧图表
    _add_chart_to_slide(
        slide, data,
        position=(Inches(m), Inches(area_top)),
        size=(Inches(chart_w), Inches(area_h)),
        theme=theme,
    )

    # 右侧表格 — 自动从 table_df 转换
    import pandas as pd
    headers = data.get("headers")
    rows = data.get("rows")
    if not rows and "table_df" in data:
        tdf = data["table_df"]
        if isinstance(tdf, pd.DataFrame):
            headers = headers or list(tdf.columns)
            rows = tdf.values.tolist()
    if not rows and "df" in data:
        tdf = data["df"]
        if isinstance(tdf, pd.DataFrame):
            headers = headers or list(tdf.columns)
            rows = tdf.values.tolist()

    if headers and rows:
        table_x = m + chart_w + gap
        n_rows = len(rows)
        table_h = min(n_rows * 0.45 + 0.45, area_h)
        add_table(
            slide,
            headers=headers,
            rows=rows,
            x=table_x, y=area_top, w=table_w, h=table_h,
            theme=theme,
            font_name=theme["body_font"],
            font_size=10,
            ib_style=True,
        )
