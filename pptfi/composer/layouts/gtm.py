"""GTM 页面布局族 — 通用市场报告版式（General Theme for Markets），
提炼自国际投行市场指南类公开报告的页面骨架。

页面构成（对应该类报告的惯用骨架）：
- 左上角蓝色箭头 + 页标题（20pt 左对齐）+ 全宽细分隔线
- 右上角页签：[GTM | 市场 | 页码] 三格框
- 左侧竖向章节签（按章节配色，白字竖排）
- 1 / 2 / 1+2 三种面板编排；面板 = 粗体小标题 + 灰色单位行 + pptchartengine 图表
- 面板内可叠加迷你数据表（GTM 的 Avg/2Q25 小表，行文字可按系列着色）
- 底部来源行（7.5pt 灰）+ 右下品牌字 + 左下页码

面板里的 ``chart`` 字段就是 pptchartengine 的声明式 spec（data 支持内联
dict / CSV 路径 / DataFrame），模型只需产出一个 job JSON 即可编排整页。
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ..helpers import add_line, add_text

# GTM 页面几何（16:9，13.333 × 7.5 英寸）
_MARGIN_L = 0.85
_MARGIN_R = 0.45
_TITLE_Y = 0.38
_RULE_Y = 0.96
_CONTENT_TOP = 1.18
_CONTENT_BOTTOM = 6.82
_FOOTER_Y = 6.98
_GAP = 0.42          # 双栏面板间距
_DIVIDER_X_FRACTION = 0.495

FONT = "微软雅黑"

# 骨架配色默认值集中在 palette.GTM_CHROME；主题可逐项覆盖（如 able_warm）。
from ..palette import GTM_CHROME as _CHROME


def _c(theme, key):
    """读骨架色：优先主题 token，回退 palette.GTM_CHROME 默认。"""
    if theme and key in theme and theme[key] is not None:
        return theme[key]
    return _CHROME[key]


def _section_color(theme, section, override=None):
    """章节签/章节强调色：override > 主题 section_colors > 主题 chart_accent。"""
    sect_map = _c(theme, "section_colors")
    return (override
            or sect_map.get(str(section).lower())
            or sect_map.get(str(section))
            or _c(theme, "chart_accent"))


def _sw(theme):
    """当前画布宽度（英寸）。随 aspect_ratio 自适应：16:9→13.333，4:3→10.0。"""
    return float((theme or {}).get("slide_w", 13.333))


def layout_gtm_panels(slide, data, theme):
    """GTM 面板页。

    data keys:
        title (str): 页标题
        section (str, optional): 章节名（决定左侧竖签文字与颜色）
        section_color (str, optional): 覆盖章节色
        tag (str, optional): 角标市场代码，默认 "GTM"
        page_num (int/str, optional): 角标页码
        panels (list[dict]): 1~3 个面板
            {title, subtitle, chart: <pptchartengine spec>, table: {...}}
            3 个面板时第 1 个占左半全高，2/3 在右侧上下排
        source (str, optional): 底部来源行
        brand (str, optional): 右下品牌字
    """
    _draw_chrome(slide, data, theme)

    sw = _sw(theme)
    panels = data.get("panels") or []
    rects = _panel_rects(len(panels), sw)
    if len(panels) >= 2:
        _draw_divider(slide, sw, theme)

    for panel, rect in zip(panels, rects):
        _render_panel(slide, panel, rect, theme)


# ============================================================================
# 页面骨架
# ============================================================================

def _draw_chrome(slide, data, theme):
    sw = _sw(theme)

    # 箭头标（强调色跟随主题）
    chevron = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(0.22), Inches(_TITLE_Y + 0.07), Inches(0.42), Inches(0.30))
    chevron.fill.solid()
    chevron.fill.fore_color.rgb = RGBColor.from_string(_c(theme, "chart_accent"))
    chevron.line.fill.background()
    chevron.shadow.inherit = False

    # 页标题
    add_text(slide, str(data.get("title", "")),
             x=_MARGIN_L, y=_TITLE_Y, w=sw - 4.0, h=0.5,
             font_size=20, font_name=FONT, color=_c(theme, "ink"), bold=True)

    # 全宽细分隔线
    add_line(slide, _MARGIN_L, _RULE_Y, sw - _MARGIN_L - _MARGIN_R, color=_c(theme, "rule"), width=1)

    # 右上角页签 [GTM | 市场 | 页码]
    cells = [str(data.get("tag", "GTM"))]
    if data.get("market"):
        cells.append(str(data["market"]))
    if data.get("page_num") is not None:
        cells.append(str(data["page_num"]))
    _draw_tag_boxes(slide, cells, right=sw - _MARGIN_R, y=_TITLE_Y + 0.06, theme=theme)

    # 左侧竖向章节签
    section = data.get("section")
    if section:
        _draw_section_tab(slide, str(section),
                          _section_color(theme, section, data.get("section_color")), theme)

    # 底部来源行 + 品牌 + 页码
    if data.get("source"):
        add_text(slide, str(data["source"]),
                 x=_MARGIN_L, y=_FOOTER_Y, w=8.6, h=0.42,
                 font_size=7.5, font_name=FONT, color=_c(theme, "source"))
    if data.get("brand"):
        add_text(slide, str(data["brand"]),
                 x=sw - 3.6, y=_FOOTER_Y, w=3.1, h=0.34,
                 font_size=13, font_name=FONT, color=_c(theme, "ink"), bold=True, align="right")
    if data.get("page_num") is not None:
        add_text(slide, str(data["page_num"]),
                 x=0.22, y=7.08, w=0.5, h=0.3,
                 font_size=10, font_name=FONT, color=_c(theme, "ink_muted"))


def _draw_tag_boxes(slide, cells, *, right, y, theme=None):
    box_h = 0.30
    paper, rule, ink = _c(theme, "paper"), _c(theme, "rule"), _c(theme, "ink")
    widths = [max(0.52, 0.16 + 0.105 * len(c)) for c in cells]
    x = right - sum(widths)
    for text, w in zip(cells, widths):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if text == cells[0] else MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(box_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(paper)
        shape.line.color.rgb = RGBColor.from_string(rule)
        shape.line.width = Pt(0.75)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.text = text
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.size = Pt(10)
        run.font.name = FONT
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(ink)
        x += w


def _draw_section_tab(slide, text, color, theme=None):
    w, h = 0.34, max(1.1, 0.32 + 0.21 * len(text))
    tab = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.18), Inches(_CONTENT_TOP + 0.25), Inches(w), Inches(h))
    tab.fill.solid()
    tab.fill.fore_color.rgb = RGBColor.from_string(color)
    tab.line.fill.background()
    tab.shadow.inherit = False
    tab.rotation = 0
    tf = tab.text_frame
    tf.word_wrap = False
    tf.text = "\n".join(text)  # 中文逐字竖排
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(10)
            run.font.name = FONT
            run.font.bold = True
            run.font.color.rgb = RGBColor.from_string(_c(theme, "paper"))


def _draw_divider(slide, sw=13.333, theme=None):
    x = _MARGIN_L + (sw - _MARGIN_L - _MARGIN_R) * _DIVIDER_X_FRACTION
    line = slide.shapes.add_connector(1, Inches(x), Inches(_CONTENT_TOP),
                                      Inches(x), Inches(_CONTENT_BOTTOM))
    line.line.color.rgb = RGBColor.from_string(_c(theme, "hairline"))
    line.line.width = Pt(0.75)
    line.shadow.inherit = False


# ============================================================================
# 面板
# ============================================================================

def _panel_rects(n, sw=13.333):
    """返回 n 个面板的 (x, y, w, h)。"""
    full_w = sw - _MARGIN_L - _MARGIN_R
    full_h = _CONTENT_BOTTOM - _CONTENT_TOP
    half_w = full_w * _DIVIDER_X_FRACTION - _GAP / 2
    right_x = _MARGIN_L + full_w * _DIVIDER_X_FRACTION + _GAP / 2
    right_w = full_w - full_w * _DIVIDER_X_FRACTION - _GAP / 2

    if n <= 1:
        return [(_MARGIN_L, _CONTENT_TOP, full_w, full_h)]
    if n == 2:
        return [
            (_MARGIN_L, _CONTENT_TOP, half_w, full_h),
            (right_x, _CONTENT_TOP, right_w, full_h),
        ]
    half_h = (full_h - 0.30) / 2
    if n == 3:
        return [
            (_MARGIN_L, _CONTENT_TOP, half_w, full_h),
            (right_x, _CONTENT_TOP, right_w, half_h),
            (right_x, _CONTENT_TOP + half_h + 0.30, right_w, half_h),
        ]
    # 4 面板：2×2 网格
    return [
        (_MARGIN_L, _CONTENT_TOP, half_w, half_h),
        (right_x, _CONTENT_TOP, right_w, half_h),
        (_MARGIN_L, _CONTENT_TOP + half_h + 0.30, half_w, half_h),
        (right_x, _CONTENT_TOP + half_h + 0.30, right_w, half_h),
    ]


def _render_panel(slide, panel, rect, theme):
    from pptchartengine import render_chart

    x, y, w, h = rect
    cursor = y

    if panel.get("title"):
        add_text(slide, str(panel["title"]), x=x, y=cursor, w=w, h=0.28,
                 font_size=12, font_name=FONT, color=_c(theme, "ink"), bold=True)
        cursor += 0.28
    if panel.get("subtitle"):
        add_text(slide, str(panel["subtitle"]), x=x, y=cursor, w=w, h=0.22,
                 font_size=9.5, font_name=FONT, color=_c(theme, "ink_muted"))
        cursor += 0.24

    chart_spec = panel.get("chart")
    if chart_spec:
        spec = dict(chart_spec)
        # 面板标题由页面层渲染，图表内不再重复
        spec.pop("title", None)
        spec.pop("subtitle", None)
        spec["position"] = [x, cursor]
        spec["chart_size"] = [w, y + h - cursor]
        render_chart(slide, spec)

    table = panel.get("table")
    if table:
        _render_mini_table(slide, table, rect=(x, cursor, w, y + h - cursor), theme=theme)


def _render_mini_table(slide, table, *, rect, theme=None):
    """面板内迷你数据表（GTM 的 Avg / 2Q25 小表）。

    table keys:
        columns (list[str]): 表头（第一格通常留空）
        rows (list[list]): 数据行，首列为行标签
        row_colors (list[str], optional): 各行标签颜色（与系列色呼应）
        at (list[float], optional): 表左上角在面板图区内的比例坐标，默认 [0.35, 0.02]
        width (float, optional): 表宽（英寸），默认按列数估算
    """
    px, py, pw, ph = rect
    columns = table.get("columns") or []
    rows = table.get("rows") or []
    if not rows or not columns:
        return
    row_colors = table.get("row_colors") or []
    fx, fy = table.get("at", [0.35, 0.02])
    n_cols = len(columns)
    n_rows = len(rows) + 1
    width = float(table.get("width", min(pw * 0.62, 0.85 + 0.75 * (n_cols - 1))))
    height = 0.21 * n_rows
    paper, ink, ink_muted = _c(theme, "paper"), _c(theme, "ink"), _c(theme, "ink_muted")

    gfx = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(px + pw * fx), Inches(py + ph * fy),
        Inches(width), Inches(height))
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False
    # 首列放行标签，列宽加权
    tbl.columns[0].width = Inches(width * 0.40)
    for c in range(1, n_cols):
        tbl.columns[c].width = Inches(width * 0.60 / (n_cols - 1))

    def style(cell, text, *, bold, color, align):
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(paper)
        cell.margin_left = Pt(3)
        cell.margin_right = Pt(3)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)
        tf = cell.text_frame
        tf.text = str(text)
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = Pt(8.5)
        run.font.name = FONT
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)

    for c, header in enumerate(columns):
        style(tbl.cell(0, c), header, bold=True, color=ink,
              align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    for r, row in enumerate(rows, start=1):
        label_color = (row_colors[r - 1] if r - 1 < len(row_colors) and row_colors[r - 1]
                       else ink).lstrip("#").upper()
        for c, value in enumerate(row):
            style(tbl.cell(r, c), value,
                  bold=(c == 0), color=label_color if c == 0 else ink_muted,
                  align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


# ============================================================================
# GTM 封面页
# ============================================================================

def layout_gtm_cover(slide, data, theme):
    """GTM 封面：大箭头 + 主标题 + 副标题 + 数据截止日 + 品牌。

    data keys: title, subtitle, date, brand, market
    """
    # 顶部品牌条
    sw = _sw(theme)
    ink, ink_muted = _c(theme, "ink"), _c(theme, "ink_muted")
    add_text(slide, str(data.get("brand", "")),
             x=_MARGIN_L, y=0.55, w=8.0, h=0.45,
             font_size=17, font_name=FONT, color=ink, bold=True)
    add_line(slide, _MARGIN_L, 1.12, sw - _MARGIN_L - _MARGIN_R, color=_c(theme, "rule"), width=1)

    # 大箭头组（GTM 封面的标志图形，三段渐深随主题）
    for dx, alpha_color in zip((0.0, 0.55, 1.1), _c(theme, "cover_chevrons")):
        chevron = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, Inches(_MARGIN_L + dx), Inches(2.35), Inches(1.05), Inches(1.5))
        chevron.fill.solid()
        chevron.fill.fore_color.rgb = RGBColor.from_string(alpha_color)
        chevron.line.fill.background()
        chevron.shadow.inherit = False

    # 主标题 / 副标题
    text_w = sw - _MARGIN_L - _MARGIN_R
    add_text(slide, str(data.get("title", "")),
             x=_MARGIN_L, y=4.15, w=text_w, h=1.0,
             font_size=34, font_name=FONT, color=ink, bold=True)
    if data.get("subtitle"):
        add_text(slide, str(data["subtitle"]),
                 x=_MARGIN_L, y=5.05, w=text_w, h=0.5,
                 font_size=16, font_name=FONT, color=ink_muted)

    # 市场 + 数据截止日
    meta_parts = []
    if data.get("market"):
        meta_parts.append(str(data["market"]))
    if data.get("date"):
        meta_parts.append(f"数据截至 {data['date']}")
    if meta_parts:
        add_text(slide, "  |  ".join(meta_parts),
                 x=_MARGIN_L, y=6.55, w=10.0, h=0.4,
                 font_size=12, font_name=FONT, color=ink_muted)


# ============================================================================
# GTM 目录页
# ============================================================================

def layout_gtm_toc(slide, data, theme):
    """GTM 目录：章节色块 + 标题…页码 两栏列表。

    data keys:
        title (str): 默认 "目录"
        items (list): [{"section": str, "entries": [{"title", "page"}]}]
                      省略时由 deck 工作流按后续页面自动生成
    """
    _draw_chrome(slide, {**data, "title": data.get("title", "目录"),
                         "section": None, "page_num": data.get("page_num")}, theme)

    items = data.get("items") or []
    col_w = (_sw(theme) - _MARGIN_L - _MARGIN_R - _GAP) / 2
    col_x = [_MARGIN_L, _MARGIN_L + col_w + _GAP]
    cursor = [_CONTENT_TOP + 0.1, _CONTENT_TOP + 0.1]
    col = 0

    for group in items:
        section = str(group.get("section", ""))
        entries = group.get("entries") or []
        block_h = 0.42 + 0.30 * len(entries)
        if cursor[col] + block_h > _CONTENT_BOTTOM and col == 0:
            col = 1
        x, y = col_x[col], cursor[col]

        color = _section_color(theme, section)
        chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x), Inches(y + 0.04), Inches(0.16), Inches(0.22))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor.from_string(color)
        chip.line.fill.background()
        chip.shadow.inherit = False
        add_text(slide, section, x=x + 0.26, y=y, w=col_w - 0.3, h=0.3,
                 font_size=13, font_name=FONT, color=_c(theme, "ink"), bold=True)
        y += 0.38

        for entry in entries:
            add_text(slide, str(entry.get("title", "")), x=x + 0.26, y=y, w=col_w - 1.0, h=0.26,
                     font_size=10.5, font_name=FONT, color=_c(theme, "rule"))
            add_text(slide, str(entry.get("page", "")), x=x + col_w - 0.7, y=y, w=0.6, h=0.26,
                     font_size=10.5, font_name=FONT, color=_c(theme, "ink_muted"), align="right")
            y += 0.30
        cursor[col] = y + 0.22


# ============================================================================
# GTM 资产收益 quilt 矩阵（GTM 最具标志性的一页）
# ============================================================================

def layout_gtm_quilt(slide, data, theme):
    """资产收益排序矩阵：每列一年，按当年收益降序排格，单元格按资产固定着色。

    data keys:
        title/subtitle/section/source/...: 页面骨架字段
        df (DataFrame): 长表，含 年份列/资产列/收益列
        year_col/asset_col/value_col: 列名（默认 年份/资产/收益率）
        colors (dict, optional): {资产: "#RRGGBB"} 指定配色
    """
    _draw_chrome(slide, data, theme)

    df = data["df"]
    year_col = data.get("year_col", "年份")
    asset_col = data.get("asset_col", "资产")
    value_col = data.get("value_col", "收益率")

    years = sorted(df[year_col].astype(str).unique())
    assets = list(df[asset_col].unique())
    quilt_pal = _c(theme, "quilt_palette")
    paper, ink = _c(theme, "paper"), _c(theme, "ink")
    palette = {a: (data.get("colors") or {}).get(a, quilt_pal[i % len(quilt_pal)])
               for i, a in enumerate(assets)}

    n_rows = len(assets) + 1
    n_cols = len(years)
    x, y = _MARGIN_L, _CONTENT_TOP + 0.15
    w = _sw(theme) - _MARGIN_L - _MARGIN_R
    h = _CONTENT_BOTTOM - y

    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False

    def style(cell, text, *, fill, color, bold, size=9):
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(str(fill).lstrip("#").upper())
        for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(cell, margin, Pt(2))
        tf = cell.text_frame
        tf.word_wrap = True
        tf.text = ""
        lines = text if isinstance(text, list) else [text]
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(line)
            run.font.size = Pt(size if i == 0 else size - 0.5)
            run.font.name = FONT
            run.font.bold = bold and i == 0
            run.font.color.rgb = RGBColor.from_string(color)

    for c, year in enumerate(years):
        style(tbl.cell(0, c), year, fill=paper, color=ink, bold=True, size=10)
        sub = df[df[year_col].astype(str) == year].sort_values(value_col, ascending=False)
        for r, (_, row) in enumerate(sub.iterrows(), start=1):
            if r >= n_rows:
                break
            asset = row[asset_col]
            value = float(row[value_col])
            style(tbl.cell(r, c), [str(asset), f"{value*100:+.1f}%"],
                  fill=palette.get(asset, _c(theme, "ink_muted")), color=paper, bold=True)


# ============================================================================
# GTM 左图右点评页
# ============================================================================

def layout_gtm_chart_text(slide, data, theme):
    """左侧面板图 + 右侧观点要点（投研「图 + 点评」版式）。

    data keys:
        panel (dict): 同 gtm_panels 的面板（title/subtitle/chart/table）
        bullets (list[str]): 右侧要点
        bullets_title (str): 右栏标题，默认 "要点"
        其余为页面骨架字段
    """
    _draw_chrome(slide, data, theme)
    sw = _sw(theme)
    _draw_divider(slide, sw, theme)

    full_w = sw - _MARGIN_L - _MARGIN_R
    left_w = full_w * _DIVIDER_X_FRACTION - _GAP / 2
    right_x = _MARGIN_L + full_w * _DIVIDER_X_FRACTION + _GAP / 2
    right_w = full_w - full_w * _DIVIDER_X_FRACTION - _GAP / 2

    panel = data.get("panel") or {}
    _render_panel(slide, panel, (_MARGIN_L, _CONTENT_TOP, left_w, _CONTENT_BOTTOM - _CONTENT_TOP), theme)

    section = str(data.get("section") or "")
    accent = _section_color(theme, section, data.get("section_color"))
    y = _CONTENT_TOP
    add_text(slide, str(data.get("bullets_title", "要点")),
             x=right_x, y=y, w=right_w, h=0.3,
             font_size=12, font_name=FONT, color=_c(theme, "ink"), bold=True)
    y += 0.45

    for bullet in data.get("bullets") or []:
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(right_x), Inches(y + 0.06), Inches(0.11), Inches(0.11))
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor.from_string(accent)
        marker.line.fill.background()
        marker.shadow.inherit = False
        box = add_text(slide, str(bullet),
                       x=right_x + 0.24, y=y, w=right_w - 0.3, h=0.9,
                       font_size=11, font_name=FONT, color=_c(theme, "rule"))
        # 按文字长度估算行数推进光标（每行约 28 个全角字符）
        lines = max(1, (len(str(bullet)) + 27) // 28)
        y += 0.06 + 0.26 * lines + 0.18
