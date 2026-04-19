"""绘图辅助函数 — 在 slide 上添加文本、形状、表格、图片、页眉/页脚"""

import io
import re
from pathlib import Path
from urllib.request import urlopen, Request
import base64

from PIL import Image as PILImage
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_str: str) -> RGBColor:
    """'1E2761' -> RGBColor"""
    return RGBColor.from_string(hex_str)


def set_slide_bg(slide, color_hex: str):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_hex)


def add_text(slide, text: str, x, y, w, h, *,
             font_size=14, font_name="微软雅黑", color="212121",
             bold=False, italic=False,
             align="left", valign="top",
             word_wrap=True, margin_pt=None):
    """在 slide 上添加文本框，返回 shape"""
    left = Inches(x) if isinstance(x, (int, float)) else x
    top = Inches(y) if isinstance(y, (int, float)) else y
    width = Inches(w) if isinstance(w, (int, float)) else w
    height = Inches(h) if isinstance(h, (int, float)) else h

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    if margin_pt is not None:
        tf.margin_left = Pt(margin_pt)
        tf.margin_right = Pt(margin_pt)
        tf.margin_top = Pt(margin_pt)
        tf.margin_bottom = Pt(margin_pt)

    # 段落
    p = tf.paragraphs[0]
    p.text = text

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)

    # 字体
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.color.rgb = _rgb(color)
    run.font.bold = bold
    run.font.italic = italic

    # 垂直对齐
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        anchor_map = {"top": "t", "middle": "ctr", "bottom": "b"}
        bodyPr.set("anchor", anchor_map.get(valign, "t"))

    return txBox


def add_rich_text(slide, parts, x, y, w, h, *,
                  font_name="微软雅黑", align="left", valign="top",
                  line_spacing=1.2, margin_pt=None):
    """添加富文本框 — parts 为 [{"text": ..., "size": ..., "color": ..., "bold": ...}, ...]
    支持 newline=True 换行"""
    left = Inches(x) if isinstance(x, (int, float)) else x
    top = Inches(y) if isinstance(y, (int, float)) else y
    width = Inches(w) if isinstance(w, (int, float)) else w
    height = Inches(h) if isinstance(h, (int, float)) else h

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if margin_pt is not None:
        tf.margin_left = Pt(margin_pt)
        tf.margin_right = Pt(margin_pt)
        tf.margin_top = Pt(margin_pt)
        tf.margin_bottom = Pt(margin_pt)

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    p = tf.paragraphs[0]
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)

    for i, part in enumerate(parts):
        if part.get("newline") and i > 0:
            p = tf.add_paragraph()
            p.alignment = align_map.get(align, PP_ALIGN.LEFT)
            if line_spacing != 1.0:
                p.line_spacing = Pt(int(part.get("size", 14) * line_spacing))

        run = p.add_run()
        run.text = part["text"]
        run.font.size = Pt(part.get("size", 14))
        run.font.name = part.get("font", font_name)
        run.font.color.rgb = _rgb(part.get("color", "212121"))
        run.font.bold = part.get("bold", False)
        run.font.italic = part.get("italic", False)

    return txBox


def add_shape(slide, shape_type, x, y, w, h, *,
              fill=None, line_color=None, line_width=None):
    """添加形状"""
    left = Inches(x) if isinstance(x, (int, float)) else x
    top = Inches(y) if isinstance(y, (int, float)) else y
    width = Inches(w) if isinstance(w, (int, float)) else w
    height = Inches(h) if isinstance(h, (int, float)) else h

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = _rgb(line_color)
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()

    return shape


def add_rect(slide, x, y, w, h, **kwargs):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, **kwargs)


def add_rounded_rect(slide, x, y, w, h, **kwargs):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, **kwargs)


def add_line(slide, x, y, w, *, color="E0E0E0", width=1):
    """添加水平线"""
    left = Inches(x) if isinstance(x, (int, float)) else x
    top = Inches(y) if isinstance(y, (int, float)) else y
    w_emu = Inches(w) if isinstance(w, (int, float)) else w

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w_emu, Pt(width))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


# ── Picture ────────────────────────────────────────────────────────────

def _resolve_image_source(image_src):
    """将图片来源统一为 (BytesIO stream, (width, height)) 元组

    支持三种来源:
    - 本地文件路径 (str / Path)
    - URL (http:// / https://)
    - base64 字符串 ("data:image/...;base64,..." 或裸 base64)

    返回 (stream, (img_w, img_h))。stream 可直接传给 python-pptx。
    """
    src = str(image_src)

    if src.startswith(("http://", "https://")):
        req = Request(src, headers={"User-Agent": "pptfi/1.0"})
        resp = urlopen(req, timeout=30)
        raw_bytes = resp.read()
    elif src.startswith("data:"):
        # data URI: "data:image/png;base64,iVBOR..."
        b64_data = src.split(";base64,", 1)[1]
        raw_bytes = base64.b64decode(b64_data)
    else:
        # 先尝试本地文件，不存在则当作裸 base64
        p = Path(src).expanduser().resolve()
        if p.exists():
            raw_bytes = p.read_bytes()
        else:
            try:
                raw_bytes = base64.b64decode(src)
            except Exception:
                raise FileNotFoundError(f"图片文件不存在: {p}")

    # 用独立 BytesIO 读尺寸，避免 PIL 关闭影响主 stream
    with PILImage.open(io.BytesIO(raw_bytes)) as pil:
        size = pil.size

    return io.BytesIO(raw_bytes), size


def add_picture(slide, image_src, x, y, w, h, *, sizing="stretch"):
    """在 slide 上插入图片，返回 Picture shape

    Args:
        image_src: 图片来源 — 本地路径 / URL / base64 字符串
        x, y, w, h: 目标区域（英寸或 Emu）
        sizing:
            "stretch"  — 拉伸填满目标区域（可能变形）
            "contain"  — 等比缩放适配，居中显示（可能留白）
            "cover"    — 等比缩放裁切填满（无变形无留白）
    """
    left = Inches(x) if isinstance(x, (int, float)) else x
    top = Inches(y) if isinstance(y, (int, float)) else y
    width = Inches(w) if isinstance(w, (int, float)) else w
    height = Inches(h) if isinstance(h, (int, float)) else h

    stream, (img_w, img_h) = _resolve_image_source(image_src)

    if sizing == "stretch":
        return slide.shapes.add_picture(stream, left, top, width, height)

    target_ratio = w / h if isinstance(w, (int, float)) and isinstance(h, (int, float)) else (width / height)
    img_ratio = img_w / img_h

    if sizing == "contain":
        if img_ratio > target_ratio:
            # 图片更宽 → 宽度撑满，高度居中
            final_w = w if isinstance(w, (int, float)) else width
            final_h = (w / img_ratio) if isinstance(w, (int, float)) else int(width / img_ratio)
        else:
            # 图片更高 → 高度撑满，宽度居中
            final_h = h if isinstance(h, (int, float)) else height
            final_w = (h * img_ratio) if isinstance(h, (int, float)) else int(height * img_ratio)

        if isinstance(w, (int, float)):
            offset_x = x + (w - final_w) / 2
            offset_y = y + (h - final_h) / 2
            return slide.shapes.add_picture(
                stream,
                Inches(offset_x), Inches(offset_y),
                Inches(final_w), Inches(final_h))
        else:
            offset_x = left + (width - final_w) // 2
            offset_y = top + (height - final_h) // 2
            return slide.shapes.add_picture(stream, offset_x, offset_y, final_w, final_h)

    if sizing == "cover":
        # 先拉伸填满目标区域
        pic = slide.shapes.add_picture(stream, left, top, width, height)
        # 通过 srcRect crop 恢复宽高比
        if img_ratio > target_ratio:
            # 图片更宽 → 裁左右
            crop = (1 - target_ratio / img_ratio) / 2
            _set_crop(pic, left=crop, right=crop)
        elif img_ratio < target_ratio:
            # 图片更高 → 裁上下
            crop = (1 - img_ratio / target_ratio) / 2
            _set_crop(pic, top=crop, bottom=crop)
        return pic

    raise ValueError(f"未知 sizing 模式: {sizing}")


def _set_crop(pic, *, left=0, right=0, top=0, bottom=0):
    """通过 OOXML srcRect 设置图片裁切"""
    from pptx.oxml.ns import qn
    from lxml import etree

    blipFill = pic._element.find(qn("p:blipFill"))
    if blipFill is None:
        return

    srcRect = blipFill.find(qn("a:srcRect"))
    if srcRect is None:
        srcRect = etree.SubElement(blipFill, qn("a:srcRect"))

    # OOXML srcRect 用千分比（1000 = 1%，100000 = 100%）
    if left:
        srcRect.set("l", str(int(left * 100000)))
    if right:
        srcRect.set("r", str(int(right * 100000)))
    if top:
        srcRect.set("t", str(int(top * 100000)))
    if bottom:
        srcRect.set("b", str(int(bottom * 100000)))


# ── Content Page Setup ─────────────────────────────────────────────────

def setup_content_page(slide, data, theme):
    """内容页通用初始化：背景 + 页眉标题

    返回常用布局参数元组 (m, sw, content_y, content_w, footer_y)。
    """
    m = theme["margin"]
    sw = theme.get("slide_w", 13.333)
    content_y = theme.get("content_y", 1.00)
    content_w = theme.get("content_w", sw - 2 * m)
    footer_y = theme.get("footer_y", 7.10)

    set_slide_bg(slide, theme["bg_light"])
    add_page_header(slide, data["title"], theme)

    return m, sw, content_y, content_w, footer_y


# ── Page Header / Footer (投行级标准) ──────────────────────────────

def add_page_header(slide, title, theme):
    """投行级标准页眉: 16pt 标题 + 全宽分隔线

    使用 theme tokens: page_title_size, header_y, divider_y, divider_h, content_w
    """
    m = theme["margin"]
    sw = theme.get("slide_w", 13.333)
    page_title_size = theme.get("page_title_size", 16)
    header_y = theme.get("header_y", 0.25)
    divider_y = theme.get("divider_y", 0.80)
    divider_h = theme.get("divider_h", 0.015)
    content_w = theme.get("content_w", sw - 2 * m)

    # 页面标题
    add_text(slide, title,
             x=m, y=header_y, w=content_w, h=0.5,
             font_size=page_title_size, color=theme["text_dark"],
             bold=True, font_name=theme["header_font"])

    # 全宽分隔线
    add_rect(slide, m, divider_y, content_w, divider_h, fill=theme["primary"])


def add_page_footer(slide, theme, source="", page_num=None, brand=""):
    """投行级标准页脚: 左侧来源 8pt + 右侧页码/品牌

    使用 theme tokens: footer_size, footer_y, content_w
    """
    m = theme["margin"]
    sw = theme.get("slide_w", 13.333)
    footer_size = theme.get("footer_size", 8)
    footer_y = theme.get("footer_y", 7.10)
    content_w = theme.get("content_w", sw - 2 * m)

    # 页脚分隔细线
    add_rect(slide, m, footer_y - 0.08, content_w, 0.005, fill=theme["border"])

    # 左侧: 来源
    if source:
        source_text = f"Source: {source}" if not source.startswith("Source") else source
        add_text(slide, source_text,
                 x=m, y=footer_y, w=content_w * 0.7, h=0.25,
                 font_size=footer_size, color=theme["text_muted"],
                 italic=True, font_name=theme["body_font"])

    # 右侧: 页码/品牌
    right_parts = []
    if brand:
        right_parts.append(brand)
    if page_num is not None:
        right_parts.append(str(page_num))
    if right_parts:
        add_text(slide, "  |  ".join(right_parts),
                 x=m + content_w * 0.3, y=footer_y, w=content_w * 0.7, h=0.25,
                 font_size=footer_size, color=theme["text_muted"],
                 align="right", font_name=theme["body_font"])


# ── Table (投行级) ──────────────────────────────────────────────────────

_NUM_PATTERN = re.compile(r'^[\s]*[-+]?[\d,]+\.?\d*\s*[%xX倍万亿]?\s*$')


def _is_numeric_cell(val: str) -> bool:
    """检测单元格值是否为数值型（用于右对齐）"""
    if not val:
        return False
    return bool(_NUM_PATTERN.match(val))


def add_table(slide, headers, rows, x, y, w, h, *,
              theme=None, header_bg=None, header_color="FFFFFF",
              font_name="微软雅黑", font_size=11,
              highlight_row=None, highlight_bg="EBF5FF",
              highlight_cols=None,
              ib_style=True):
    """添加带样式的表格

    Args:
        highlight_cols: 高亮列索引列表 (如 [2, 3])
        ib_style: True 启用投行级样式（无垂直边框、数值右对齐、斑马纹 token）
    """
    left = Inches(x) if isinstance(x, (int, float)) else x
    top = Inches(y) if isinstance(y, (int, float)) else y
    width = Inches(w) if isinstance(w, (int, float)) else w
    height = Inches(h) if isinstance(h, (int, float)) else h

    n_rows = len(rows) + 1
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 均匀列宽
    col_w = int(width / n_cols)
    for j in range(n_cols):
        table.columns[j].width = col_w

    # 颜色 tokens
    hdr_bg = header_bg or (theme.get("table_header_bg", theme["primary"]) if theme else "36454F")
    zebra_even = theme.get("table_zebra_even", "F8FAFC") if theme else "F8FAFC"
    zebra_odd = theme.get("table_zebra_odd", "FFFFFF") if theme else "FFFFFF"
    highlight_col_bg = theme.get("table_zebra_even", "EBF5FF") if theme else "EBF5FF"

    # 预扫描: 检测哪些列是数值列（用于自动右对齐）
    numeric_cols = set()
    if ib_style and rows:
        for j in range(n_cols):
            num_count = sum(1 for row in rows if len(row) > j and _is_numeric_cell(str(row[j] or "")))
            if num_count > len(rows) * 0.5:
                numeric_cols.add(j)

    # 表头
    for j, hdr in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(hdr)
        align = "right" if (ib_style and j in numeric_cols) else "center"
        _style_cell(cell, bg=hdr_bg, color=header_color,
                     font_name=font_name, font_size=font_size, bold=True,
                     align=align)

    # 数据行
    for i, row in enumerate(rows):
        is_highlight = (highlight_row is not None and i == highlight_row)
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell_text = str(val) if val is not None else ""
            cell.text = cell_text

            # 背景色优先级: highlight_row > highlight_cols > zebra
            if is_highlight:
                bg = highlight_bg
            elif highlight_cols and j in highlight_cols:
                bg = highlight_col_bg
            else:
                bg = zebra_even if i % 2 == 0 else zebra_odd

            # 对齐: 数值列右对齐, 首列左对齐, 其余居中
            if ib_style and j in numeric_cols:
                align = "right"
            elif j == 0:
                align = "left"
            else:
                align = "center"

            _style_cell(cell, bg=bg, color="212121",
                         font_name=font_name, font_size=font_size,
                         bold=is_highlight,
                         align=align)

    # 投行级: 移除垂直边框，仅保留水平线
    if ib_style:
        _remove_vertical_borders(table, n_rows, n_cols, hdr_bg)

    return table_shape


def _remove_vertical_borders(table, n_rows, n_cols, header_color):
    """移除垂直边框，仅保留表头上下 + 表格底部水平线"""
    from pptx.oxml.ns import qn
    from lxml import etree

    for i in range(n_rows):
        for j in range(n_cols):
            cell = table.cell(i, j)
            tc = cell._tc
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn("a:tcPr"))

            # 清除左右边框
            for border_name in ["lnL", "lnR"]:
                border = tcPr.find(qn(f"a:{border_name}"))
                if border is not None:
                    tcPr.remove(border)
                border = etree.SubElement(tcPr, qn(f"a:{border_name}"))
                border.set("w", "0")
                noFill = etree.SubElement(border, qn("a:noFill"))

            # 设置水平线: 表头底部 + 表格底部用主色细线，其余极浅
            for border_name in ["lnT", "lnB"]:
                border = tcPr.find(qn(f"a:{border_name}"))
                if border is not None:
                    tcPr.remove(border)

                is_header_bottom = (i == 0 and border_name == "lnB")
                is_table_bottom = (i == n_rows - 1 and border_name == "lnB")
                is_header_top = (i == 0 and border_name == "lnT")

                border = etree.SubElement(tcPr, qn(f"a:{border_name}"))
                if is_header_top or is_header_bottom or is_table_bottom:
                    border.set("w", "12700")  # 1pt
                    solidFill = etree.SubElement(border, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", header_color)
                else:
                    border.set("w", "6350")  # 0.5pt
                    solidFill = etree.SubElement(border, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", "E2E8F0")


def _style_cell(cell, *, bg, color, font_name, font_size, bold=False, align="left"):
    """设置单元格样式"""
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(bg)

    # 内边距
    cell.margin_left = Pt(6)
    cell.margin_right = Pt(6)
    cell.margin_top = Pt(4)
    cell.margin_bottom = Pt(4)

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align_map.get(align, PP_ALIGN.LEFT)
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.color.rgb = _rgb(color)
            run.font.bold = bold
