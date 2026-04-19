"""TemplateAssembler — Flow A 增强：单页模板拼装

设计师做好每种布局的单页 .pptx 模板，组合引擎：
  1. 打开各单页模板
  2. 拷贝幻灯片到目标 PPT
  3. 对每页做文本/图表替换

用法:
    assembler = TemplateAssembler()
    assembler.add_slide_from_template("layouts/title.pptx", text_data={"标题": "报告标题"})
    assembler.add_slide_from_template("layouts/chart.pptx", text_data={...}, chart_configs={...})
    assembler.add_slide_from_template("layouts/table.pptx", text_data={...})
    assembler.save("output.pptx")

也支持 layout_book 模式（一个 .pptx 包含多种布局，按页名选取）:
    assembler = TemplateAssembler.from_layout_book("layout_book.pptx")
    assembler.add_slide_by_name("标题页", text_data={...})
    assembler.add_slide_by_name("图表页", text_data={...}, chart_configs={...})
    assembler.save("output.pptx")
"""

import sys
import copy
from pathlib import Path
from typing import Dict, List, Optional, Any

from pptx import Presentation
from pptx.util import Inches
from lxml import etree


class TemplateAssembler:
    """单页模板拼装器"""

    def __init__(self):
        self.prs = None
        self._page_count = 0
        self._layout_book = None  # layout book 模式下的源 prs

    @classmethod
    def from_layout_book(cls, layout_book_path):
        """从 layout book（多页布局模板库）创建

        layout book 是一个 .pptx 文件，每页代表一种布局，
        页面中用占位符标记可替换内容。
        """
        assembler = cls()
        assembler._layout_book = Presentation(str(layout_book_path))
        # 创建空的目标 prs，继承 layout book 的尺寸
        assembler.prs = Presentation()
        assembler.prs.slide_width = assembler._layout_book.slide_width
        assembler.prs.slide_height = assembler._layout_book.slide_height
        return assembler

    def add_slide_from_template(self, template_path, *,
                                 text_data: Dict[str, str] = None,
                                 chart_configs: Dict = None):
        """从单页 .pptx 模板添加一页

        Args:
            template_path: 单页 .pptx 模板路径
            text_data: 文本占位符替换 {"{变量名}": "值"}
            chart_configs: 图表替换配置（与 TemplateReplacer 格式兼容）
        """
        src_prs = Presentation(str(template_path))

        if self.prs is None:
            self.prs = Presentation()
            self.prs.slide_width = src_prs.slide_width
            self.prs.slide_height = src_prs.slide_height

        # 拷贝第一张幻灯片
        src_slide = src_prs.slides[0]
        new_slide = self._copy_slide(src_prs, src_slide)

        # 替换内容
        self._apply_replacements(new_slide, text_data, chart_configs)

        self._page_count += 1
        name = Path(template_path).stem
        print(f"  ✅ 第 {self._page_count} 页: {name}", file=sys.stderr)
        return self

    def add_slide_by_name(self, slide_name: str, *,
                           text_data: Dict[str, str] = None,
                           chart_configs: Dict = None):
        """从 layout book 中按幻灯片名称/索引添加一页

        slide_name 匹配规则:
        1. 精确匹配幻灯片中的标题文本
        2. 尝试作为索引（"0", "1", ...）
        """
        if self._layout_book is None:
            raise ValueError("需要先用 from_layout_book() 初始化")

        src_slide = self._find_slide(slide_name)
        if src_slide is None:
            available = self._list_slide_names()
            raise ValueError(f"找不到幻灯片 '{slide_name}'。可用: {available}")

        new_slide = self._copy_slide(self._layout_book, src_slide)
        self._apply_replacements(new_slide, text_data, chart_configs)

        self._page_count += 1
        print(f"  ✅ 第 {self._page_count} 页: {slide_name}", file=sys.stderr)
        return self

    def add_blank_slide(self):
        """添加空白页"""
        if self.prs is None:
            self.prs = Presentation()
        blank_layout = self.prs.slide_layouts[-1]
        self.prs.slides.add_slide(blank_layout)
        self._page_count += 1
        return self

    def save(self, path):
        """保存"""
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        print(f"\n✅ PPT 已保存: {path} ({self._page_count} 页)", file=sys.stderr)
        return str(path)

    # ========================================================================
    # 内部方法
    # ========================================================================

    def _copy_slide(self, src_prs, src_slide):
        """将源幻灯片拷贝到目标 prs（通过 XML 深拷贝）"""
        # 获取空白布局
        blank_layout = self.prs.slide_layouts[-1]
        new_slide = self.prs.slides.add_slide(blank_layout)

        # 清除新幻灯片默认内容
        sp_tree = new_slide._element.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}spTree',
        )

        # 替换整个 cSld 内容
        src_cSld = src_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
        )
        dst_cSld = new_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
        )

        if src_cSld is not None and dst_cSld is not None:
            # 深拷贝 spTree（形状树）
            src_spTree = src_cSld.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}spTree'
            )
            dst_spTree = dst_cSld.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}spTree'
            )

            # 如果找不到带 pml 命名空间的 spTree，尝试通用方式
            if src_spTree is None:
                for child in src_cSld:
                    if 'spTree' in child.tag:
                        src_spTree = child
                        break
            if dst_spTree is None:
                for child in dst_cSld:
                    if 'spTree' in child.tag:
                        dst_spTree = child
                        break

            if src_spTree is not None and dst_spTree is not None:
                # 移除目标的所有 sp 子元素，保留 nvGrpSpPr 和 grpSpPr
                to_remove = []
                for child in dst_spTree:
                    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if tag not in ('nvGrpSpPr', 'grpSpPr'):
                        to_remove.append(child)
                for child in to_remove:
                    dst_spTree.remove(child)

                # 拷贝源的形状
                for child in src_spTree:
                    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if tag not in ('nvGrpSpPr', 'grpSpPr'):
                        dst_spTree.append(copy.deepcopy(child))

        # 拷贝背景
        src_bg = src_cSld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if src_bg is None:
            for child in src_cSld:
                if 'bg' in child.tag and 'spTree' not in child.tag:
                    src_bg = child
                    break

        if src_bg is not None:
            dst_bg = dst_cSld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
            if dst_bg is None:
                for child in dst_cSld:
                    if 'bg' in child.tag and 'spTree' not in child.tag:
                        dst_bg = child
                        break

            if dst_bg is not None:
                dst_cSld.remove(dst_bg)
            # 插入到 spTree 之前
            dst_cSld.insert(0, copy.deepcopy(src_bg))

        # 拷贝关联的图片等媒体资源
        self._copy_slide_rels(src_prs, src_slide, new_slide)

        return new_slide

    def _copy_slide_rels(self, src_prs, src_slide, dst_slide):
        """拷贝幻灯片的关系文件（图片、嵌入对象等）"""
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        for rel in src_slide.part.rels.values():
            # 跳过布局关系
            if rel.reltype == RT.SLIDE_LAYOUT:
                continue

            try:
                if rel.is_external:
                    dst_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                else:
                    # 拷贝嵌入的部件（如图片）
                    target_part = rel.target_part
                    dst_slide.part.rels.get_or_add(rel.reltype, target_part)
            except Exception:
                pass  # 跳过无法拷贝的关系

    def _apply_replacements(self, slide, text_data, chart_configs):
        """对幻灯片应用文本和图表替换"""
        if not text_data and not chart_configs:
            return

        # 文本替换
        if text_data:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._replace_text_in_shape(shape, text_data)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                self._replace_text_in_paragraph(paragraph, text_data)

        # 图表替换
        if chart_configs:
            self._replace_charts(slide, chart_configs)

    def _replace_text_in_shape(self, shape, data):
        """替换形状中的文本占位符"""
        for paragraph in shape.text_frame.paragraphs:
            self._replace_text_in_paragraph(paragraph, data)

    def _replace_text_in_paragraph(self, paragraph, data):
        """替换段落中的 {占位符}"""
        full_text = paragraph.text
        if '{' not in full_text:
            return

        for key, value in data.items():
            placeholder = f"{{{key}}}" if not key.startswith('{') else key
            if placeholder in full_text:
                full_text = full_text.replace(placeholder, str(value))

        # 重写 run（保留第一个 run 的格式）
        if paragraph.runs:
            first_run = paragraph.runs[0]
            # 清除后续 run
            for run in paragraph.runs[1:]:
                run.text = ""
            first_run.text = full_text

    def _replace_charts(self, slide, chart_configs):
        """替换图表占位符（复用 TemplateReplacer 逻辑）"""
        from pptfi.chart_builder.api import create_combo_chart

        shapes_to_process = list(slide.shapes)
        for shape in shapes_to_process:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            # 匹配 {@ 图:xxx} 或 {@图|xxx}
            for chart_name, config in chart_configs.items():
                markers = [f"{{@图:{chart_name}}}", f"{{@图|{chart_name}}}"]
                if text in markers:
                    # 用图表替换此形状
                    left, top = shape.left, shape.top
                    width, height = shape.width, shape.height

                    # 删除占位符形状
                    sp = shape._element
                    sp.getparent().remove(sp)

                    # 创建图表
                    create_combo_chart(
                        slide=slide,
                        df=config["df"],
                        categories_col=config["categories_col"],
                        series_config=config["series_config"],
                        position=(left, top),
                        size=(width, height),
                        style_config=config.get("style_config"),
                        layout_config=config.get("layout_config"),
                    )
                    break

    def _find_slide(self, name: str):
        """在 layout book 中按名称或索引查找幻灯片"""
        # 尝试按索引
        try:
            idx = int(name)
            if 0 <= idx < len(self._layout_book.slides):
                return self._layout_book.slides[idx]
        except ValueError:
            pass

        # 按标题文本匹配
        for slide in self._layout_book.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if name in shape.text_frame.text:
                        return slide
        return None

    def _list_slide_names(self):
        """列出 layout book 中所有幻灯片的标题"""
        names = []
        for i, slide in enumerate(self._layout_book.slides):
            title = f"[{i}]"
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = f"[{i}] {shape.text_frame.text.strip()[:30]}"
                    break
            names.append(title)
        return names
