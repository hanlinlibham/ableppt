"""
PPT 渲染器 - 从 JSON 重建 PPTX 文件

从解析器生成的 JSON 重建完整的 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from typing import Dict, List
import json
import pandas as pd

from ..chart_builder import create_combo_chart
from ..chart_builder.layout import ChartLayoutConfig, LegendConfig, ValueAxisConfig
from ..chart_builder.styles import StyleConfig
from ..chart_builder.date_axis import DateAxisConfig
from pptx.enum.chart import XL_LEGEND_POSITION


class PPTRenderer:
    """PPT 渲染器"""
    
    def __init__(self, json_data: Dict):
        """
        初始化渲染器
        
        Args:
            json_data: 解析器生成的 JSON 数据
        """
        self.json_data = json_data
        self.prs = Presentation()
        
        # ⭐ 设置幻灯片尺寸为 16:9（从元数据中读取）
        slide_width = json_data.get("metadata", {}).get("slide_width")
        slide_height = json_data.get("metadata", {}).get("slide_height")
        
        if slide_width and slide_height:
            from pptx.util import Emu
            self.prs.slide_width = Emu(slide_width)
            self.prs.slide_height = Emu(slide_height)
            print(f"✅ 幻灯片尺寸: {Emu(slide_width).inches:.2f}\" × {Emu(slide_height).inches:.2f}\"")
        else:
            # 默认使用 16:9
            from pptx.util import Inches
            self.prs.slide_width = Inches(13.33)  # 16:9 标准宽度
            self.prs.slide_height = Inches(7.5)
            print(f"✅ 使用默认 16:9 尺寸")
    
    def render(self) -> Presentation:
        """
        从 JSON 渲染 PPT
        
        Returns:
            Presentation 对象
        """
        for slide_data in self.json_data["slides"]:
            self._render_slide(slide_data)
        
        return self.prs
    
    def _render_slide(self, slide_data: Dict):
        """渲染单个幻灯片"""
        # 创建空白幻灯片
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        for element in slide_data["elements"]:
            self._render_element(slide, element)
    
    def _render_element(self, slide, element: Dict):
        """渲染单个元素"""
        element_type = element.get("type")
        position = element.get("position", {})
        data = element.get("data", {})
        
        if element_type == "chart":
            self._render_chart(slide, data, position)
        elif element_type == "table":
            self._render_table(slide, data, position)
        elif element_type == "textbox":
            self._render_textbox(slide, data, position)
        elif element_type == "group":
            self._render_group(slide, data, position)
    
    def _render_chart(self, slide, chart_data: Dict, position: Dict):
        """渲染图表"""
        try:
            # 重建 DataFrame
            df = pd.DataFrame(chart_data["data"])
            
            # 转换日期列
            categories_col = chart_data["categories_col"]
            if categories_col in df.columns:
                df[categories_col] = pd.to_datetime(df[categories_col])
            
            series_config = chart_data["series_config"]
            layout_info = chart_data.get("layout_info", {})
            
            # 提取数据范围信息
            data_cols = [col for col in df.columns if col != categories_col]
            
            # 配置样式
            style_config = StyleConfig(
                color_scheme="aim00",
                line_width_pt=2.0,
                marker_style="none",
            )
            
            # 配置布局
            legend_pos_map = {
                "-4160": XL_LEGEND_POSITION.TOP,
                "-4107": XL_LEGEND_POSITION.BOTTOM,
                "-4131": XL_LEGEND_POSITION.RIGHT,
                "-4129": XL_LEGEND_POSITION.LEFT,
            }
            legend_pos_str = str(layout_info.get("legend", {}).get("position", -4160))
            legend_position = legend_pos_map.get(legend_pos_str, XL_LEGEND_POSITION.TOP)
            
            layout_config = ChartLayoutConfig(
                title=chart_data.get("title", ""),
                legend_config=LegendConfig(
                    position=legend_position,
                    font_size_pt=9,
                    font_name="黑体",
                ),
                value_axis_config=ValueAxisConfig(
                    number_format="0%",
                    font_size_pt=9,
                    font_name="黑体",
                    has_major_gridlines=False,
                ),
                secondary_value_axis_config=ValueAxisConfig(
                    number_format="#,##0",
                    font_size_pt=9,
                    font_name="黑体",
                    has_major_gridlines=False,
                ),
            )
            
            # 日期轴配置
            label_interval = len(df) // 7
            date_axis_config = DateAxisConfig(
                base_unit='days',
                major_unit=label_interval,
                number_format='yyyy/mm',
            )
            layout_config.date_axis_config = date_axis_config
            
            # 创建图表
            chart = create_combo_chart(
                slide=slide,
                df=df,
                categories_col=categories_col,
                series_config=series_config,
                position=(Emu(position["left"]), Emu(position["top"])),
                size=(Emu(position["width"]), Emu(position["height"])),
                style_config=style_config,
                layout_config=layout_config,
            )
            
            print(f"  ✅ 图表已渲染: {chart_data.get('title', '未命名')}")
            
        except Exception as e:
            print(f"  ❌ 图表渲染失败: {e}")
            import traceback
            traceback.print_exc()
    
    def _render_table(self, slide, table_data: Dict, position: Dict):
        """渲染表格"""
        try:
            rows = table_data["rows"]
            cols = table_data["columns"]
            data = table_data["data"]
            
            # 创建表格
            table_shape = slide.shapes.add_table(
                rows, cols,
                Emu(position["left"]),
                Emu(position["top"]),
                Emu(position["width"]),
                Emu(position["height"])
            )
            
            table = table_shape.table
            
            # 填充数据
            for i, row_data in enumerate(data):
                for j, cell_text in enumerate(row_data):
                    cell = table.cell(i, j)
                    cell.text = cell_text
                    # 设置字体
                    cell.text_frame.paragraphs[0].font.name = "黑体"
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
            
            print(f"  ✅ 表格已渲染: {rows}行 × {cols}列")
            
        except Exception as e:
            print(f"  ❌ 表格渲染失败: {e}")
    
    def _render_textbox(self, slide, textbox_data: Dict, position: Dict):
        """渲染文本框"""
        try:
            # 创建文本框
            textbox = slide.shapes.add_textbox(
                Emu(position["left"]),
                Emu(position["top"]),
                Emu(position["width"]),
                Emu(position["height"])
            )
            
            text_frame = textbox.text_frame
            
            # 添加段落
            paragraphs = textbox_data.get("paragraphs", [])
            for i, para_data in enumerate(paragraphs):
                if i == 0:
                    # 使用默认的第一个段落
                    para = text_frame.paragraphs[0]
                else:
                    para = text_frame.add_paragraph()
                
                # 设置段落级别
                if para_data.get("level"):
                    para.level = para_data["level"]
                
                # 设置对齐方式
                alignment_str = para_data.get("alignment")
                if alignment_str:
                    try:
                        # 尝试解析对齐方式
                        if "CENTER" in alignment_str:
                            para.alignment = PP_ALIGN.CENTER
                        elif "RIGHT" in alignment_str:
                            para.alignment = PP_ALIGN.RIGHT
                        elif "LEFT" in alignment_str:
                            para.alignment = PP_ALIGN.LEFT
                    except:
                        pass
                
                # 添加 runs（保留原始格式）
                runs = para_data.get("runs", [])
                if runs:
                    for run_data in runs:
                        run = para.add_run()
                        run.text = run_data["text"]
                        
                        # ⭐ 应用原始字体信息
                        if run_data.get("font_name"):
                            # 如果原始字体不可用，回退到黑体
                            try:
                                run.font.name = run_data["font_name"]
                            except:
                                run.font.name = "黑体"
                        else:
                            run.font.name = "黑体"
                        
                        if run_data.get("font_size"):
                            run.font.size = Pt(run_data["font_size"])
                        
                        if run_data.get("bold"):
                            run.font.bold = True
                        
                        if run_data.get("italic"):
                            run.font.italic = True
                        
                        if run_data.get("underline"):
                            run.font.underline = True
                        
                        # 应用颜色
                        if run_data.get("color"):
                            try:
                                from pptx.util import RGBColor
                                color_hex = run_data["color"].lstrip('#')
                                r = int(color_hex[0:2], 16)
                                g = int(color_hex[2:4], 16)
                                b = int(color_hex[4:6], 16)
                                run.font.color.rgb = RGBColor(r, g, b)
                            except:
                                pass
                else:
                    # 如果没有 runs，直接设置段落文本
                    para.text = para_data["text"]
                    para.font.name = "黑体"
                    para.font.size = Pt(9)
            
            preview = textbox_data["full_text"][:30] + "..." if len(textbox_data["full_text"]) > 30 else textbox_data["full_text"]
            print(f'  ✅ 文本框已渲染: "{preview}"')
            
        except Exception as e:
            print(f"  ❌ 文本框渲染失败: {e}")
            import traceback
            traceback.print_exc()
    
    def _render_group(self, slide, group_data: Dict, position: Dict):
        """渲染组合形状（递归处理内部元素）"""
        # 注意：python-pptx 不直接支持创建组合形状
        # 我们将组合形状的元素展开，直接渲染到幻灯片上
        for element in group_data.get("elements", []):
            self._render_element(slide, element)
        
        print(f"  ✅ 组合形状已渲染: {group_data.get('element_count', 0)} 个子元素")
    
    def save(self, output_path: Path):
        """
        保存 PPT 文件
        
        Args:
            output_path: 输出路径
        """
        self.prs.save(str(output_path))
        print(f"\n✅ PPT 已保存: {output_path}")


def render_from_json_file(json_path: Path, output_pptx_path: Path):
    """
    从 JSON 文件渲染 PPT
    
    Args:
        json_path: JSON 文件路径
        output_pptx_path: 输出的 PPTX 文件路径
    """
    with open(json_path, 'r', encoding='utf-8') as f:
        json_data = json.load(f)
    
    renderer = PPTRenderer(json_data)
    renderer.render()
    renderer.save(output_pptx_path)


if __name__ == "__main__":
    # 测试
    from pathlib import Path
    
    json_path = Path(__file__).parent.parent.parent / "output" / "aim01_parsed.json"
    output_path = Path(__file__).parent.parent.parent / "output" / "aim01_rebuilt.pptx"
    
    render_from_json_file(json_path, output_path)

