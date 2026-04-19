"""
文本渲染器
"""

from typing import Dict, Any, List
from jinja2 import Template
from pptx.slide import Slide
from pptfi.models.job import TextSpec
from pptfi.utils.template_utils import find_shape_by_name


class TextRenderer:
    """文本渲染器"""

    @staticmethod
    def render(slide: Slide, texts: List[TextSpec], context: Dict[str, Any]) -> None:
        """
        渲染文本到幻灯片

        Args:
            slide: 幻灯片对象
            texts: 文本配置列表
            context: 模板变量上下文
        """
        for text_spec in texts or []:
            shape = find_shape_by_name(slide, text_spec.target)

            if not shape:
                print(f"警告: 未找到文本框 '{text_spec.target}'")
                continue

            if not shape.has_text_frame:
                print(f"警告: '{text_spec.target}' 不是文本框")
                continue

            # 使用 Jinja2 渲染文本
            value = Template(text_spec.value).render(**context)

            # 清除现有文本并添加新文本
            text_frame = shape.text_frame
            text_frame.clear()

            # 添加段落和文本
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = value
