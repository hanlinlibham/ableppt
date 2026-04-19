"""
PPT 模板替换器 - 支持文本、图表、表格的智能替换

占位符格式：
- {变量名} - 简单文本替换
- {@图:图表配置键} - 替换为图表
- {@表:表格配置键} - 替换为表格数据
"""

import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.util import Pt, Inches, Emu
from pathlib import Path
from typing import Dict, Any, Callable

from ..chart_builder import create_combo_chart
from ..chart_builder.layout import ChartLayoutConfig, LegendConfig, ValueAxisConfig
from ..chart_builder.styles import StyleConfig
from ..chart_builder.date_axis import DateAxisConfig
from pptx.enum.chart import XL_LEGEND_POSITION


class TemplateReplacer:
    """PPT 模板替换器"""
    
    def __init__(self, template_path: Path, data_provider: Callable = None):
        """
        初始化模板替换器
        
        Args:
            template_path: 模板文件路径
            data_provider: 数据提供函数（可选）
                          签名：data_provider(key: str) -> Any
        """
        self.template_path = template_path
        self.prs = Presentation(str(template_path))
        self.data_provider = data_provider or {}
        
        # 占位符模式
        self.text_pattern = re.compile(r'\{([^@}][^}]*)\}')  # {变量名}
        self.chart_pattern = re.compile(r'\{@图[:|]([^}]+)\}')  # {@图:配置键} 或 {@图|配置键}
        self.table_pattern = re.compile(r'\{@表[:|]([^}]+)\}')  # {@表:配置键} 或 {@表|配置键}
    
    def replace(self, data: Dict[str, Any] = None, chart_configs: Dict = None) -> Presentation:
        """
        执行替换
        
        Args:
            data: 文本替换数据字典 {变量名: 值}
            chart_configs: 图表配置字典 {配置键: 图表配置}
        
        Returns:
            替换后的 Presentation 对象
        """
        data = data or {}
        chart_configs = chart_configs or {}
        
        print(f"\n🔄 开始替换模板...")
        print(f"  - 文本数据: {len(data)} 个变量")
        print(f"  - 图表配置: {len(chart_configs)} 个")
        
        for slide_idx, slide in enumerate(self.prs.slides):
            print(f"\n处理幻灯片 {slide_idx + 1}...")
            self._process_slide(slide, data, chart_configs)
        
        print(f"\n✅ 模板替换完成！")
        return self.prs
    
    def _process_slide(self, slide, data: Dict, chart_configs: Dict):
        """处理单个幻灯片"""
        shapes_to_process = list(slide.shapes)  # 复制列表，避免迭代时修改
        
        for shape in shapes_to_process:
            self._process_shape(slide, shape, data, chart_configs)
    
    def _process_shape(self, slide, shape, data: Dict, chart_configs: Dict):
        """处理单个形状"""
        
        # 处理组合形状
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                self._process_shape(slide, subshape, data, chart_configs)
            return
        
        # 处理表格
        if hasattr(shape, 'has_table') and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    self._replace_text_in_textframe(cell.text_frame, data)
            return
        
        # 处理文本框
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            text = shape.text_frame.text
            
            # 检查是否包含图表占位符
            chart_match = self.chart_pattern.search(text)
            if chart_match:
                config_key = chart_match.group(1)
                if config_key in chart_configs:
                    self._replace_with_chart(slide, shape, chart_configs[config_key])
                    print(f"  ✅ 图表已替换: {config_key}")
                    return
            
            # 检查是否包含表格占位符
            table_match = self.table_pattern.search(text)
            if table_match:
                config_key = table_match.group(1)
                # TODO: 实现表格替换
                print(f"  ⚠️  表格替换暂未实现: {config_key}")
                return
            
            # 普通文本替换
            self._replace_text_in_textframe(shape.text_frame, data)
    
    def _replace_text_in_textframe(self, text_frame, data: Dict):
        """
        替换文本框中的占位符（保留字体格式）
        
        完全采用 v4_generator.py 的逐字符处理方法
        """
        for paragraph in text_frame.paragraphs:
            # ⭐ 收集所有字符及其字体对象（保留原始 Font 对象引用）
            char_list = []
            for run in paragraph.runs:
                for char in run.text:
                    char_list.append({
                        'char': char,
                        'font': run.font  # 直接保存 Font 对象引用
                    })
            
            if not char_list:
                continue
            
            # 拼接完整文本
            full_text = ''.join([item['char'] for item in char_list])
            
            # 查找占位符
            matches = list(self.text_pattern.finditer(full_text))
            if not matches:
                continue
            
            # 构建新的字符列表
            new_char_list = []
            idx = 0
            
            for match in matches:
                start, end = match.span()
                placeholder_key = match.group(1)
                
                # 添加占位符前的文本
                while idx < start:
                    new_char_list.append(char_list[idx])
                    idx += 1
                
                # 获取替换值
                replacement_value = self._get_value(placeholder_key, data)
                replacement_text = str(replacement_value) if replacement_value is not None else ''
                
                # ⭐ 获取占位符的字体（使用占位符开始位置前一个字符的字体）
                if idx > 0:
                    placeholder_font = char_list[idx - 1]['font']
                elif len(paragraph.runs) > 0:
                    placeholder_font = paragraph.runs[0].font
                else:
                    placeholder_font = None
                
                # ⭐ 检查是否有百分号跟随（v4_generator.py 的百分比处理）
                percent_style = None
                if end < len(char_list) and char_list[end]['char'] == '%':
                    percent_style = char_list[end]['font']
                    # 不跳过 %，让它保留在原位
                
                # 添加替换文本（使用占位符的字体）
                for char in replacement_text:
                    new_char_list.append({
                        'char': char,
                        'font': percent_style or placeholder_font
                    })
                
                idx = end
            
            # 添加剩余文本
            while idx < len(char_list):
                new_char_list.append(char_list[idx])
                idx += 1
            
            # 重建段落（使用 v4_generator.py 的方法）
            self._rebuild_paragraph_v4(paragraph, new_char_list)
    
    def _rebuild_paragraph_v4(self, paragraph, new_char_list):
        """
        重建段落 - 完全采用 v4_generator.py 的方法
        
        关键：使用 copy_font 和 fonts_are_equal 来保留所有字体属性
        """
        if not new_char_list:
            return
        
        # 清空所有 runs
        for run in paragraph.runs:
            run.text = ''
        
        # 创建第一个 run
        current_run = paragraph.add_run()
        current_run.text = new_char_list[0]['char']
        self._copy_font(current_run.font, new_char_list[0]['font'])
        
        # 遍历剩余字符
        for char_info in new_char_list[1:]:
            if not self._fonts_are_equal_v4(char_info['font'], current_run.font):
                # 字体不同，创建新 run
                current_run = paragraph.add_run()
                current_run.text = char_info['char']
                self._copy_font(current_run.font, char_info['font'])
            else:
                # 字体相同，追加到当前 run
                current_run.text += char_info['char']
    
    def _copy_font(self, target_font, source_font):
        """
        复制字体样式 - 完全采用 v4_generator.py 的实现
        
        包括所有属性和颜色类型的处理
        """
        # 复制字体的基础属性
        attributes = ['name', 'size', 'bold', 'italic', 'underline', 'strike', 'subscript', 'superscript']
        for attr in attributes:
            if hasattr(source_font, attr):
                try:
                    setattr(target_font, attr, getattr(source_font, attr))
                except:
                    pass
        
        # ⭐ 复制颜色属性（支持所有颜色类型）
        try:
            if source_font.color.type == MSO_COLOR_TYPE.RGB:
                target_font.color.rgb = source_font.color.rgb
            elif source_font.color.type == MSO_COLOR_TYPE.SCHEME:
                target_font.color.theme_color = source_font.color.theme_color
                if hasattr(source_font.color, 'brightness'):
                    target_font.color.brightness = source_font.color.brightness
            elif source_font.color.type == MSO_COLOR_TYPE.PRESET:
                target_font.color.theme_color = source_font.color.theme_color
            elif source_font.color.type == MSO_COLOR_TYPE.SYSTEM:
                target_font.color.rgb = source_font.color.rgb
        except:
            pass
    
    def _fonts_are_equal_v4(self, font1, font2) -> bool:
        """
        比较两个字体是否相同 - 完全采用 v4_generator.py 的实现
        
        扩展更多的属性和颜色类型
        """
        # 比较字体的基础属性
        attributes = ['name', 'size', 'bold', 'italic', 'underline', 'strike', 'subscript', 'superscript']
        for attr in attributes:
            if getattr(font1, attr, None) != getattr(font2, attr, None):
                return False
        
        # 比较颜色类型和颜色值
        try:
            if font1.color.type != font2.color.type:
                return False
            
            if font1.color.type == MSO_COLOR_TYPE.RGB:
                return font1.color.rgb == font2.color.rgb
            elif font1.color.type == MSO_COLOR_TYPE.SCHEME:
                return (font1.color.theme_color == font2.color.theme_color and
                       getattr(font1.color, 'brightness', None) == getattr(font2.color, 'brightness', None))
            elif font1.color.type == MSO_COLOR_TYPE.PRESET:
                return font1.color.theme_color == font2.color.theme_color
            elif font1.color.type == MSO_COLOR_TYPE.SYSTEM:
                return font1.color.rgb == font2.color.rgb
        except:
            pass
        
        return True
    
    def _replace_with_chart(self, slide, shape, chart_config: Dict):
        """
        用图表替换占位符形状
        
        ⭐ 确保图表的边界和原始矩形完全一致
        
        Args:
            slide: 幻灯片
            shape: 占位符形状（矩形或文本框）
            chart_config: 图表配置，包含：
                - df: DataFrame
                - categories_col: 分类列名
                - series_config: 系列配置
                - style_config: 样式配置（可选）
                - layout_config: 布局配置（可选）
        """
        # ⭐ 精确获取占位符的位置和大小（EMU 单位）
        position = (shape.left, shape.top)
        size = (shape.width, shape.height)
        
        print(f"    → 占位符位置: left={shape.left.inches:.2f}\", top={shape.top.inches:.2f}\"")
        print(f"    → 占位符大小: width={shape.width.inches:.2f}\", height={shape.height.inches:.2f}\"")
        
        # 删除占位符形状
        sp = shape.element
        sp.getparent().remove(sp)
        
        # 获取图表配置
        df = chart_config.get('df')
        categories_col = chart_config.get('categories_col')
        series_config = chart_config.get('series_config', [])
        style_config = chart_config.get('style_config')
        layout_config = chart_config.get('layout_config')
        
        if df is None or categories_col is None:
            print(f"  ❌ 图表配置缺少必要参数: df 或 categories_col")
            return
        
        # ⭐ 创建图表（使用完全相同的位置和大小）
        chart = create_combo_chart(
            slide=slide,
            df=df,
            categories_col=categories_col,
            series_config=series_config,
            position=position,  # 精确的 EMU 位置
            size=size,          # 精确的 EMU 大小
            style_config=style_config,
            layout_config=layout_config,
        )
        
        print(f"    ✅ 图表已创建，边界与原矩形完全一致")
    
    
    def _get_value(self, key: str, data: Dict) -> Any:
        """
        获取替换值
        
        优先级：
        1. data 字典
        2. data_provider 函数
        """
        # 从 data 字典获取
        if key in data:
            return data[key]
        
        # 从 data_provider 获取
        if callable(self.data_provider):
            try:
                return self.data_provider(key)
            except:
                pass
        elif isinstance(self.data_provider, dict):
            return self.data_provider.get(key)
        
        # 未找到，返回原始占位符
        return f'{{{key}}}'
    
    def save(self, output_path: Path):
        """保存替换后的 PPT"""
        self.prs.save(str(output_path))
        print(f"\n✅ 已保存: {output_path}")
    
    def extract_placeholders(self, by_page: bool = False) -> Dict:
        """
        提取所有占位符
        
        Args:
            by_page: 是否按页面分组（默认 False）
        
        Returns:
            如果 by_page=False:
                {
                    'text_placeholders': [...],
                    'chart_placeholders': [...],
                    'table_placeholders': [...]
                }
            
            如果 by_page=True:
                {
                    'pages': [
                        {
                            'page_number': 1,
                            'text_placeholders': [...],
                            'chart_placeholders': [...],
                            'table_placeholders': [...]
                        },
                        ...
                    ],
                    'all_text': [...],  # 所有文本占位符（去重）
                    'all_chart': [...],  # 所有图表占位符（去重）
                }
        """
        if not by_page:
            # 原有逻辑：全局去重
            text_placeholders = set()
            chart_placeholders = set()
            table_placeholders = set()
            
            for slide in self.prs.slides:
                for shape in slide.shapes:
                    placeholders = self._extract_from_shape(shape)
                    text_placeholders.update(placeholders['text'])
                    chart_placeholders.update(placeholders['chart'])
                    table_placeholders.update(placeholders['table'])
            
            return {
                'text_placeholders': sorted(list(text_placeholders)),
                'chart_placeholders': sorted(list(chart_placeholders)),
                'table_placeholders': sorted(list(table_placeholders)),
            }
        else:
            # ⭐ 新逻辑：按页面分组
            pages = []
            all_text = set()
            all_chart = set()
            all_table = set()
            
            for slide_idx, slide in enumerate(self.prs.slides):
                page_data = {
                    'page_number': slide_idx + 1,
                    'text_placeholders': [],
                    'chart_placeholders': [],
                    'table_placeholders': [],
                }
                
                text_set = set()
                chart_set = set()
                table_set = set()
                
                for shape in slide.shapes:
                    placeholders = self._extract_from_shape(shape)
                    text_set.update(placeholders['text'])
                    chart_set.update(placeholders['chart'])
                    table_set.update(placeholders['table'])
                
                page_data['text_placeholders'] = sorted(list(text_set))
                page_data['chart_placeholders'] = sorted(list(chart_set))
                page_data['table_placeholders'] = sorted(list(table_set))
                
                pages.append(page_data)
                
                # 累积全局占位符
                all_text.update(text_set)
                all_chart.update(chart_set)
                all_table.update(table_set)
            
            return {
                'pages': pages,
                'all_text': sorted(list(all_text)),
                'all_chart': sorted(list(all_chart)),
                'all_table': sorted(list(all_table)),
            }
    
    def _extract_from_shape(self, shape) -> Dict:
        """从形状中提取占位符"""
        result = {'text': set(), 'chart': set(), 'table': set()}
        
        # 处理组合形状
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                sub_result = self._extract_from_shape(subshape)
                result['text'].update(sub_result['text'])
                result['chart'].update(sub_result['chart'])
                result['table'].update(sub_result['table'])
            return result
        
        # 处理文本框和表格
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            text = shape.text_frame.text
            
            # 文本占位符
            for match in self.text_pattern.finditer(text):
                result['text'].add(match.group(1))
            
            # 图表占位符
            for match in self.chart_pattern.finditer(text):
                result['chart'].add(match.group(1))
            
            # 表格占位符
            for match in self.table_pattern.finditer(text):
                result['table'].add(match.group(1))
        
        if hasattr(shape, 'has_table') and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text_frame.text
                    for match in self.text_pattern.finditer(text):
                        result['text'].add(match.group(1))
        
        return result


def create_template_from_ppt(ppt_path: Path, output_path: Path, 
                             placeholders: Dict[str, str]):
    """
    从现有 PPT 创建模板（添加占位符）
    
    Args:
        ppt_path: 原始 PPT 路径
        output_path: 输出模板路径
        placeholders: 替换映射 {原始文本: 占位符}
    """
    prs = Presentation(str(ppt_path))
    
    print(f"\n🔧 创建模板: {ppt_path.name} → {output_path.name}")
    
    for slide in prs.slides:
        for shape in slide.shapes:
            _add_placeholders_to_shape(shape, placeholders)
    
    prs.save(str(output_path))
    print(f"✅ 模板已保存: {output_path}")


def _add_placeholders_to_shape(shape, placeholders: Dict[str, str]):
    """向形状添加占位符"""
    
    # 处理组合形状
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            _add_placeholders_to_shape(subshape, placeholders)
        return
    
    # 处理文本框
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for original, placeholder in placeholders.items():
                    run.text = run.text.replace(original, placeholder)
    
    # 处理表格
    if hasattr(shape, 'has_table') and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for original, placeholder in placeholders.items():
                            run.text = run.text.replace(original, placeholder)


if __name__ == "__main__":
    # 测试：提取 aim02.pptx 中的占位符
    from pathlib import Path
    
    aim02_path = Path(__file__).parent.parent.parent / "aim" / "aim02.pptx"
    
    replacer = TemplateReplacer(aim02_path)
    placeholders = replacer.extract_placeholders()
    
    print("\n📋 提取的占位符:")
    print("="*80)
    print(f"\n文本占位符 ({len(placeholders['text_placeholders'])} 个):")
    for p in placeholders['text_placeholders']:
        print(f"  - {{{p}}}")
    
    if placeholders['chart_placeholders']:
        print(f"\n图表占位符 ({len(placeholders['chart_placeholders'])} 个):")
        for p in placeholders['chart_placeholders']:
            print(f"  - {{@图:{p}}}")
    
    if placeholders['table_placeholders']:
        print(f"\n表格占位符 ({len(placeholders['table_placeholders'])} 个):")
        for p in placeholders['table_placeholders']:
            print(f"  - {{@表:{p}}}")
