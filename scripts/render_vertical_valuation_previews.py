#!/usr/bin/env python3
"""Render the 4 page-specific vertical valuation presets and export PowerPoint previews."""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from ableppt import (
    create_range_snapshot_chart,
    get_asx200_sector_valuation_snapshot_preset,
    get_msci_emu_sector_valuation_snapshot_preset,
    get_msci_japan_sector_valuation_snapshot_preset,
    get_sp500_sector_valuation_snapshot_preset,
)


ROOT = Path(__file__).resolve().parent.parent
EXPORT_SCRIPT = ROOT / "scripts" / "export_powerpoint_pdf.py"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
CHART_POS = (Inches(0.75), Inches(1.45))
CHART_SIZE = (Inches(11.85), Inches(5.35))


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _add_page_title(slide, title: str, subtitle: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.30), Inches(12.0), Inches(0.55))
    title_frame = title_box.text_frame
    title_frame.text = title
    run = title_frame.paragraphs[0].runs[0]
    run.font.bold = True
    run.font.name = "Aptos"
    run.font.size = Pt(26)

    subtitle_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.84), Inches(12.0), Inches(0.28))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    sub_run = subtitle_frame.paragraphs[0].runs[0]
    sub_run.font.name = "Aptos"
    sub_run.font.size = Pt(12)


def _render_one(df: pd.DataFrame, preset_fn, output_pptx: Path) -> None:
    preset = preset_fn(df)
    prs = _new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_page_title(slide, preset["title"], preset["subtitle"])
    create_range_snapshot_chart(
        slide=slide,
        df=preset["df"],
        categories_col=preset["categories_col"],
        min_col=preset["min_col"],
        max_col=preset["max_col"],
        average_col=preset["average_col"],
        current_col=preset["current_col"],
        orientation=preset["orientation"],
        position=CHART_POS,
        size=CHART_SIZE,
        range_color=preset["range_color"],
        average_color=preset["average_color"],
        current_color=preset["current_color"],
        number_format=preset["number_format"],
        show_average_ticks=preset["show_average_ticks"],
        show_current_markers=preset["show_current_markers"],
        show_current_labels=preset["show_current_labels"],
        axis_break=preset["axis_break"],
    )
    prs.save(output_pptx)


def _export_preview(pptx_path: Path, out_dir: Path) -> dict:
    pdf_path = out_dir / f"{pptx_path.stem}.pdf"
    png_dir = out_dir / f"{pptx_path.stem}-png"
    result = subprocess.run(
        [sys.executable, str(EXPORT_SCRIPT), str(pptx_path), str(pdf_path), "--png-dir", str(png_dir)],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(result.stdout or result.stderr)
    return json.loads(result.stdout)


def main() -> int:
    output_dir = ROOT / "output" / "vertical-valuation-previews"
    output_dir.mkdir(parents=True, exist_ok=True)

    jobs = [
        ("asx200", ROOT / "data" / "asx200_sector_valuation_snapshot.csv", get_asx200_sector_valuation_snapshot_preset),
        ("sp500", ROOT / "data" / "sp500_sector_valuation_snapshot.csv", get_sp500_sector_valuation_snapshot_preset),
        ("msci_emu", ROOT / "data" / "msci_emu_sector_valuation_snapshot.csv", get_msci_emu_sector_valuation_snapshot_preset),
        ("msci_japan", ROOT / "data" / "msci_japan_sector_valuation_snapshot.csv", get_msci_japan_sector_valuation_snapshot_preset),
    ]

    payload = {"status": "ok", "previews": []}
    for name, csv_path, preset_fn in jobs:
        df = pd.read_csv(csv_path)
        pptx_path = output_dir / f"{name}.pptx"
        _render_one(df, preset_fn, pptx_path)
        preview = _export_preview(pptx_path, output_dir)
        payload["previews"].append(
            {
                "name": name,
                "csv": str(csv_path),
                "pptx": str(pptx_path),
                "pdf": preview["pdf"],
                "pngs": preview["pngs"],
            }
        )

    print(json.dumps(payload, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
