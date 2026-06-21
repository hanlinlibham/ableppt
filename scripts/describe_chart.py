#!/usr/bin/env python3
"""解析已有 PPT 中的图表，输出可读摘要（不输出 XML）

用法: python describe_chart.py <input.pptx> [--slide N]
输出: JSON 摘要到 stdout
"""

import sys
import json
import argparse
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Emu
from ableppt.chart_builder import ChartParser


def describe_chart(chart, chart_idx):
    """提取单个图表的可读摘要"""
    try:
        series_config, df, categories_col, layout_info = ChartParser(chart).parse()
    except Exception as e:
        return {"chart_index": chart_idx, "error": str(e)}

    # 构建系列摘要
    series_summary = []
    for sc in series_config:
        col_key = sc.get("key", sc.get("name", "unknown"))
        data_points = len(df[col_key]) if col_key in df.columns else 0
        series_summary.append({
            "name": sc.get("name", col_key),
            "type": sc.get("type", "unknown"),
            "axis": sc.get("axis", "primary"),
            "data_points": data_points,
        })

    # 日期范围
    date_range = None
    if categories_col and categories_col in df.columns:
        col = df[categories_col]
        try:
            first = col.iloc[0]
            last = col.iloc[-1]
            if hasattr(first, "strftime"):
                date_range = f"{first.strftime('%Y-%m-%d')} ~ {last.strftime('%Y-%m-%d')}"
            else:
                date_range = f"{first} ~ {last}"
        except Exception:
            pass

    title = ""
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text
    except Exception:
        pass

    return {
        "chart_index": chart_idx,
        "title": title,
        "series": series_summary,
        "categories_col": categories_col,
        "date_range": date_range,
        "data_rows": len(df),
    }


def main():
    parser = argparse.ArgumentParser(description="描述 PPT 中的图表")
    parser.add_argument("input", help="输入 .pptx 文件路径")
    parser.add_argument("--slide", type=int, help="只解析第 N 张幻灯片（从 1 开始）")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"错误: 文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    # 重定向 print 到 stderr
    old_stdout = sys.stdout
    sys.stdout = sys.stderr

    prs = Presentation(str(input_path))
    result = {"slides": []}

    for slide_idx, slide in enumerate(prs.slides):
        if args.slide is not None and (slide_idx + 1) != args.slide:
            continue

        slide_info = {"slide_index": slide_idx, "charts": []}
        chart_idx = 0

        for shape in slide.shapes:
            if not hasattr(shape, "has_chart") or not shape.has_chart:
                continue

            summary = describe_chart(shape.chart, chart_idx)
            # 添加位置信息
            summary["position"] = {
                "left_inches": round(shape.left / 914400, 2),
                "top_inches": round(shape.top / 914400, 2),
                "width_inches": round(shape.width / 914400, 2),
                "height_inches": round(shape.height / 914400, 2),
            }
            slide_info["charts"].append(summary)
            chart_idx += 1

        result["slides"].append(slide_info)

    sys.stdout = old_stdout
    json.dump(result, sys.stdout, ensure_ascii=False, indent=2)
    print()


if __name__ == "__main__":
    main()
