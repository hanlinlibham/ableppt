import copy
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

from pptfi import (
    PageComposer,
    PptEngine,
    chart_engine_info,
    demo_waterfall,
    describe_chart,
    generate_ppt,
    parse_bubble,
    parse_range_snapshot,
    parse_scatter,
    parse_template,
    parse_waterfall,
    render_bubble,
    render_family,
    render_range_snapshot,
    render_scatter,
    render_waterfall,
    render_job,
)


PROJECT_ROOT = Path(__file__).resolve().parent.parent


def run_cli(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, "-m", "pptfi", *args],
        cwd=PROJECT_ROOT,
        text=True,
        capture_output=True,
        check=False,
    )


def assert_widescreen_presentation(path: Path) -> None:
    prs = Presentation(str(path))
    width_in = round(prs.slide_width / 914400, 3)
    height_in = round(prs.slide_height / 914400, 3)
    assert width_in == 13.333
    assert height_in == 7.5


def test_sdk_top_level_exports():
    assert callable(parse_template)
    assert callable(generate_ppt)
    assert callable(chart_engine_info)
    assert callable(render_waterfall)
    assert callable(render_scatter)
    assert callable(render_bubble)
    assert callable(render_family)
    assert callable(render_range_snapshot)
    assert PptEngine.__name__ == "PptEngine"
    assert PageComposer.__name__ == "PageComposer"


def test_cli_help():
    result = run_cli("--help")
    assert result.returncode == 0
    assert "Financial PowerPoint automation CLI with chart-engine compatibility helpers" in result.stdout
    assert "parse-template" in result.stdout
    assert "render" in result.stdout
    assert "chart-engine-info" in result.stdout
    assert "render-waterfall" in result.stdout
    assert "parse-waterfall" in result.stdout
    assert "demo-waterfall" in result.stdout
    assert "render-scatter" in result.stdout
    assert "parse-scatter" in result.stdout
    assert "render-bubble" in result.stdout
    assert "parse-bubble" in result.stdout
    assert "render-range-snapshot" in result.stdout
    assert "parse-range-snapshot" in result.stdout
    assert "render-family" in result.stdout


def test_cli_chart_engine_info_reports_compatibility_layer():
    result = run_cli("chart-engine-info")
    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["compatibility_layer"]["enabled"] is True
    assert payload["operations"]["create_combo_chart"]["module"].startswith("ablechart")
    assert payload["operations"]["create_waterfall_chart"]["module"].startswith("ablechart")
    assert payload["operations"]["create_scatter_chart"]["module"].startswith("ablechart")
    assert payload["operations"]["create_bubble_chart"]["module"].startswith("ablechart")
    assert payload["operations"]["create_range_snapshot_chart"]["module"].startswith("ablechart")


def test_cli_parse_template():
    template = PROJECT_ROOT / "aim" / "aim00.pptx"
    result = run_cli("parse-template", str(template))
    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert "pages" in payload
    assert "all_chart" in payload


def test_cli_generate_real_ppt(tmp_path: Path):
    template = PROJECT_ROOT / "aim" / "aim00.pptx"
    config_path = tmp_path / "config.json"
    output_path = tmp_path / "output.pptx"
    config_path.write_text(
        json.dumps({"text_data": {"标题": "CLI 验证"}, "chart_configs": {}}, ensure_ascii=False),
        encoding="utf-8",
    )

    result = run_cli("generate", str(template), str(config_path), str(output_path))
    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["status"] == "ok"
    assert output_path.exists()
    assert output_path.stat().st_size > 1000


def test_sdk_describe_chart_recovers_categories_col_for_composer_output(tmp_path: Path):
    job_path = PROJECT_ROOT / "job_demo_company_a.json"
    raw_job = json.loads(job_path.read_text(encoding="utf-8"))
    job = copy.deepcopy(raw_job)
    job["output"]["path"] = str(tmp_path / "company_a.pptx")
    local_job_path = tmp_path / "job.json"
    local_job_path.write_text(json.dumps(job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(local_job_path)
    report = describe_chart(result["output"], slide=3)
    charts = report["slides"][0]["charts"]

    assert charts
    assert charts[0]["categories_col"] in {"年份", "日期", "分类"}


def test_sdk_and_cli_waterfall_demo(tmp_path: Path):
    output = tmp_path / "waterfall-demo.pptx"
    sdk_result = demo_waterfall(output)
    assert Path(sdk_result["output"]).exists()
    assert_widescreen_presentation(output)

    parsed = parse_waterfall(output)
    assert parsed["categories_col"] == "项目"
    assert parsed["rows"][0]["measure"] == "total"

    cli_output = tmp_path / "waterfall-cli.pptx"
    cli_result = run_cli("demo-waterfall", str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert Path(payload["output"]).exists()

    parsed_cli = run_cli("parse-waterfall", str(cli_output))
    assert parsed_cli.returncode == 0, parsed_cli.stderr
    parsed_payload = json.loads(parsed_cli.stdout)
    assert parsed_payload["categories_col"] == "项目"


def test_sdk_and_cli_render_waterfall_from_spec(tmp_path: Path):
    config = PROJECT_ROOT / "waterfall_demo.json"
    output = tmp_path / "waterfall-spec-sdk.pptx"
    sdk_result = render_waterfall(config, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["summary"]["prepared_rows"] == 6
    assert_widescreen_presentation(output)

    prs = Presentation(str(output))
    assert len(prs.slides) == 1
    chart_shapes = [shape for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert chart_shapes

    cli_output = tmp_path / "waterfall-spec-cli.pptx"
    cli_result = run_cli("render-waterfall", str(config), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert Path(payload["output"]).exists()
    assert payload["summary"]["total_categories"] == ["期初收益", "期末收益"]


def test_sdk_and_cli_render_semantic_family(tmp_path: Path):
    config = {
        "family": "performance_compare",
        "title": "基金 vs 基准",
        "rows": [
            {"日期": "2025-01-31", "基金": 0.01, "沪深300": 0.0},
            {"日期": "2025-02-28", "基金": 0.03, "沪深300": 0.01},
            {"日期": "2025-03-31", "基金": 0.02, "沪深300": 0.015},
        ],
        "categories_col": "日期",
        "series": [
            {"key": "基金", "name": "基金", "role": "fund", "type": "line"},
            {"key": "沪深300", "name": "沪深300", "role": "benchmark", "type": "line"},
        ],
    }
    config_path = tmp_path / "family.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "family-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "performance_compare"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "family-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "performance_compare"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_matrix_family(tmp_path: Path):
    config = {
        "family": "ranked_tile_matrix",
        "title": "股票风格箱",
        "row_col": "市值",
        "column_col": "风格",
        "value_col": "占比",
        "rows": [
            {"市值": "大盘", "风格": "价值", "占比": 51},
            {"市值": "大盘", "风格": "平衡", "占比": 9},
            {"市值": "大盘", "风格": "成长", "占比": 23},
            {"市值": "中盘", "风格": "价值", "占比": 0},
            {"市值": "中盘", "风格": "平衡", "占比": 6},
            {"市值": "中盘", "风格": "成长", "占比": 11},
        ],
    }
    config_path = tmp_path / "matrix.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "matrix-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "ranked_tile_matrix"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "matrix-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "ranked_tile_matrix"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_table_plus_chart_composite(tmp_path: Path):
    config = {
        "family": "table_plus_chart_composite",
        "title": "基金换手率及重仓股",
        "chart_family": "score_overlay",
        "chart_spec": {
            "df": [
                {"报告期": "2024-12-31", "换手率": 0.41, "同类中位数": 0.36, "相对得分": 70},
                {"报告期": "2025-03-31", "换手率": 0.39, "同类中位数": 0.34, "相对得分": 68},
                {"报告期": "2025-06-30", "换手率": 0.45, "同类中位数": 0.33, "相对得分": 74}
            ],
            "categories_col": "报告期",
            "raw_series": [
                {"key": "换手率", "name": "换手率(左轴)", "type": "line", "role": "raw"},
                {"key": "同类中位数", "name": "同类基金中位数", "type": "line", "role": "median"}
            ],
            "score_series": [
                {"key": "相对得分", "name": "换手率_相对得分(右轴)", "type": "line", "role": "score"}
            ],
            "raw_number_format": "0%",
            "score_number_format": "0"
        },
        "headers": ["报告期", "重仓股市值占基金净值比", "占股票市值比"],
        "table_rows": [
            ["2024-12-31", "54.2%", "57.8%"],
            ["2025-03-31", "55.1%", "58.6%"],
            ["2025-06-30", "56.4%", "60.2%"]
        ]
    }
    config_path = tmp_path / "composite.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "composite-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "table_plus_chart_composite"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "composite-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "table_plus_chart_composite"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_factor_attribution_panel(tmp_path: Path):
    config = {
        "family": "factor_attribution_panel",
        "title": "股票组合的收益因子分解",
        "chart_family": "attribution_decomposition",
        "chart_spec": {
            "df": [
                {"阶段": "期初收益", "贡献": 8.5, "度量": "total"},
                {"阶段": "全A因子", "贡献": 1.4, "度量": "relative"},
                {"阶段": "个股选择", "贡献": 1.1, "度量": "relative"},
                {"阶段": "期末收益", "贡献": 11.0, "度量": "total"}
            ],
            "categories_col": "阶段",
            "value_col": "贡献",
            "measure_col": "度量"
        },
        "summary_items": [
            {"label": "最大正贡献", "value": "全A因子 +1.4", "accent": "10B981"},
            {"label": "Alpha 来源", "value": "个股选择 +1.1", "accent": "5679CC"}
        ],
        "bullets": [
            "系统性因子仍是主要贡献来源。",
            "个股选择为正，具备 alpha 贡献。"
        ]
    }
    config_path = tmp_path / "panel.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "panel-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "factor_attribution_panel"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "panel-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "factor_attribution_panel"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_regime_table_panel(tmp_path: Path):
    config = {
        "family": "regime_table_panel",
        "title": "调仓时机",
        "chart_spec": {
            "df": [
                {"日期": "2025-01-31", "基金": 0.010, "沪深300": 0.004},
                {"日期": "2025-02-28", "基金": 0.030, "沪深300": 0.018},
                {"日期": "2025-03-31", "基金": 0.027, "沪深300": 0.015}
            ],
            "categories_col": "日期",
            "series_entries": [
                {"key": "基金", "name": "易方达蓝筹精选混合", "role": "fund", "type": "line"},
                {"key": "沪深300", "name": "沪深300", "role": "benchmark", "type": "line"}
            ],
            "events": [
                {"start": "2025-01-31", "end": "2025-02-28", "label": "快速上涨", "color": "E07138"},
                {"start": "2025-03-31", "end": "2025-03-31", "label": "震荡市场", "color": "E0C470"}
            ],
            "show_event_labels": False
        },
        "table_headers": ["起止日期", "当前基金", "沪深300"],
        "table_rows": [
            ["2018-01-23 ~ 2018-12-18", "-3.96%", "-4.55%"],
            ["2018-12-19 ~ 2019-04-03", "27.23%", "28.57%"]
        ]
    }
    config_path = tmp_path / "regime.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "regime-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "regime_table_panel"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "regime-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "regime_table_panel"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_manager_timeline_profile(tmp_path: Path):
    config = {
        "family": "manager_timeline_profile",
        "title": "基金经理变更",
        "manager_name": "张坤",
        "tenure_start": "2018-09-05",
        "tenure_days": 2783,
        "tenure_return": "76.43%",
        "chart_spec": {
            "df": [
                {"日期": "2024-12-31", "张坤": 0.42, "沪深300": 0.28},
                {"日期": "2025-03-31", "张坤": 0.44, "沪深300": 0.29}
            ],
            "categories_col": "日期",
            "series_entries": [
                {"key": "张坤", "name": "张坤", "role": "fund", "type": "line"},
                {"key": "沪深300", "name": "沪深300", "role": "benchmark", "type": "line"}
            ]
        }
    }
    config_path = tmp_path / "manager.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "manager-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "manager_timeline_profile"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "manager-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "manager_timeline_profile"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_award_timeline_panel(tmp_path: Path):
    config = {
        "family": "award_timeline_panel",
        "title": "获奖情况",
        "headers": ["获奖时间", "获得奖项", "获奖时经理"],
        "rows": []
    }
    config_path = tmp_path / "award.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "award-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "award_timeline_panel"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "award-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "award_timeline_panel"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_selection_timing_grid(tmp_path: Path):
    config = {
        "family": "selection_timing_grid",
        "title": "选股择时",
        "section_titles": ["选股能力", "择时能力"],
        "level_labels": ["不显著", "一般", "强"],
        "row_labels": ["成立以来", "近两年", "近三年", "近五年"],
        "sections": [
            ["一般", "一般", "强", "强"],
            ["不显著", "一般", "一般", "强"]
        ]
    }
    config_path = tmp_path / "selection_grid.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "selection-grid-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "selection_timing_grid"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "selection-grid-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "selection_timing_grid"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_holding_detail_panel(tmp_path: Path):
    config = {
        "family": "holding_detail",
        "title": "股票持仓详情",
        "subtitle": "前十大重仓股票",
        "headers": ["股票代码", "股票名称", "占基金资产净值比例"],
        "rows": [
            ["600519.SH", "贵州茅台", "9.82%"],
            ["00700.HK", "腾讯控股", "8.73%"]
        ]
    }
    config_path = tmp_path / "holding.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "holding-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "holding_detail"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "holding-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "holding_detail"
    assert Path(payload["output"]).exists()


def test_sdk_and_cli_render_dual_chart_panel(tmp_path: Path):
    config = {
        "family": "dual_chart_panel",
        "title": "股票资产收益贡献",
        "left_chart_family": "distribution_plus_history",
        "left_chart_spec": {
            "df": [
                {"报告期": "2024-12-31", "消费": 0.55, "医药": 0.14},
                {"报告期": "2025-03-31", "消费": 0.58, "医药": 0.12}
            ],
            "mode": "history",
            "categories_col": "报告期",
            "series_columns": ["消费", "医药"],
            "chart_type": "area"
        },
        "right_chart_family": "performance_compare",
        "right_chart_spec": {
            "df": [
                {"行业": "消费", "基金": 0.18, "基准": 0.12},
                {"行业": "医药", "基金": 0.09, "基准": 0.05}
            ],
            "categories_col": "行业",
            "series_entries": [
                {"key": "基金", "name": "基金收益", "role": "fund", "type": "bar"},
                {"key": "基准", "name": "基准收益", "role": "benchmark", "type": "bar"}
            ],
            "number_format": "0%"
        }
    }
    config_path = tmp_path / "dual.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "dual-sdk.pptx"
    sdk_result = render_family(config_path, output)
    assert Path(sdk_result["output"]).exists()
    assert sdk_result["family"] == "dual_chart_panel"
    assert_widescreen_presentation(output)

    cli_output = tmp_path / "dual-cli.pptx"
    cli_result = run_cli("render-family", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    payload = json.loads(cli_result.stdout)
    assert payload["family"] == "dual_chart_panel"
    assert Path(payload["output"]).exists()


def test_packaged_research_shell_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_research_shell_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "research-shell.pptx")
    job_path = tmp_path / "research-shell.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()
    prs = Presentation(str(result["output"]))
    width_in = round(prs.slide_width / 914400, 3)
    height_in = round(prs.slide_height / 914400, 3)
    assert width_in == 10.0
    assert height_in == 7.5


def test_packaged_dashboard_shell_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_dashboard_shell_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "dashboard-shell.pptx")
    job_path = tmp_path / "dashboard-shell.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()
    prs = Presentation(str(result["output"]))
    width_in = round(prs.slide_width / 914400, 3)
    height_in = round(prs.slide_height / 914400, 3)
    assert width_in == 13.333
    assert height_in == 7.5


def test_packaged_factsheet_shell_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_factsheet_shell_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "factsheet-shell.pptx")
    job_path = tmp_path / "factsheet-shell.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()
    prs = Presentation(str(result["output"]))
    width_in = round(prs.slide_width / 914400, 3)
    height_in = round(prs.slide_height / 914400, 3)
    assert width_in == 10.0
    assert height_in == 7.5


def test_packaged_summary_shell_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_summary_shell_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "summary-shell.pptx")
    job_path = tmp_path / "summary-shell.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()
    prs = Presentation(str(result["output"]))
    width_in = round(prs.slide_width / 914400, 3)
    height_in = round(prs.slide_height / 914400, 3)
    assert width_in == 13.333
    assert height_in == 7.5


def test_packaged_section_cover_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_section_cover_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "section-cover.pptx")
    job_path = tmp_path / "section-cover.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()
    prs = Presentation(str(result["output"]))
    assert round(prs.slide_width / 914400, 3) == 10.0
    assert round(prs.slide_height / 914400, 3) == 7.5


def test_packaged_chapter_divider_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_chapter_divider_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "chapter-divider.pptx")
    job_path = tmp_path / "chapter-divider.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()
    prs = Presentation(str(result["output"]))
    assert round(prs.slide_width / 914400, 3) == 13.333
    assert round(prs.slide_height / 914400, 3) == 7.5


def test_packaged_profile_factsheet_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_profile_factsheet_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "profile-factsheet.pptx")
    job_path = tmp_path / "profile-factsheet.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()
    prs = Presentation(str(result["output"]))
    assert round(prs.slide_width / 914400, 3) == 10.0
    assert round(prs.slide_height / 914400, 3) == 7.5


def test_composer_waterfall_layout_round_trip(tmp_path: Path):
    csv_path = tmp_path / "waterfall.csv"
    csv_path.write_text(
        "阶段,贡献,度量\n期初收益,8.5,total\n权益贡献,2.1,relative\n债券贡献,1.3,relative\n汇率拖累,-1.8,relative\n期末收益,10.1,total\n",
        encoding="utf-8",
    )
    output_path = tmp_path / "composer-waterfall.pptx"
    job = {
        "mode": "composer",
        "theme": "able_finance",
        "datasources": {
            "wf": {"type": "csv", "path": str(csv_path)},
        },
        "pages": [
            {
                "layout": "waterfall",
                "data": {
                    "title": "组合收益归因桥",
                    "source": "wf",
                    "categories_col": "阶段",
                    "value_col": "贡献",
                    "measure_col": "度量",
                    "insight": "权益和债券贡献推动收益抬升，但汇率因素构成主要拖累。",
                    "show_connectors": True,
                    "show_value_labels": True,
                },
            }
        ],
        "output": {"path": str(output_path)},
    }
    job_path = tmp_path / "job.json"
    job_path.write_text(json.dumps(job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()

    parsed = parse_waterfall(result["output"])
    assert parsed["categories_col"] == "阶段"
    assert parsed["rows"][0]["度量"] == "total"
    assert parsed["rows"][3]["贡献"] == -1.8


def test_packaged_composer_waterfall_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_waterfall_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "packaged-waterfall.pptx")
    job_path = tmp_path / "job.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()

    parsed = parse_waterfall(result["output"])
    assert parsed["categories_col"] == "阶段"
    assert parsed["rows"][-1]["度量"] == "total"


def test_packaged_composer_scatter_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_scatter_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "packaged-scatter.pptx")
    job_path = tmp_path / "job.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()

    parsed = parse_scatter(result["output"])
    assert parsed["chart_family"] == "scatter"
    assert parsed["x_col"] == "波动率"
    assert parsed["y_col"] == "收益率"


def test_packaged_composer_bubble_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_bubble_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "packaged-bubble.pptx")
    job_path = tmp_path / "job.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()

    parsed = parse_bubble(result["output"])
    assert parsed["chart_family"] == "bubble"
    assert parsed["size_col"] == "规模"


def test_packaged_composer_range_snapshot_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_range_snapshot_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "packaged-range-snapshot.pptx")
    if "datasources" in raw_job and "valuation" in raw_job["datasources"]:
        raw_job["datasources"]["valuation"]["path"] = str(PROJECT_ROOT / "data" / "range_snapshot_valuation.csv")
    job_path = tmp_path / "job.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()

    parsed = parse_range_snapshot(result["output"])
    assert parsed["chart_family"] == "range_snapshot"
    assert parsed["min_col"] == "range_min"
    assert parsed["current_col"] == "current"
    assert parsed["orientation"] == "horizontal"
    assert parsed["rows"][0]["market"] == "China"


def test_packaged_composer_range_snapshot_sector_job(tmp_path: Path):
    raw_job = json.loads((PROJECT_ROOT / "job_range_snapshot_sector_demo.json").read_text(encoding="utf-8"))
    raw_job["output"]["path"] = str(tmp_path / "packaged-range-snapshot-sector.pptx")
    if "datasources" in raw_job and "sector_valuation" in raw_job["datasources"]:
        raw_job["datasources"]["sector_valuation"]["path"] = str(PROJECT_ROOT / "data" / "range_snapshot_sector_valuation.csv")
    job_path = tmp_path / "job_sector.json"
    job_path.write_text(json.dumps(raw_job, ensure_ascii=False, indent=2), encoding="utf-8")

    result = render_job(job_path)
    assert Path(result["output"]).exists()

    parsed = parse_range_snapshot(result["output"])
    assert parsed["chart_family"] == "range_snapshot"
    assert parsed["orientation"] == "vertical"
    assert parsed["rows"][0]["sector"] == "Energy"


def test_sdk_and_cli_render_scatter_and_bubble(tmp_path: Path):
    scatter_config = PROJECT_ROOT / "scatter_demo.json"
    bubble_config = PROJECT_ROOT / "bubble_demo.json"

    scatter_output = tmp_path / "scatter-sdk.pptx"
    scatter_result = render_scatter(scatter_config, scatter_output)
    assert Path(scatter_result["output"]).exists()
    assert_widescreen_presentation(scatter_output)
    parsed_scatter = parse_scatter(scatter_output)
    assert parsed_scatter["chart_family"] == "scatter"
    assert parsed_scatter["x_col"] == "波动率"
    assert parsed_scatter["y_col"] == "收益率"

    bubble_output = tmp_path / "bubble-sdk.pptx"
    bubble_result = render_bubble(bubble_config, bubble_output)
    assert Path(bubble_result["output"]).exists()
    assert_widescreen_presentation(bubble_output)
    parsed_bubble = parse_bubble(bubble_output)
    assert parsed_bubble["chart_family"] == "bubble"
    assert parsed_bubble["size_col"] == "规模"

    cli_scatter_output = tmp_path / "scatter-cli.pptx"
    cli_scatter = run_cli("render-scatter", str(scatter_config), str(cli_scatter_output))
    assert cli_scatter.returncode == 0, cli_scatter.stderr
    scatter_payload = json.loads(cli_scatter.stdout)
    assert Path(scatter_payload["output"]).exists()

    cli_parse_scatter = run_cli("parse-scatter", str(cli_scatter_output))
    assert cli_parse_scatter.returncode == 0, cli_parse_scatter.stderr
    parsed_scatter_payload = json.loads(cli_parse_scatter.stdout)
    assert parsed_scatter_payload["chart_family"] == "scatter"

    cli_bubble_output = tmp_path / "bubble-cli.pptx"
    cli_bubble = run_cli("render-bubble", str(bubble_config), str(cli_bubble_output))
    assert cli_bubble.returncode == 0, cli_bubble.stderr
    bubble_payload = json.loads(cli_bubble.stdout)
    assert Path(bubble_payload["output"]).exists()

    cli_parse_bubble = run_cli("parse-bubble", str(cli_bubble_output))
    assert cli_parse_bubble.returncode == 0, cli_parse_bubble.stderr
    parsed_bubble_payload = json.loads(cli_parse_bubble.stdout)
    assert parsed_bubble_payload["chart_family"] == "bubble"


def test_sdk_and_cli_render_range_snapshot(tmp_path: Path):
    csv_path = tmp_path / "range_snapshot.csv"
    csv_path.write_text(
        "market,range_min,range_max,average,current\n"
        "China,7.8,24.5,11.2,13.8\n"
        "EM,7.5,18.1,11.7,14.0\n"
        "Europe,8.2,19.7,13.1,14.6\n"
        "U.S.,10.2,22.7,16.0,22.7\n",
        encoding="utf-8",
    )
    config = {
        "csv_path": str(csv_path),
        "categories_col": "market",
        "min_col": "range_min",
        "max_col": "range_max",
        "average_col": "average",
        "current_col": "current",
        "orientation": "horizontal",
        "title": "Global valuations",
    }
    config_path = tmp_path / "range_snapshot.json"
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")

    sdk_output = tmp_path / "range-snapshot-sdk.pptx"
    sdk_result = render_range_snapshot(config_path, sdk_output)
    assert Path(sdk_result["output"]).exists()
    assert_widescreen_presentation(sdk_output)
    parsed_sdk = parse_range_snapshot(sdk_output)
    assert parsed_sdk["chart_family"] == "range_snapshot"
    assert parsed_sdk["min_col"] == "range_min"
    assert parsed_sdk["max_col"] == "range_max"
    assert parsed_sdk["average_col"] == "average"
    assert parsed_sdk["current_col"] == "current"
    assert parsed_sdk["orientation"] == "horizontal"
    assert parsed_sdk["rows"][0]["market"] == "China"

    cli_output = tmp_path / "range-snapshot-cli.pptx"
    cli_result = run_cli("render-range-snapshot", str(config_path), str(cli_output))
    assert cli_result.returncode == 0, cli_result.stderr
    cli_payload = json.loads(cli_result.stdout)
    assert Path(cli_payload["output"]).exists()

    parsed_cli = run_cli("parse-range-snapshot", str(cli_output))
    assert parsed_cli.returncode == 0, parsed_cli.stderr
    parsed_cli_payload = json.loads(parsed_cli.stdout)
    assert parsed_cli_payload["chart_family"] == "range_snapshot"
    assert parsed_cli_payload["orientation"] == "horizontal"
    assert parsed_cli_payload["rows"][-1]["market"] == "U.S."
